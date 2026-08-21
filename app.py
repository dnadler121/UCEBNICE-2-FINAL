import os, json, unicodedata, random, html, re, base64, uuid, urllib.parse, urllib.request, urllib.error, zipfile, shutil, importlib.util, tempfile, threading, hmac, hashlib, ast, math, io
from pathlib import Path
from datetime import datetime
from flask import Flask, render_template, render_template_string, request, jsonify, session, redirect, url_for, send_from_directory, send_file, flash
from flask_sqlalchemy import SQLAlchemy
from sqlalchemy import inspect, text
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from markupsafe import Markup, escape
from dotenv import load_dotenv

BASE = Path(__file__).resolve().parent
load_dotenv(BASE / '.env')

# Na Renderu ukládáme všechna uživatelská data na Persistent Disk.
# Při lokálním spuštění zůstávají data ve složce projektu.
DATA_DIR = Path('/var/data') if Path('/var/data').is_dir() else BASE
DATA_DIR.mkdir(parents=True, exist_ok=True)

UPLOADS = DATA_DIR / 'uploads'
UPLOADS.mkdir(parents=True, exist_ok=True)

INTERACTIVE_LESSONS = DATA_DIR / 'interactive_lessons'
INTERACTIVE_LESSONS.mkdir(parents=True, exist_ok=True)

DB_PATH = DATA_DIR / 'montessori.db'

app = Flask(__name__)
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'change-me')
# Nový název cookie odřízne staré přihlášení z předchozích testovacích verzí.
# Když aplikaci spustíš poprvé, vždy tě pošle na přihlášení.
app.config['SESSION_COOKIE_NAME'] = 'montessori_engine_v1_2_role_login'
app.config["SQLALCHEMY_DATABASE_URI"] = "sqlite:///" + str(DB_PATH)
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)

# --- CZ/EN interface -------------------------------------------------------
I18N = {
    'cs': {},
    'en': {
        'Přihlášení':'Log in','Odhlásit':'Log out','Nepřihlášeno':'Not logged in',
        'Učitel':'Teacher','Student':'Student','Editor':'Editor','Studenti':'Students','Databáze':'Database',
        'Vyber předmět, školu a ročník, téma a lekci':'Choose a subject, school and grade, topic and lesson',
        'DIGITÁLNÍ UČEBNICE':'DIGITAL TEXTBOOK','Vyber si předmět a začni pracovat':'Choose a subject and start learning',
        'Jedno přihlášení, všechny lekce na jednom místě. Každá lekce se automaticky zařadí podle předmětu, školy, ročníku a tématu.':'One login, all lessons in one place. Each lesson is organised by subject, school, grade and topic.',
        'PŘEDMĚTY':'SUBJECTS','Biologie a občanská výchova':'Biology and Civics',
        'Výklady z HTML, otázky, aktivity a závěrečné testy.':'HTML lessons, questions, activities and final tests.',
        'INTERAKTIVNÍ LEKCE':'INTERACTIVE LESSONS','Matematika':'Mathematics','Informatika':'Computer Science',
        'Školy a ročníky, témata a samostatné matematické podaplikace.':'Schools and grades, topics and interactive mathematics activities.',
        'Algoritmy, programování, data a další praktická témata.':'Algorithms, programming, data and other practical topics.',
        'lekcí':'lessons','Otevřít →':'Open →','Učitelská správa':'Teacher administration',
        'Vytváření lekcí, studenti a výsledky':'Create lessons, manage students and results','HTML lekce':'HTML lesson',
        'Import matematika / informatika':'Import Mathematics / Computer Science','Výsledky':'Results',
        'Uživatelské jméno':'Username','Heslo':'Password','Vstoupit do aplikace':'Enter application',
        'Po přihlášení se teprve zobrazí předměty, bloky a lekce. Učitel navíc uvidí editor.':'After logging in, subjects, topics and lessons will appear. Teachers will also see the editor.',
        'Učitelský účet se načítá z':'The teacher account is loaded from','Studenty vytváří učitel v editoru.':'Students are created by the teacher in the editor.',
        'Zpět':'Back','Další':'Next','Pokračovat':'Continue','Dokončit':'Finish','Odevzdat':'Submit','Zkusit znovu':'Try again',
        'Správně':'Correct','Špatně':'Incorrect','Otázka':'Question','Otázky':'Questions','Výklad':'Lesson','Aktivita':'Activity',
        'Zajímavost':'Did you know?','Závěrečný test':'Final test','Hotovo':'Done','Uložit':'Save','Zrušit':'Cancel','Smazat':'Delete','Upravit':'Edit',
        'Přidat':'Add','Název':'Title','Téma':'Topic','Ročník':'Grade','Škola':'School','Předmět':'Subject','Popis':'Description',
        'Nová lekce':'New lesson','Vytvořit':'Create','Zveřejněno':'Published','Ano':'Yes','Ne':'No','Hledat':'Search',
        'Jazyk':'Language','Čeština':'Czech','Angličtina':'English','Jazyk je během rozpracované práce uzamčen.':'The language is locked while work is in progress.',
        'Objevuj, hledej ve studijním materiálu a postupuj krok za krokem.':'Explore, use the study material and progress step by step.',
        'Ukončit a uložit':'Finish and save','Teď to zkus sám':'Now try it yourself','Ověření se odemkne po lekci':'The assessment unlocks after the lesson',
        'Studijní materiál':'Study material','Vše potřebné najdeš tady.':'Everything you need is here.','MOJE CESTA':'MY JOURNEY','Aktuální úkol':'Current task',
        'OTÁZKA K VÝKLADU':'QUESTION ABOUT THE LESSON','Napiš odpověď':'Type your answer','Zkontrolovat':'Check','Další část':'Next section',
        'Jak pracovat':'How to work','Tip':'Tip','Když si nejsi jistý/á, hledej ve studijním materiálu vedle.':'If you are unsure, look in the study material next to the task.',
        'POZORUJ VIDEO':'WATCH THE VIDEO','ZASTAV A UKAŽ':'PAUSE AND POINT','NAJDI NA OBRÁZKU':'FIND IN THE PICTURE','SKLÁDAČKA / KARTIČKY':'PUZZLE / CARDS',
        'MISE VE SKUTEČNÉM SVĚTĚ':'REAL-WORLD MISSION','Zkontrolovat kartičky':'Check cards','Zkontrolovat skládačku':'Check puzzle','Zkontrolovat a pokračovat':'Check and continue',
        'Když si nejsi jistý/á, vše potřebné můžeš najít ve studijním materiálu vedle.':'If you are unsure, you can find everything you need in the study material next to the task.'
    }
}

def current_lang():
    lang = session.get('lang', 'cs')
    return lang if lang in ('cs','en') else 'cs'

def tr(value):
    if value is None: return ''
    return I18N.get(current_lang(), {}).get(str(value), str(value))

@app.route('/language/<lang>')
def set_language(lang):
    if lang not in ('cs','en'):
        return redirect(request.referrer or url_for('index'))
    # Jazyk lze přepínat i během studentské lekce. Obsah, otázky i aktivity
    # se znovu načtou v právě zvoleném jazyce.
    session.pop('language_locked', None)
    session.pop('work_lang', None)
    session['lang'] = lang
    session.modified = True
    return redirect(request.referrer or url_for('index'))

def lock_language():
    u = current_user()
    if u and u.role == 'student':
        session['language_locked'] = True
        session['work_lang'] = current_lang()
        session.modified = True

def unlock_language():
    session.pop('language_locked', None)
    session.pop('work_lang', None)
    session.modified = True

class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    name = db.Column(db.String(160), nullable=False)
    role = db.Column(db.String(20), default='student')
    password_hash = db.Column(db.String(255), nullable=False)

class Subject(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(120), nullable=False)
    icon = db.Column(db.String(20), default='🌱')

class Grade(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    subject_id = db.Column(db.Integer, db.ForeignKey('subject.id'), nullable=False)
    name = db.Column(db.String(60), nullable=False)
    subject = db.relationship('Subject', backref='grades')

class Block(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    grade_id = db.Column(db.Integer, db.ForeignKey('grade.id'), nullable=False)
    title = db.Column(db.String(160), nullable=False)
    order = db.Column(db.Integer, default=1)
    grade = db.relationship('Grade', backref='blocks')

class Lesson(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    block_id = db.Column(db.Integer, db.ForeignKey('block.id'), nullable=False)
    title = db.Column(db.String(200), nullable=False)
    title_en = db.Column(db.String(200), default='')
    tip = db.Column(db.Text, default='')
    tip_en = db.Column(db.Text, default='')
    hero_image = db.Column(db.String(255), default='')
    hero_image_en = db.Column(db.String(255), default='')
    order = db.Column(db.Integer, default=1)
    is_published = db.Column(db.Boolean, default=True)
    block = db.relationship('Block', backref='lessons')

class Section(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    heading = db.Column(db.String(200), nullable=False)
    heading_en = db.Column(db.String(200), default='')
    text = db.Column(db.Text, default='')
    text_en = db.Column(db.Text, default='')
    interest = db.Column(db.Text, default='')
    interest_en = db.Column(db.Text, default='')
    image = db.Column(db.String(255), default='')
    image_en = db.Column(db.String(255), default='')
    activity = db.Column(db.Text, default='')
    activity_en = db.Column(db.Text, default='')
    order = db.Column(db.Integer, default=1)
    lesson = db.relationship('Lesson', backref='sections')

class InlineImage(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    section_id = db.Column(db.Integer, db.ForeignKey('section.id'), nullable=False)
    file = db.Column(db.String(255), nullable=False)
    caption = db.Column(db.String(255), default='')
    order = db.Column(db.Integer, default=1)
    section = db.relationship('Section', backref='inline_images')

class Question(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    section_id = db.Column(db.Integer, db.ForeignKey('section.id'), nullable=True)
    area = db.Column(db.String(20), default='study')  # study/final
    qtype = db.Column(db.String(30), default='choice')
    question = db.Column(db.Text, nullable=False)
    options_json = db.Column(db.Text, default='[]')
    correct_json = db.Column(db.Text, default='0')
    roots_json = db.Column(db.Text, default='[]')
    hint = db.Column(db.Text, default='')
    order = db.Column(db.Integer, default=1)
    lang = db.Column(db.String(2), default='cs')
    lesson = db.relationship('Lesson', backref='questions')
    section = db.relationship('Section', backref='questions')

class PracticalActivity(db.Model):
    """Univerzální praktická aktivita pro Biologii / Občanskou výchovu."""
    id = db.Column(db.Integer, primary_key=True)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    section_id = db.Column(db.Integer, db.ForeignKey('section.id'), nullable=True)
    order = db.Column(db.Integer, default=1)
    lang = db.Column(db.String(2), default='cs')
    activity_type = db.Column(db.String(40), nullable=False, default='find_image')
    title = db.Column(db.String(220), default='Praktická aktivita')
    title_en = db.Column(db.String(220), default='')
    prompt = db.Column(db.Text, default='')
    prompt_en = db.Column(db.Text, default='')
    config_json = db.Column(db.Text, default='{}')
    config_en_json = db.Column(db.Text, default='{}')
    image_file = db.Column(db.String(255), default='')
    video_file = db.Column(db.String(255), default='')
    include_final = db.Column(db.Boolean, default=True)
    lesson = db.relationship('Lesson', backref='practical_activities')
    section = db.relationship('Section', backref='practical_activities')


class StudentSectionProgress(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    section_id = db.Column(db.Integer, db.ForeignKey('section.id'), nullable=False)
    read_complete = db.Column(db.Boolean, default=False)
    completed = db.Column(db.Boolean, default=False)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow)


class StudentActivityProgress(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    activity_id = db.Column(db.Integer, db.ForeignKey('practical_activity.id'), nullable=False)
    context = db.Column(db.String(20), default='study')
    completed = db.Column(db.Boolean, default=False)
    attempts = db.Column(db.Integer, default=0)
    answer_json = db.Column(db.Text, default='{}')
    updated_at = db.Column(db.DateTime, default=datetime.utcnow)


class StudyQuestionProgress(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    question_id = db.Column(db.Integer, db.ForeignKey('question.id'), nullable=False)
    completed = db.Column(db.Boolean, default=False)
    attempts = db.Column(db.Integer, default=0)
    answer_json = db.Column(db.Text, default='{}')
    updated_at = db.Column(db.DateTime, default=datetime.utcnow)


class FinalItemProgress(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    item_key = db.Column(db.String(80), nullable=False)
    completed = db.Column(db.Boolean, default=False)
    attempts = db.Column(db.Integer, default=0)
    answer_json = db.Column(db.Text, default='{}')
    updated_at = db.Column(db.DateTime, default=datetime.utcnow)


class Result(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    percent = db.Column(db.Integer, default=0)
    grade = db.Column(db.Integer, default=5)
    score = db.Column(db.Integer, default=0)
    total = db.Column(db.Integer, default=0)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    focus_lost = db.Column(db.Integer, default=0)
    status = db.Column(db.String(60), default='dokončeno')
    user = db.relationship('User')
    lesson = db.relationship('Lesson')


class InteractiveLesson(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    slug = db.Column(db.String(120), unique=True, nullable=False)
    subject = db.Column(db.String(40), nullable=False)
    school = db.Column(db.String(160), nullable=False)
    grade_name = db.Column(db.String(120), nullable=False)
    topic = db.Column(db.String(160), nullable=False)
    title = db.Column(db.String(220), nullable=False)
    description = db.Column(db.Text, default='')
    icon = db.Column(db.String(20), default='📘')
    package_dir = db.Column(db.String(255), nullable=False)
    is_published = db.Column(db.Boolean, default=True)
    imported_at = db.Column(db.DateTime, default=datetime.utcnow)


class InteractiveResult(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    interactive_lesson_id = db.Column(db.Integer, db.ForeignKey('interactive_lesson.id'), nullable=False)
    percent = db.Column(db.Integer, default=100)
    grade = db.Column(db.Integer, default=1)
    completed_at = db.Column(db.DateTime, default=datetime.utcnow)
    focus_lost = db.Column(db.Integer, default=0)
    status = db.Column(db.String(60), default='dokončeno')
    user = db.relationship('User')
    interactive_lesson = db.relationship('InteractiveLesson')


class InteractiveProgress(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    interactive_lesson_id = db.Column(db.Integer, db.ForeignKey('interactive_lesson.id'), nullable=False)
    completed = db.Column(db.Boolean, default=False)
    current_grade = db.Column(db.Integer, default=5)
    last_completed_at = db.Column(db.DateTime, nullable=True)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow)
    user = db.relationship('User')
    interactive_lesson = db.relationship('InteractiveLesson')


class InformaticsLesson(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    school = db.Column(db.String(160), default='')
    grade_name = db.Column(db.String(120), default='')
    topic = db.Column(db.String(180), default='')
    title = db.Column(db.String(220), nullable=False)
    intro = db.Column(db.Text, default='')
    html_original = db.Column(db.String(255), default='')
    html_stored = db.Column(db.String(255), default='')
    is_published = db.Column(db.Boolean, default=True)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


class InformaticsTask(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    lesson_id = db.Column(db.Integer, db.ForeignKey('informatics_lesson.id'), nullable=False)
    order = db.Column(db.Integer, default=1)
    title = db.Column(db.String(220), nullable=False)
    assignment = db.Column(db.Text, default='')
    source_original = db.Column(db.String(255), nullable=False)
    source_stored = db.Column(db.String(255), nullable=False)
    file_type = db.Column(db.String(40), default='')
    analysis_json = db.Column(db.Text, default='{}')
    checks_json = db.Column(db.Text, default='[]')
    image_file = db.Column(db.String(255), default='')
    lesson = db.relationship('InformaticsLesson', backref='tasks')


class InformaticsWorkFile(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    task_id = db.Column(db.Integer, db.ForeignKey('informatics_task.id'), nullable=False)
    token = db.Column(db.String(255), nullable=False, unique=True)
    original_name = db.Column(db.String(255), nullable=False)
    stored_name = db.Column(db.String(255), nullable=False)
    downloaded_at = db.Column(db.DateTime, default=datetime.utcnow)
    user = db.relationship('User')
    task = db.relationship('InformaticsTask')


class InformaticsSubmission(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    task_id = db.Column(db.Integer, db.ForeignKey('informatics_task.id'), nullable=False)
    original_name = db.Column(db.String(255), nullable=False)
    stored_name = db.Column(db.String(255), nullable=False)
    feedback_json = db.Column(db.Text, default='[]')
    percent = db.Column(db.Integer, default=0)
    grade = db.Column(db.Integer, default=5)
    status = db.Column(db.String(60), default='kontrola')
    focus_lost = db.Column(db.Integer, default=0)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    user = db.relationship('User')
    task = db.relationship('InformaticsTask')


class MathLesson(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    school = db.Column(db.String(160), default='')
    grade_name = db.Column(db.String(120), default='')
    topic = db.Column(db.String(180), default='')
    title = db.Column(db.String(220), nullable=False)
    html_original = db.Column(db.String(255), default='')
    html_stored = db.Column(db.String(255), default='')
    is_published = db.Column(db.Boolean, default=True)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


class MathExample(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    lesson_id = db.Column(db.Integer, db.ForeignKey('math_lesson.id'), nullable=False)
    order = db.Column(db.Integer, default=1)
    title = db.Column(db.String(220), default='')
    problem = db.Column(db.Text, nullable=False)
    # Volitelné slovní zadání pro žáka. 'problem' zůstává čistý matematický vzor pro engine.
    prose_problem = db.Column(db.Text, default='')
    # Volitelný obrázek / náčrt přímo k tomuto matematickému příkladu.
    # Soubor se ukládá do trvalé složky uploads, v DB je pouze jeho název.
    image_stored = db.Column(db.String(255), default='')
    variant_enabled = db.Column(db.Boolean, default=False)
    variant_values = db.Column(db.Text, default='')
    variant_condition = db.Column(db.Text, default='')
    variant_min = db.Column(db.Float, default=1)
    variant_max = db.Column(db.Float, default=30)
    variant_step = db.Column(db.Float, default=1)
    # Obecná pravidla pro kvalitu dopočítaných výsledků náhodné varianty.
    variant_result_kind = db.Column(db.String(20), default='any')
    variant_result_sign = db.Column(db.String(20), default='any')
    variant_result_min = db.Column(db.Float, nullable=True)
    variant_result_max = db.Column(db.Float, nullable=True)
    variant_result_decimals = db.Column(db.Integer, default=-1)
    lesson = db.relationship('MathLesson', backref='examples')


class MathStep(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    example_id = db.Column(db.Integer, db.ForeignKey('math_example.id'), nullable=False)
    order = db.Column(db.Integer, default=1)
    instruction = db.Column(db.Text, nullable=False)
    expected = db.Column(db.Text, nullable=False)
    hint = db.Column(db.Text, default='')
    # Volitelný jeden nebo více obecných vzorců pro dopočet výsledků z n1, n2, ...
    # Více vzorců odděl středníkem / novým řádkem; postupně nahradí čísla ve vzorové odpovědi.
    result_formula = db.Column(db.Text, default='')
    result_decimals = db.Column(db.Integer, default=2)
    example = db.relationship('MathExample', backref='steps')


class MathAttempt(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    lesson_id = db.Column(db.Integer, db.ForeignKey('math_lesson.id'), nullable=False)
    completed_steps = db.Column(db.Integer, default=0)
    total_steps = db.Column(db.Integer, default=0)
    percent = db.Column(db.Integer, default=0)
    grade = db.Column(db.Integer, default=5)
    status = db.Column(db.String(60), default='rozpracováno')
    focus_lost = db.Column(db.Integer, default=0)
    answers_json = db.Column(db.Text, default='[]')
    # Stabilně uložené náhodné varianty pro jednotlivé příklady této lekce.
    # JSON: {"example_id": [n1, n2, ...]}
    variant_json = db.Column(db.Text, default='{}')
    updated_at = db.Column(db.DateTime, default=datetime.utcnow)
    user = db.relationship('User')
    lesson = db.relationship('MathLesson')


class LessonFocusSession(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    lesson_kind = db.Column(db.String(20), nullable=False)  # html / interactive
    lesson_key = db.Column(db.String(160), nullable=False)
    count = db.Column(db.Integer, default=0)
    terminated = db.Column(db.Boolean, default=False)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow)


class StudentProgress(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    current_step = db.Column(db.Integer, default=0)
    status = db.Column(db.String(40), default='rozpracováno')
    updated_at = db.Column(db.DateTime, default=datetime.utcnow)
    user = db.relationship('User')
    lesson = db.relationship('Lesson')


class StudentLessonReset(db.Model):
    """Jednorázový příznak, že učitel resetoval studentovi konkrétní HTML lekci.

    Slouží hlavně k vyčištění starého průběhu uloženého v session v prohlížeči
    studenta při jeho příštím otevření lekce.
    """
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


def ensure_informatics_columns():
    """Doplní nové sloupce do starší lokální/Render databáze bez mazání dat."""
    try:
        insp = inspect(db.engine)
        if 'informatics_lesson' not in insp.get_table_names():
            return
        cols = {c['name'] for c in insp.get_columns('informatics_lesson')}
        with db.engine.begin() as conn:
            if 'html_original' not in cols:
                conn.execute(text("ALTER TABLE informatics_lesson ADD COLUMN html_original VARCHAR(255) DEFAULT ''"))
            if 'html_stored' not in cols:
                conn.execute(text("ALTER TABLE informatics_lesson ADD COLUMN html_stored VARCHAR(255) DEFAULT ''"))
    except Exception as exc:
        print('Informatics migration warning:', exc)


def strip_accents(s):
    return ''.join(c for c in unicodedata.normalize('NFD', str(s).lower()) if unicodedata.category(c) != 'Mn').strip()

def grade_from_percent(percent):
    if percent >= 90: return 1
    if percent >= 75: return 2
    if percent >= 60: return 3
    if percent >= 40: return 4
    return 5

def current_user():
    uid = session.get('uid')
    return db.session.get(User, uid) if uid else None


def touch_progress(lesson_id, step=0, status='rozpracováno'):
    u = current_user()
    if not u or u.role != 'student':
        return
    pr = StudentProgress.query.filter_by(user_id=u.id, lesson_id=lesson_id).first()
    if not pr:
        pr = StudentProgress(user_id=u.id, lesson_id=lesson_id)
        db.session.add(pr)
    pr.current_step = int(step or 0)
    pr.status = status
    pr.updated_at = datetime.utcnow()
    db.session.commit()

def consume_pending_lesson_reset(lesson_id):
    """Vyčistí starý session průběh po učitelském resetu lekce.

    Trvalý databázový průběh je smazán už při stisku Reset/Smazat v administraci.
    Tohle odstraní ještě lokální session data v prohlížeči studenta, aby se z nich
    výsledek po přihlášení nemohl znovu vytvořit.
    """
    u = current_user()
    if not u or u.role != 'student':
        return False
    marker = StudentLessonReset.query.filter_by(user_id=u.id, lesson_id=lesson_id).first()
    if not marker:
        return False

    partial = session.get('html_partial_progress', {})
    if str(lesson_id) in partial:
        partial.pop(str(lesson_id), None)
        session['html_partial_progress'] = partial

    completed = session.get('completed_steps', {})
    if str(lesson_id) in completed:
        completed.pop(str(lesson_id), None)
        session['completed_steps'] = completed

    session.modified = True
    db.session.delete(marker)
    db.session.commit()
    return True


def reset_student_lesson_progress(user_id, lesson_id):
    """Kompletní reset jedné HTML lekce pro jednoho studenta."""
    Result.query.filter_by(user_id=user_id, lesson_id=lesson_id).delete(synchronize_session=False)
    StudentProgress.query.filter_by(user_id=user_id, lesson_id=lesson_id).delete(synchronize_session=False)
    StudentSectionProgress.query.filter_by(user_id=user_id, lesson_id=lesson_id).delete(synchronize_session=False)
    StudyQuestionProgress.query.filter_by(user_id=user_id, lesson_id=lesson_id).delete(synchronize_session=False)
    FinalItemProgress.query.filter_by(user_id=user_id, lesson_id=lesson_id).delete(synchronize_session=False)

    activity_ids = [a.id for a in PracticalActivity.query.filter_by(lesson_id=lesson_id).all()]
    if activity_ids:
        StudentActivityProgress.query.filter(
            StudentActivityProgress.user_id == user_id,
            StudentActivityProgress.activity_id.in_(activity_ids)
        ).delete(synchronize_session=False)

    LessonFocusSession.query.filter_by(
        user_id=user_id, lesson_kind='html', lesson_key=str(lesson_id)
    ).delete(synchronize_session=False)

    # Při příštím otevření lekce studentem smažeme také stará session data
    # uložená v jeho vlastním prohlížeči.
    StudentLessonReset.query.filter_by(user_id=user_id, lesson_id=lesson_id).delete(synchronize_session=False)
    db.session.add(StudentLessonReset(user_id=user_id, lesson_id=lesson_id))
    db.session.commit()


def last_result_for_student(user_id):
    return Result.query.filter_by(user_id=user_id).order_by(Result.created_at.desc()).first()

def current_progress_for_student(user_id):
    return StudentProgress.query.filter_by(user_id=user_id).order_by(StudentProgress.updated_at.desc()).first()

def last_informatics_submission_for_student(user_id):
    return InformaticsSubmission.query.filter(
        InformaticsSubmission.user_id == user_id,
        InformaticsSubmission.status != 'kontrola'
    ).order_by(InformaticsSubmission.created_at.desc()).first()

def last_math_attempt_for_student(user_id):
    return MathAttempt.query.filter(
        MathAttempt.user_id == user_id,
        MathAttempt.status != 'rozpracováno'
    ).order_by(MathAttempt.updated_at.desc()).first()

def student_overview_rows():
    rows = []
    for stu in User.query.filter_by(role='student').order_by(User.name).all():
        pr = current_progress_for_student(stu.id)
        res = last_result_for_student(stu.id)
        inf = last_informatics_submission_for_student(stu.id)
        math = last_math_attempt_for_student(stu.id)
        inf_grade = informatics_grade_from_percent(inf.percent) if inf else None
        rows.append({
            'student': stu,
            'progress': pr,
            'result': res,
            'informatics': inf,
            'informatics_grade': inf_grade,
            'math': math,
        })
    return rows

def require_login():
    if not current_user(): return redirect(url_for('login'))

def require_teacher():
    r = require_login()
    if r: return r
    if current_user().role != 'teacher': return redirect(url_for('dashboard'))

def role_home():
    u = current_user()
    if not u:
        return redirect(url_for('login'))
    if u.role == 'teacher':
        return redirect(url_for('teacher_home'))
    return redirect(url_for('dashboard'))

def completed_steps_for(lesson_id):
    user = current_user()
    lesson = db.session.get(Lesson, lesson_id)
    if user and user.role == 'student' and lesson:
        section_ids = [s.id for s in sorted(lesson.sections, key=lambda x: x.order)]
        completed_ids = {r.section_id for r in StudentSectionProgress.query.filter_by(user_id=user.id, lesson_id=lesson_id, completed=True).all()}
        return {idx for idx, sid in enumerate(section_ids) if sid in completed_ids}
    done = session.get('completed_steps', {})
    return set(done.get(str(lesson_id), []))

def mark_step_complete(lesson_id, step):
    lesson = db.session.get(Lesson, lesson_id)
    user = current_user()
    if user and user.role == 'student' and lesson:
        sections = sorted(lesson.sections, key=lambda x: x.order)
        if 0 <= int(step) < len(sections):
            sec = sections[int(step)]
            row = StudentSectionProgress.query.filter_by(user_id=user.id, lesson_id=lesson_id, section_id=sec.id).first()
            if not row:
                row = StudentSectionProgress(user_id=user.id, lesson_id=lesson_id, section_id=sec.id)
                db.session.add(row)
            row.read_complete = True
            row.completed = True
            row.updated_at = datetime.utcnow()
            db.session.commit()
            return
    done = session.setdefault('completed_steps', {})
    key = str(lesson_id)
    arr = set(done.get(key, []))
    arr.add(int(step))
    done[key] = sorted(arr)
    session.modified = True

def lesson_ready_for_test(lesson):
    data = lesson_to_dict(lesson)
    needed = set(range(len(data['sections'])))
    return needed.issubset(completed_steps_for(lesson.id))

def activity_to_dict(a):
    is_en = current_lang() == 'en'
    raw_cfg = (a.config_en_json or '') if is_en else (a.config_json or '')
    if is_en and not str(raw_cfg).strip():
        raw_cfg = a.config_json or '{}'
    try:
        config = json.loads(raw_cfg or '{}')
    except Exception:
        config = {}
    return {
        'id': a.id, 'type': a.activity_type,
        'title': (a.title_en or a.title) if is_en else a.title,
        'prompt': (a.prompt_en or a.prompt) if is_en else a.prompt,
        'config': config, 'image': a.image_file or '', 'video': a.video_file or '',
        'include_final': bool(a.include_final), 'order': a.order,
    }

def activity_completed(activity_id, context='study'):
    user = current_user()
    if not user or user.role != 'student':
        return False
    row = StudentActivityProgress.query.filter_by(user_id=user.id, activity_id=activity_id, context=context).first()
    return bool(row and row.completed)

def slovni_hodnoceni(percent):
    percent = int(percent or 0)
    if percent >= 90: return 'Výborně zvládnuto'
    if percent >= 75: return 'Zvládnuto'
    if percent >= 60: return 'Většinu už zvládám'
    if percent >= 40: return 'Ještě objevuji'
    return 'Začínám se orientovat'

def q_to_dict(q):
    try:
        options = json.loads(q.options_json or '[]')
    except Exception:
        options = []
    d = {'id': q.id, 'type': q.qtype, 'question': q.question, 'options': options, 'hint': q.hint}
    if q.qtype in ['choice','image_choice']:
        d['correct'] = json.loads(q.correct_json or '0')
        if q.qtype == 'image_choice':
            d['images'] = options
            d['labels'] = [''] * len(options)
    else:
        # Krátká odpověď může mít také obrázek. Kvůli kompatibilitě ho ukládáme do options_json jako {"image": "soubor.jpg"}.
        d['roots'] = json.loads(q.roots_json or '[]')
        d['image'] = options.get('image','') if isinstance(options, dict) else ''
    return d

def lesson_to_dict(lesson):
    lang = current_lang()
    is_en = lang == 'en'
    sections = []
    all_activities = sorted([a for a in lesson.practical_activities if (a.lang or 'cs') != 'en'], key=lambda x: x.order)
    for sec in sorted(lesson.sections, key=lambda x:x.order):
        sec_acts = [activity_to_dict(a) for a in all_activities if a.section_id == sec.id]
        # Česká sada určuje STRUKTURU lekce. Angličtina pouze lokalizuje stejné
        # kroky podle pořadí. Tím přepnutí CZ/EN nikdy nepřidá/neztratí video,
        # aktivitu ani jiný typ kroku jen proto, že je v EN JSONu jiný počet otázek.
        cs_qs = [q for q in sorted(sec.questions, key=lambda x:(x.order, x.id)) if q.area=='study' and (q.lang or 'cs') == 'cs']
        en_qs = [q for q in sorted(sec.questions, key=lambda x:(x.order, x.id)) if q.area=='study' and (q.lang or 'cs') == 'en']
        if is_en:
            qs = [en_qs[i] if i < len(en_qs) else q for i, q in enumerate(cs_qs)]
        else:
            qs = cs_qs
        sections.append({
            'id': sec.id,
            'heading': (sec.heading_en or '') if is_en else sec.heading,
            'text': (sec.text_en or '') if is_en else sec.text,
            'interest': (sec.interest_en or '') if is_en else sec.interest,
            'image': (sec.image_en or '') if is_en else sec.image,
            'activity': (sec.activity_en or '') if is_en else sec.activity,
            'questions': [q_to_dict(q) for q in qs],
            'activities': sec_acts,
        })
    subject = lesson.block.grade.subject
    grade = lesson.block.grade
    final_questions = [q for sec in sections for q in sec['questions']]
    final_activities = [activity_to_dict(a) for a in all_activities if a.include_final]
    return {
        '_id': lesson.id, '_slug': lesson.id, 'subject': tr(subject.name) if is_en else subject.name, 'icon': subject.icon, 'grade': grade.name,
        'block': lesson.block.title,
        'title': (lesson.title_en or '') if is_en else lesson.title,
        'tip': (lesson.tip_en or '') if is_en else lesson.tip,
        'hero_image': (lesson.hero_image_en or '') if is_en else lesson.hero_image,
        'sections': sections,
        'final_test': final_questions,
        'final_activities': final_activities,
        'content_language': lang,
        'has_content': bool(((lesson.title_en or '').strip() or any((x.get('text') or '').strip() or x.get('questions') or x.get('activities') for x in sections)) if is_en else True),
    }

def lesson_gallery(lesson):
    """Vrátí obrázky dostupné v editoru pro konkrétní lekci."""
    seen = []
    def add(name):
        if name and name not in seen:
            seen.append(name)
    if lesson:
        add(lesson.hero_image)
        for sec in lesson.sections:
            add(sec.image)
            for q in lesson.questions:
                try:
                    opts = json.loads(q.options_json or '[]')
                except Exception:
                    opts = []
                if q.qtype == 'image_choice':
                    for img in opts:
                        add(img)
                elif q.qtype == 'text' and isinstance(opts, dict):
                    add(opts.get('image',''))
    return seen

def course_from_lesson(lesson):
    if lesson:
        sub = lesson.block.grade.subject
        return {'subject': sub.name, 'grade': lesson.block.grade.name, 'block': lesson.block.title, 'icon': sub.icon}
    return {'subject': 'Montessori', 'grade': '', 'block': '', 'icon': '🌱'}

def cleanup_empty_curriculum():
    """Odstraní prázdné bloky/ročníky/předměty po trvalém smazání lekcí.

    Díky tomu v navigaci nezůstávají staré odkazy nebo prázdné nadpisy po
    lekcích, které učitel už smazal.
    """
    # Mažeme od nejnižší úrovně nahoru. Flush zajistí, že následující dotaz
    # už vidí aktuální stav i v rámci stejné transakce.
    for block in Block.query.all():
        if Lesson.query.filter_by(block_id=block.id).count() == 0:
            db.session.delete(block)
    db.session.flush()
    for grade in Grade.query.all():
        if Block.query.filter_by(grade_id=grade.id).count() == 0:
            db.session.delete(grade)
    db.session.flush()
    for subject in Subject.query.all():
        if Grade.query.filter_by(subject_id=subject.id).count() == 0:
            db.session.delete(subject)
    db.session.flush()


def visible_lessons():
    return Lesson.query.filter_by(is_published=True).join(Block).join(Grade).join(Subject).order_by(Subject.name, Grade.name, Block.order, Lesson.order).all()

@app.context_processor
def inject():
    u = current_user()
    last = None
    if u:
        res = Result.query.filter_by(user_id=u.id).order_by(Result.created_at.desc()).first()
        if res: last = {'lesson': res.lesson.title, 'percent': res.percent, 'grade': res.grade, 'score': res.score, 'total': res.total}
    return {'user': u, 'last_result': last, 'lang': current_lang(), 't': tr, 'language_locked': bool(session.get('language_locked'))}


def normalize_subject(value):
    value = strip_accents(value).replace(' ', '-')
    aliases = {'matematika': 'matematika', 'math': 'matematika',
               'informatika': 'informatika', 'ict': 'informatika'}
    return aliases.get(value, value)


def safe_package_slug(value):
    value = strip_accents(value)
    return re.sub(r'[^a-z0-9]+', '-', value).strip('-')[:120]


def interactive_groups_for(subject_kind):
    subject_value = 'matematika' if subject_kind == 'matematika' else 'informatika'
    lessons = InteractiveLesson.query.filter_by(
        subject=subject_value, is_published=True
    ).order_by(
        InteractiveLesson.school,
        InteractiveLesson.grade_name,
        InteractiveLesson.topic,
        InteractiveLesson.title
    ).all()

    grouped = {}
    for lesson in lessons:
        key = (lesson.school, lesson.grade_name)
        group = grouped.setdefault(key, {
            'school': lesson.school,
            'grade': lesson.grade_name,
            'topics': {}
        })
        group['topics'].setdefault(lesson.topic, []).append(lesson)
    return list(grouped.values())


def safe_extract_zip(zip_file, target):
    target = target.resolve()
    for member in zip_file.infolist():
        member_path = (target / member.filename).resolve()
        if target not in member_path.parents and member_path != target:
            raise ValueError('Balíček obsahuje nepovolenou cestu.')
    zip_file.extractall(target)


def find_package_root(temp_dir):
    candidates = list(temp_dir.rglob('lesson.json'))
    if len(candidates) != 1:
        raise ValueError('Balíček musí obsahovat právě jeden soubor lesson.json.')
    return candidates[0].parent


def load_interactive_module(lesson):
    module_file = BASE / lesson.package_dir / 'lesson_app.py'
    if not module_file.exists():
        return None
    module_name = f'interactive_{lesson.slug}_{module_file.stat().st_mtime_ns}'
    spec = importlib.util.spec_from_file_location(module_name, module_file)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def upsert_interactive_progress(lesson, percent=100, grade=None, focus_lost=0):
    user = current_user()
    if not user or user.role != 'student':
        return
    if grade is None:
        grade = grade_from_percent(percent)
    progress = InteractiveProgress.query.filter_by(
        user_id=user.id,
        interactive_lesson_id=lesson.id
    ).first()
    if not progress:
        progress = InteractiveProgress(
            user_id=user.id,
            interactive_lesson_id=lesson.id
        )
        db.session.add(progress)
    progress.completed = True
    progress.current_grade = int(grade)
    progress.last_completed_at = datetime.utcnow()
    progress.updated_at = datetime.utcnow()
    db.session.add(InteractiveResult(
        user_id=user.id,
        interactive_lesson_id=lesson.id,
        percent=int(percent),
        grade=int(grade),
        focus_lost=int(focus_lost or 0),
        status='dokončeno'
    ))
    db.session.commit()

def get_focus_session(kind, key, create=False):
    user = current_user()
    if not user or user.role != 'student':
        return None
    row = LessonFocusSession.query.filter_by(
        user_id=user.id, lesson_kind=str(kind), lesson_key=str(key)
    ).first()
    if not row and create:
        row = LessonFocusSession(
            user_id=user.id, lesson_kind=str(kind), lesson_key=str(key), count=0
        )
        db.session.add(row)
        db.session.flush()
    return row


def get_focus_count(kind, key):
    user = current_user()
    if not user:
        return 0
    row = LessonFocusSession.query.filter_by(
        user_id=user.id,
        lesson_kind=str(kind),
        lesson_key=str(key)
    ).first()
    return int(row.count or 0) if row else 0


def focus_attempt_marker(kind, key):
    return f'{str(kind)}:{str(key)}'


def begin_focus_attempt(kind, key):
    """Při prvním otevření lekce v novém přihlášení začne počítání od nuly.

    Přechody mezi částmi stejné lekce ani obnovení stránky počítadlo nemažou.
    Po automatickém ukončení nebo řádném dokončení může další otevření
    stejné lekce začít jako nový pokus.
    """
    user = current_user()
    if not user or user.role != 'student':
        return

    marker = focus_attempt_marker(kind, key)
    active = list(session.get('focus_active_attempts', []))
    if marker in active:
        return

    row = get_focus_session(kind, key, create=False)
    if row:
        db.session.delete(row)
        db.session.flush()

    active.append(marker)
    session['focus_active_attempts'] = active
    session.modified = True
    db.session.commit()


def end_focus_attempt(kind, key):
    marker = focus_attempt_marker(kind, key)
    active = list(session.get('focus_active_attempts', []))
    if marker in active:
        active.remove(marker)
        session['focus_active_attempts'] = active
        session.modified = True


def consume_focus_count(kind, key):
    row = get_focus_session(kind, key, create=False)
    count = int(row.count) if row else 0
    if row:
        db.session.delete(row)
    end_focus_attempt(kind, key)
    return count


@app.route('/api/focus-lost', methods=['POST'])
def api_focus_lost():
    r = require_login()
    if r:
        return jsonify({'ok': False, 'error': 'login'}), 401
    user = current_user()
    if user.role != 'student':
        return jsonify({'ok': True, 'ignored': True})

    data = request.get_json(silent=True) or {}
    kind = str(data.get('kind', '')).strip()
    key = str(data.get('key', '')).strip()
    if kind not in ('html', 'interactive', 'informatics', 'math') or not key:
        return jsonify({'ok': False, 'error': 'Neplatná lekce.'}), 400

    row = get_focus_session(kind, key, create=True)
    if row.terminated:
        return jsonify({'ok': True, 'count': row.count, 'terminated': True,
                        'redirect': url_for('focus_terminated')})

    row.count = min(3, int(row.count or 0) + 1)
    row.updated_at = datetime.utcnow()
    row.terminated = row.count >= 3

    if row.terminated:
        if kind == 'html':
            lesson_item = db.session.get(Lesson, int(key)) if key.isdigit() else None
            if lesson_item:
                if final_attempt_is_active(lesson_item.id):
                    # Třetí opuštění ukončí právě běžící závěrečný pokus.
                    # Učitel dostane procenta dosažená v tomto posledním pokusu.
                    persist_final_result(lesson_item, status='ukončeno po 3 opuštěních', focus_lost=3)
                    final_attempt_mark_finished(lesson_item.id)
                else:
                    touch_progress(lesson_item.id, 0, 'ukončeno po 3 opuštěních')
        elif kind == 'interactive':
            lesson_item = InteractiveLesson.query.filter_by(slug=key).first()
            if lesson_item:
                db.session.add(InteractiveResult(
                    user_id=user.id, interactive_lesson_id=lesson_item.id,
                    percent=0, grade=5, focus_lost=3,
                    status='ukončeno po 3 opuštěních'
                ))
        elif kind == 'informatics':
            task_item = db.session.get(InformaticsTask, int(key)) if str(key).isdigit() else None
            if task_item:
                last = InformaticsSubmission.query.filter_by(
                    user_id=user.id, task_id=task_item.id
                ).order_by(InformaticsSubmission.created_at.desc()).first()
                percent = last.percent if last else 0
                db.session.add(InformaticsSubmission(
                    user_id=user.id, task_id=task_item.id,
                    original_name=(last.original_name if last else 'bez_souboru'),
                    stored_name=(last.stored_name if last else ''),
                    feedback_json=(last.feedback_json if last else '[]'),
                    percent=percent, grade=informatics_grade_from_percent(percent),
                    status='ukončeno po 3 opuštěních', focus_lost=3
                ))
        elif kind == 'math':
            lesson_item = db.session.get(MathLesson, int(key)) if str(key).isdigit() else None
            if lesson_item:
                total = MathStep.query.join(MathExample).filter(MathExample.lesson_id == lesson_item.id).count()
                attempt = MathAttempt.query.filter_by(user_id=user.id, lesson_id=lesson_item.id).first()
                done = attempt.completed_steps if attempt else 0
                percent = round(done / max(total,1) * 100)
                if not attempt:
                    attempt = MathAttempt(user_id=user.id, lesson_id=lesson_item.id)
                    db.session.add(attempt)
                attempt.completed_steps = done
                attempt.total_steps = total
                attempt.percent = percent
                attempt.grade = informatics_grade_from_percent(percent)
                attempt.status = 'ukončeno po 3 opuštěních'
                attempt.focus_lost = 3
                attempt.updated_at = datetime.utcnow()
        end_focus_attempt(kind, key)
        db.session.commit()
        return jsonify({'ok': True, 'count': 3, 'terminated': True,
                        'redirect': url_for('focus_terminated')})

    db.session.commit()
    return jsonify({'ok': True, 'count': row.count, 'terminated': False})


@app.route('/lesson-ukoncena')
def focus_terminated():
    r = require_login()
    if r:
        return r
    return render_template('terminated.html', course=course_from_lesson(None), lesson=None)


@app.route('/login', methods=['GET','POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username','').strip().lower()
        password = request.form.get('password','')
        u = User.query.filter_by(username=username).first()
        if u and check_password_hash(u.password_hash, password):
            session['uid'] = u.id
            if u.role == 'teacher':
                return redirect(url_for('teacher_home'))
            return redirect(url_for('portal'))
        return render_template('login.html', course=course_from_lesson(None), error='Špatné jméno nebo heslo.')
    return render_template('login.html', course=course_from_lesson(None), error=None)

@app.route('/logout')
def logout():
    session.clear(); return redirect(url_for('login'))

@app.route('/')
def index():
    r = require_login()
    if r: return r
    return redirect(url_for('portal'))

@app.route('/portal')
def portal():
    r = require_login()
    if r: return r
    counts = {
        'bio_obc': Lesson.query.join(Block).join(Grade).join(Subject).filter(
            db.or_(Subject.name.ilike('%bio%'), Subject.name.ilike('%občan%'), Subject.name.ilike('%obcan%')),
            Lesson.is_published.is_(True)
        ).count(),
        'matematika': Lesson.query.join(Block).join(Grade).join(Subject).filter(
            Subject.name.ilike('%matemat%'), Lesson.is_published.is_(True)
        ).count() + InteractiveLesson.query.filter_by(subject='matematika', is_published=True).count()
        + MathLesson.query.filter_by(is_published=True).count(),
        'informatika': Lesson.query.join(Block).join(Grade).join(Subject).filter(
            Subject.name.ilike('%informat%'), Lesson.is_published.is_(True)
        ).count() + InteractiveLesson.query.filter_by(subject='informatika', is_published=True).count()
        + InformaticsLesson.query.filter_by(is_published=True).count(),
    }
    return render_template('portal.html', course=course_from_lesson(None), lesson=None, counts=counts)

@app.route('/catalog/<kind>')
def subject_catalog(kind):
    r = require_login()
    if r: return r
    filters = {
        'bio-obc': lambda q: q.filter(db.or_(Subject.name.ilike('%bio%'), Subject.name.ilike('%občan%'), Subject.name.ilike('%obcan%'))),
        'matematika': lambda q: q.filter(Subject.name.ilike('%matemat%')),
        'informatika': lambda q: q.filter(Subject.name.ilike('%informat%')),
    }
    if kind not in filters:
        return 'Neznámý předmět', 404
    q = Subject.query
    subjects = filters[kind](q).order_by(Subject.name).all()
    titles = {
        'bio-obc': ('Biologie a občanská výchova', '🧬'),
        'matematika': ('Matematika', '➗'),
        'informatika': ('Informatika', '💻'),
    }
    title, icon = titles[kind]
    interactive_groups = interactive_groups_for(kind) if kind in ('matematika', 'informatika') else []
    informatics_lessons = InformaticsLesson.query.filter_by(is_published=True).order_by(
        InformaticsLesson.school, InformaticsLesson.grade_name, InformaticsLesson.topic, InformaticsLesson.title
    ).all() if kind == 'informatika' else []
    math_lessons = MathLesson.query.filter_by(is_published=True).order_by(
        MathLesson.school, MathLesson.grade_name, MathLesson.topic, MathLesson.title
    ).all() if kind == 'matematika' else []
    return render_template(
        'catalog.html',
        course={'subject': title, 'grade':'', 'block':'', 'icon':icon},
        lesson=None,
        subjects=subjects,
        interactive_groups=interactive_groups,
        informatics_lessons=informatics_lessons,
        math_lessons=math_lessons,
        kind=kind,
        title=title,
        icon=icon
    )



def _json_value(value, default):
    try:
        return json.loads(value) if value else default
    except Exception:
        return default


def export_html_lessons_backup():
    """Vyexportuje biologii a občanku z databáze do verzovaného JSON souboru."""
    exported = []
    for lesson_item in Lesson.query.order_by(Lesson.id).all():
        subject = lesson_item.block.grade.subject
        grade = lesson_item.block.grade
        block = lesson_item.block
        sections = []
        for sec in sorted(lesson_item.sections, key=lambda x: x.order):
            sections.append({
                'heading': sec.heading,
                'text': sec.text,
                'interest': sec.interest,
                'image': sec.image,
                'activity': sec.activity,
                'order': sec.order,
                'inline_images': [
                    {'file': img.file, 'caption': img.caption, 'order': img.order}
                    for img in sorted(sec.inline_images, key=lambda x: x.order)
                ]
            })
        questions = []
        for q in sorted(lesson_item.questions, key=lambda x: (x.area, x.order)):
            section_order = q.section.order if q.section else None
            questions.append({
                'section_order': section_order,
                'area': q.area,
                'qtype': q.qtype,
                'question': q.question,
                'options': _json_value(q.options_json, []),
                'correct': _json_value(q.correct_json, 0),
                'roots': _json_value(q.roots_json, []),
                'hint': q.hint,
                'order': q.order,
            })
        exported.append({
            'subject': {'name': subject.name, 'icon': subject.icon},
            'grade': grade.name,
            'block': {'title': block.title, 'order': block.order},
            'lesson': {
                'title': lesson_item.title,
                'tip': lesson_item.tip,
                'hero_image': lesson_item.hero_image,
                'order': lesson_item.order,
                'is_published': lesson_item.is_published,
            },
            'sections': sections,
            'questions': questions,
        })
    target = CONTENT_BACKUP / 'html_lessons.json'
    target.write_text(json.dumps({'version': 1, 'lessons': exported}, ensure_ascii=False, indent=2), encoding='utf-8')
    return target


def restore_html_lessons_backup():
    """Na novém Renderu obnoví chybějící HTML lekce ze souboru v GitHubu."""
    source = CONTENT_BACKUP / 'html_lessons.json'
    if not source.exists():
        return 0
    try:
        payload = json.loads(source.read_text(encoding='utf-8'))
    except Exception:
        return 0
    restored = 0
    for item in payload.get('lessons', []):
        sub_data = item.get('subject', {})
        subject_name = str(sub_data.get('name', '')).strip()
        grade_name = str(item.get('grade', '')).strip()
        block_data = item.get('block', {})
        block_title = str(block_data.get('title', '')).strip()
        lesson_data = item.get('lesson', {})
        title = str(lesson_data.get('title', '')).strip()
        if not all((subject_name, grade_name, block_title, title)):
            continue
        sub = Subject.query.filter_by(name=subject_name).first()
        if not sub:
            sub = Subject(name=subject_name, icon=sub_data.get('icon', '📘'))
            db.session.add(sub); db.session.flush()
        gr = Grade.query.filter_by(subject_id=sub.id, name=grade_name).first()
        if not gr:
            gr = Grade(subject_id=sub.id, name=grade_name)
            db.session.add(gr); db.session.flush()
        bl = Block.query.filter_by(grade_id=gr.id, title=block_title).first()
        if not bl:
            bl = Block(grade_id=gr.id, title=block_title, order=int(block_data.get('order', 1) or 1))
            db.session.add(bl); db.session.flush()
        existing = Lesson.query.filter_by(block_id=bl.id, title=title).first()
        if existing:
            continue
        les = Lesson(
            block_id=bl.id, title=title, tip=lesson_data.get('tip', ''),
            hero_image=lesson_data.get('hero_image', ''), order=int(lesson_data.get('order', 1) or 1),
            is_published=bool(lesson_data.get('is_published', True))
        )
        db.session.add(les); db.session.flush()
        section_by_order = {}
        for sec_data in item.get('sections', []):
            sec = Section(
                lesson_id=les.id, heading=sec_data.get('heading', 'Výklad'), text=sec_data.get('text', ''),
                interest=sec_data.get('interest', ''), image=sec_data.get('image', ''),
                activity=sec_data.get('activity', ''), order=int(sec_data.get('order', 1) or 1)
            )
            db.session.add(sec); db.session.flush()
            section_by_order[sec.order] = sec
            for img_data in sec_data.get('inline_images', []):
                db.session.add(InlineImage(
                    section_id=sec.id, file=img_data.get('file', ''), caption=img_data.get('caption', ''),
                    order=int(img_data.get('order', 1) or 1)
                ))
        for q_data in item.get('questions', []):
            sec = section_by_order.get(q_data.get('section_order'))
            db.session.add(Question(
                lesson_id=les.id, section_id=sec.id if sec else None,
                area=q_data.get('area', 'study'), qtype=q_data.get('qtype', 'choice'),
                question=q_data.get('question', ''),
                options_json=json.dumps(q_data.get('options', []), ensure_ascii=False),
                correct_json=json.dumps(q_data.get('correct', 0), ensure_ascii=False),
                roots_json=json.dumps(q_data.get('roots', []), ensure_ascii=False),
                hint=q_data.get('hint', ''), order=int(q_data.get('order', 1) or 1)
            ))
        restored += 1
    db.session.commit()
    return restored


def restore_interactive_lessons_from_files():
    restored = 0
    for meta_file in INTERACTIVE_LESSONS.glob('*/lesson.json'):
        try:
            meta = json.loads(meta_file.read_text(encoding='utf-8'))
            slug = safe_package_slug(meta.get('slug', meta_file.parent.name))
            if InteractiveLesson.query.filter_by(slug=slug).first():
                continue
            subject = normalize_subject(meta.get('subject', ''))
            if subject not in ('matematika', 'informatika'):
                continue
            db.session.add(InteractiveLesson(
                slug=slug, subject=subject, school=str(meta.get('school', '')).strip(),
                grade_name=str(meta.get('grade', '')).strip(), topic=str(meta.get('topic', '')).strip(),
                title=str(meta.get('title', slug)).strip(), description=str(meta.get('description', '')).strip(),
                icon=str(meta.get('icon', '➗' if subject == 'matematika' else '💻')).strip(),
                package_dir=str(meta_file.parent), is_published=bool(meta.get('is_published', True)),
                imported_at=datetime.utcnow()
            ))
            restored += 1
        except Exception:
            continue
    db.session.commit()
    return restored



@app.route('/teacher/interactive/import', methods=['GET', 'POST'])
def import_interactive_lesson():
    r = require_teacher()
    if r:
        return r

    if request.method == 'GET':
        return render_template(
            'import_interactive.html',
            course=course_from_lesson(None),
            lesson=None
        )

    package_file = (
        request.files.get('package')
        or request.files.get('zip_file')
        or request.files.get('lesson_zip')
        or request.files.get('file')
    )

    if not package_file or not package_file.filename:
        flash('Vyber ZIP balíček interaktivní lekce.')
        return redirect(url_for('import_interactive_lesson'))

    if not package_file.filename.lower().endswith('.zip'):
        flash('Balíček musí být ve formátu ZIP.')
        return redirect(url_for('import_interactive_lesson'))

    temp_root = Path(tempfile.mkdtemp(prefix='ucebnice_import_'))
    try:
        zip_path = temp_root / secure_filename(package_file.filename)
        package_file.save(zip_path)

        with zipfile.ZipFile(zip_path, 'r') as archive:
            extract_dir = temp_root / 'extracted'
            extract_dir.mkdir(parents=True, exist_ok=True)
            safe_extract_zip(archive, extract_dir)

        package_root = find_package_root(extract_dir)
        meta_file = package_root / 'lesson.json'
        meta = json.loads(meta_file.read_text(encoding='utf-8-sig'))

        subject = normalize_subject(meta.get('subject', ''))
        if subject not in ('matematika', 'informatika'):
            raise ValueError('V lesson.json musí být předmět matematika nebo informatika.')

        school = str(meta.get('school', '')).strip()
        grade_name = str(meta.get('grade', '')).strip()
        topic = str(meta.get('topic', '')).strip()
        title = str(meta.get('title', '')).strip()

        if not all((school, grade_name, topic, title)):
            raise ValueError(
                'V lesson.json musí být vyplněno school, grade, topic a title.'
            )

        slug = safe_package_slug(meta.get('slug') or title)
        if not slug:
            raise ValueError('Nepodařilo se vytvořit platný název lekce.')

        existing = InteractiveLesson.query.filter_by(slug=slug).first()
        if existing:
            raise ValueError(f'Interaktivní lekce se slugem „{slug}“ už existuje.')

        if not (package_root / 'templates' / 'index.html').exists():
            raise ValueError('Balíček musí obsahovat templates/index.html.')

        if not (package_root / 'lesson_app.py').exists():
            raise ValueError('Balíček musí obsahovat lesson_app.py.')

        destination = INTERACTIVE_LESSONS / slug
        if destination.exists():
            shutil.rmtree(destination)
        shutil.copytree(package_root, destination)

        item = InteractiveLesson(
            slug=slug,
            subject=subject,
            school=school,
            grade_name=grade_name,
            topic=topic,
            title=title,
            description=str(meta.get('description', '')).strip(),
            icon=str(
                meta.get(
                    'icon',
                    '➗' if subject == 'matematika' else '💻'
                )
            ).strip(),
            package_dir=str(destination),
            is_published=bool(meta.get('is_published', True)),
            imported_at=datetime.utcnow()
        )
        db.session.add(item)
        db.session.commit()

        flash(f'Interaktivní lekce „{title}“ byla úspěšně importována.')
        return redirect(url_for('teacher_home'))

    except zipfile.BadZipFile:
        db.session.rollback()
        flash('Soubor není platný ZIP balíček.')
    except (ValueError, json.JSONDecodeError) as exc:
        db.session.rollback()
        flash(str(exc))
    except Exception as exc:
        db.session.rollback()
        flash(f'Import se nepodařil: {exc}')
    finally:
        shutil.rmtree(temp_root, ignore_errors=True)

    return redirect(url_for('import_interactive_lesson'))


@app.route('/interactive/<slug>')
def interactive_lesson(slug):
    r = require_login()
    if r:
        return r
    lesson_item = InteractiveLesson.query.filter_by(slug=slug, is_published=True).first()
    if not lesson_item:
        return 'Interaktivní lekce nebyla nalezena.', 404

    if current_user().role == 'student':
        begin_focus_attempt('interactive', slug)

    template_file = BASE / lesson_item.package_dir / 'templates' / 'index.html'
    if not template_file.exists():
        return 'Balíček lekce neobsahuje templates/index.html.', 500

    html_source = template_file.read_text(encoding='utf-8')
    if current_user().role == 'student':
        guard = render_template_string(
            '<script>window.UCEBNICE_FOCUS_GUARD={{ cfg|tojson }};</script>'
            "<script src=\"{{ url_for('static', filename='js/focus_guard.js') }}\"></script>",
            cfg={'kind': 'interactive', 'key': slug}
        )
        if '</body>' in html_source.lower():
            pos = html_source.lower().rfind('</body>')
            html_source = html_source[:pos] + guard + html_source[pos:]
        else:
            html_source += guard

    return render_template_string(
        html_source,
        package=lesson_item,
        lesson=lesson_item,
        user=current_user(),
        asset_url=lambda path: url_for('interactive_asset', slug=slug, filename=path),
        api_url=lambda action: url_for('interactive_api', slug=slug, action=action),
        complete_url=url_for('complete_interactive', slug=slug),
        portal_url=url_for('subject_catalog', kind=lesson_item.subject)
    )


@app.route('/interactive-assets/<slug>/<path:filename>')
def interactive_asset(slug, filename):
    r = require_login()
    if r:
        return r
    lesson_item = InteractiveLesson.query.filter_by(slug=slug).first_or_404()
    return send_from_directory(BASE / lesson_item.package_dir / 'static', filename)


@app.route('/interactive/<slug>/api/<action>', methods=['GET', 'POST'])
def interactive_api(slug, action):
    r = require_login()
    if r:
        return jsonify({'ok': False, 'error': 'login'}), 401

    lesson_item = InteractiveLesson.query.filter_by(slug=slug, is_published=True).first()
    if not lesson_item:
        return jsonify({'ok': False, 'error': 'lesson'}), 404

    module = load_interactive_module(lesson_item)
    if not module or not hasattr(module, 'handle'):
        return jsonify({'ok': False, 'error': 'Balíček neobsahuje lesson_app.py s funkcí handle().'}), 500

    payload = request.get_json(silent=True) if request.is_json else request.form.to_dict()
    payload = payload or {}
    try:
        result = module.handle(
            action=action,
            payload=payload,
            session=session,
            user={'id': current_user().id, 'username': current_user().username, 'name': current_user().name}
        )
        if not isinstance(result, dict):
            raise ValueError('Funkce handle() musí vrátit slovník.')
        session.modified = True
        return jsonify(result)
    except Exception as exc:
        return jsonify({'ok': False, 'error': str(exc)}), 400


@app.route('/interactive/<slug>/complete', methods=['POST'])
def complete_interactive(slug):
    r = require_login()
    if r:
        return jsonify({'ok': False, 'error': 'login'}), 401

    lesson_item = InteractiveLesson.query.filter_by(slug=slug, is_published=True).first_or_404()
    data = request.get_json(silent=True) or request.form.to_dict()
    try:
        percent = max(0, min(100, int(data.get('percent', 100))))
        grade = max(1, min(5, int(data.get('grade', grade_from_percent(percent)))))
    except (TypeError, ValueError):
        return jsonify({'ok': False, 'error': 'Neplatný výsledek.'}), 400

    focus_lost = consume_focus_count('interactive', slug)
    upsert_interactive_progress(lesson_item, percent=percent, grade=grade, focus_lost=focus_lost)
    return jsonify({
        'ok': True,
        'message': 'Dokončení lekce bylo uloženo.',
        'percent': percent,
        'grade': grade
    })


@app.route('/teacher/interactive/<int:lesson_id>/delete', methods=['POST'])
def delete_interactive_lesson(lesson_id):
    r = require_teacher()
    if r:
        return r
    item = db.session.get(InteractiveLesson, lesson_id)
    if item:
        InteractiveResult.query.filter_by(interactive_lesson_id=item.id).delete()
        InteractiveProgress.query.filter_by(interactive_lesson_id=item.id).delete()
        package_path = BASE / item.package_dir
        db.session.delete(item)
        db.session.commit()
        shutil.rmtree(package_path, ignore_errors=True)
        flash('Interaktivní lekce byla odstraněna.')
    return redirect(url_for('teacher_home'))


@app.route('/dashboard')
def dashboard():
    r = require_login()
    if r: return r
    subjects = Subject.query.order_by(Subject.name).all()
    lessons = [lesson_to_dict(l) for l in visible_lessons()]
    return render_template('dashboard.html', course=course_from_lesson(None), subjects=subjects, lessons=lessons, lesson=None)

@app.route('/lesson/<int:lesson_id>')
def lesson(lesson_id):
    r=require_login();
    if r: return r
    lesson = db.session.get(Lesson, lesson_id)
    if not lesson: return 'Lekce nenalezena', 404
    if current_user().role == 'student':
        # Přepínač CZ/EN musí zůstat dostupný i během práce v lekci.
        unlock_language()
        consume_pending_lesson_reset(lesson.id)
        begin_focus_attempt('html', lesson.id)
    step = int(request.args.get('step',0))
    data = lesson_to_dict(lesson)
    step = max(0, min(step, len(data['sections'])-1))
    related = Lesson.query.filter_by(block_id=lesson.block_id, is_published=True).order_by(Lesson.order).all()
    completed_steps = completed_steps_for(lesson.id)
    study_completed_questions = set()
    study_completed_activities = set()
    section_read_complete = False
    user = current_user()
    if user and user.role == 'student' and data.get('sections'):
        sec_id = data['sections'][step]['id']
        study_completed_questions = {r.question_id for r in StudyQuestionProgress.query.filter_by(user_id=user.id, lesson_id=lesson.id, completed=True).all()}
        study_completed_activities = {r.activity_id for r in StudentActivityProgress.query.filter_by(user_id=user.id, context='study', completed=True).all()}
        spr = StudentSectionProgress.query.filter_by(user_id=user.id, lesson_id=lesson.id, section_id=sec_id).first()
        section_read_complete = bool(spr and spr.read_complete)
    touch_progress(lesson.id, step, 'rozpracováno')
    return render_template('lesson.html', lesson=data, lessons=[lesson_to_dict(l) for l in related], course=course_from_lesson(lesson), step=step, completed_steps=completed_steps, ready_for_test=lesson_ready_for_test(lesson), study_completed_questions=study_completed_questions, study_completed_activities=study_completed_activities, section_read_complete=section_read_complete)

@app.route('/test/<int:lesson_id>')
def final_test(lesson_id):
    r=require_login();
    if r: return r
    lesson = db.session.get(Lesson, lesson_id)
    if not lesson: return 'Lekce nenalezena', 404
    if current_user().role == 'student':
        consume_pending_lesson_reset(lesson.id)
    if not lesson_ready_for_test(lesson):
        flash('Nejdřív dokonči otázky k výkladu a aktivitu. Test se odemkne až potom.')
        return redirect(url_for('lesson', lesson_id=lesson.id))
    if current_user().role == 'student':
        # Každé nové otevření závěrečného úkolu je nový pokus od 0 %.
        # Pokud student obnoví stránku nebo se vrátí do testu, předchozí
        # rozpracovaný pokus se nejprve uloží učiteli jako poslední pokus.
        if final_attempt_is_active(lesson.id):
            persist_final_result(lesson, status='pokus přerušen – zahájen nový', focus_lost=get_focus_count('html', lesson.id))
            final_attempt_mark_finished(lesson.id)
        consume_focus_count('html', lesson.id)
        reset_final_attempt(current_user().id, lesson.id)
        final_attempt_mark_active(lesson.id)
        begin_focus_attempt('html', lesson.id)
    related = Lesson.query.filter_by(block_id=lesson.block_id, is_published=True).order_by(Lesson.order).all()
    touch_progress(lesson.id, 999, 'závěrečný test – nový pokus')
    return render_template('test.html', lesson=lesson_to_dict(lesson), lessons=[lesson_to_dict(l) for l in related], course=course_from_lesson(lesson))

@app.route('/final-abort/<int:lesson_id>')
def final_abort(lesson_id):
    r = require_login()
    if r: return r
    lesson = db.session.get(Lesson, lesson_id)
    if not lesson: return 'Lekce nenalezena', 404
    if current_user().role == 'student' and final_attempt_is_active(lesson.id):
        persist_final_result(lesson, status='pokus ukončen studentem', focus_lost=get_focus_count('html', lesson.id))
        final_attempt_mark_finished(lesson.id)
        consume_focus_count('html', lesson.id)
    if request.args.get('to') == 'lesson':
        return redirect(url_for('lesson', lesson_id=lesson.id, review=1))
    return redirect(url_for('dashboard'))


@app.route('/finish/<int:lesson_id>', methods=['POST'])
def finish(lesson_id):
    r=require_login();
    if r: return r
    lesson = db.session.get(Lesson, lesson_id)
    data = lesson_to_dict(lesson)
    total = len(data['final_test']); score=0; detail=[]
    for q in data['final_test']:
        ans = request.form.get(f'q{q["id"]}','')
        ok = check_question(q, ans)
        if ok: score += 1
        detail.append(ok)
    percent = round(score/max(total,1)*100); grade = grade_from_percent(percent)
    focus_lost = consume_focus_count('html', lesson.id)
    Result.query.filter_by(user_id=current_user().id, lesson_id=lesson.id).filter(Result.status != 'dokončeno').delete(synchronize_session=False)
    partial = session.get('html_partial_progress', {})
    partial.pop(str(lesson.id), None)
    session['html_partial_progress'] = partial
    session.modified = True
    db.session.add(Result(user_id=current_user().id, lesson_id=lesson.id, percent=percent, grade=grade, score=score, total=total, focus_lost=focus_lost, status='dokončeno')); db.session.commit()
    touch_progress(lesson.id, 1000, 'dokončeno')
    unlock_language()
    return render_template('finish.html', lesson=data, course=course_from_lesson(lesson), score=score, total=total, percent=percent, grade=grade, detail=detail)

def check_question(q, ans):
    if q.get('type') in ['choice','image_choice']:
        return str(ans) == str(q.get('correct'))
    na = strip_accents(ans)
    return any(strip_accents(root) in na for root in q.get('roots',[]))

@app.route('/api/check', methods=['POST'])
def api_check():
    d = request.get_json(force=True)
    return jsonify({'ok': check_question(d.get('question',{}), d.get('answer',''))})

@app.route('/api/study-question-check', methods=['POST'])
def api_study_question_check():
    r = require_login()
    if r: return jsonify({'ok': False, 'error': 'login'}), 401
    d = request.get_json(silent=True) or {}
    q = db.session.get(Question, int(d.get('question_id', 0) or 0))
    if not q: return jsonify({'ok': False, 'error': 'question'}), 404
    answer = d.get('answer', '')
    ok = check_question(q_to_dict(q), answer)
    user = current_user()
    if user.role == 'student':
        row = StudyQuestionProgress.query.filter_by(user_id=user.id, lesson_id=q.lesson_id, question_id=q.id).first()
        if not row:
            row = StudyQuestionProgress(user_id=user.id, lesson_id=q.lesson_id, question_id=q.id)
            db.session.add(row)
        row.attempts = int(row.attempts or 0) + 1
        row.answer_json = json.dumps(answer, ensure_ascii=False)
        row.completed = bool(row.completed or ok)
        row.updated_at = datetime.utcnow()
        db.session.commit()
    return jsonify({'ok': ok, 'message': 'Správně, můžeš pokračovat.' if ok else 'Ještě ne. Zkus to znovu – vše potřebné najdeš ve studijním materiálu.'})




def _point_in_zone(x, y, zone):
    try:
        x, y = float(x), float(y)
    except Exception:
        return False
    if not isinstance(zone, dict):
        return False
    shape = str(zone.get('shape', 'rect'))
    zx, zy = float(zone.get('x', 0)), float(zone.get('y', 0))
    zw, zh = float(zone.get('w', 0)), float(zone.get('h', 0))
    if shape in ('circle', 'oval'):
        if zw <= 0 or zh <= 0:
            return False
        cx, cy = zx + zw / 2, zy + zh / 2
        return ((x-cx)/(zw/2))**2 + ((y-cy)/(zh/2))**2 <= 1
    return zx <= x <= zx + zw and zy <= y <= zy + zh


def check_practical_activity(activity, answer):
    try:
        raw_cfg = activity.config_en_json if current_lang() == 'en' and (activity.config_en_json or '').strip() else activity.config_json
        cfg = json.loads(raw_cfg or '{}')
    except Exception:
        cfg = {}
    answer = answer or {}
    typ = activity.activity_type
    if typ == 'find_image':
        return _point_in_zone(answer.get('x'), answer.get('y'), cfg.get('zone', {}))
    if typ == 'video_find':
        try:
            target = float(cfg.get('time', 0)); tol = max(.15, float(cfg.get('tolerance', .8)))
            actual = float(answer.get('time', -999))
        except Exception:
            return False
        return abs(actual-target) <= tol and _point_in_zone(answer.get('x'), answer.get('y'), cfg.get('zone', {}))
    if typ == 'video_observe':
        return str(answer.get('selected')) == str(cfg.get('correct', 0))
    if typ == 'sort':
        expected = {str(i): str(item.get('category')) for i, item in enumerate(cfg.get('items', []))}
        got = {str(k): str(v) for k, v in (answer.get('assignments') or {}).items()}
        return bool(expected) and all(got.get(k) == v for k, v in expected.items())
    if typ == 'cards':
        if cfg.get('mode') == 'categories':
            expected = {str(i): str(item.get('category')) for i, item in enumerate(cfg.get('items', []))}
            got = {str(k): str(v) for k, v in (answer.get('assignments') or {}).items()}
            return bool(expected) and all(got.get(k) == v for k, v in expected.items())
        # Obrázkové kartičky: každý nahraný obrázek má svůj očíslovaný cíl.
        # Student obrázky nevidí očíslované; přiřazuje je k názvům/cílům 1..N.
        assignments = answer.get('assignments') or {}
        cards = cfg.get('cards', [])
        if not cards:
            return False
        return all(str(assignments.get(str(i), assignments.get(i, ''))) == str(i) for i in range(len(cards)))
    if typ == 'real_world':
        answers = [strip_accents(x) for x in (answer.get('answers') or []) if str(x).strip()]
        min_items = max(1, int(cfg.get('min_items', 3) or 3))
        concepts = cfg.get('concepts') or []
        matched = set()
        for text_answer in answers:
            for idx, concept in enumerate(concepts):
                roots = concept if isinstance(concept, list) else [x.strip() for x in str(concept).split('|') if x.strip()]
                if any(strip_accents(root) in text_answer for root in roots if str(root).strip()):
                    matched.add(idx)
        return len(matched) >= min_items
    return False


def _upsert_activity_progress(activity, context, answer, ok):
    user = current_user()
    if not user or user.role != 'student':
        return
    row = StudentActivityProgress.query.filter_by(user_id=user.id, activity_id=activity.id, context=context).first()
    if not row:
        row = StudentActivityProgress(user_id=user.id, activity_id=activity.id, context=context)
        db.session.add(row)
    row.attempts = int(row.attempts or 0) + 1
    row.answer_json = json.dumps(answer or {}, ensure_ascii=False)
    row.completed = bool(row.completed or ok)
    row.updated_at = datetime.utcnow()
    db.session.commit()


def final_item_keys(lesson):
    data = lesson_to_dict(lesson)
    return [f'q:{q["id"]}' for q in data.get('final_test', [])] + [f'a:{a["id"]}' for a in data.get('final_activities', [])]


def final_attempt_mark_active(lesson_id):
    active = set(str(x) for x in session.get('final_active_attempts', []))
    active.add(str(lesson_id))
    session['final_active_attempts'] = sorted(active)
    session.modified = True


def final_attempt_mark_finished(lesson_id):
    active = set(str(x) for x in session.get('final_active_attempts', []))
    active.discard(str(lesson_id))
    session['final_active_attempts'] = sorted(active)
    session.modified = True


def final_attempt_is_active(lesson_id):
    return str(lesson_id) in {str(x) for x in session.get('final_active_attempts', [])}


def reset_final_attempt(user_id, lesson_id):
    """Nový závěrečný pokus vždy začíná od nuly.

    Maže se pouze pracovní stav závěrečného pokusu; studijní část lekce
    a poslední učitelův výsledek zůstávají zachované.
    """
    FinalItemProgress.query.filter_by(user_id=user_id, lesson_id=lesson_id).delete(synchronize_session=False)
    activity_ids = [a.id for a in PracticalActivity.query.filter_by(lesson_id=lesson_id).all()]
    if activity_ids:
        StudentActivityProgress.query.filter(
            StudentActivityProgress.user_id == user_id,
            StudentActivityProgress.context == 'final',
            StudentActivityProgress.activity_id.in_(activity_ids)
        ).delete(synchronize_session=False)
    db.session.commit()


def final_progress_snapshot(lesson, user=None):
    user = user or current_user()
    if not user or user.role != 'student':
        return {'percent': 0, 'grade': 5, 'completed': 0, 'total': 0, 'label': slovni_hodnoceni(0), 'done': False}
    keys = final_item_keys(lesson)
    completed_keys = {r.item_key for r in FinalItemProgress.query.filter_by(
        user_id=user.id, lesson_id=lesson.id, completed=True
    ).all()}
    completed = sum(1 for k in keys if k in completed_keys)
    total = len(keys)
    percent = round(completed / max(total, 1) * 100)
    grade = grade_from_percent(percent)
    return {
        'percent': percent, 'grade': grade, 'completed': completed, 'total': total,
        'label': slovni_hodnoceni(percent), 'done': bool(total and completed >= total)
    }


def final_next_item_key(lesson, user=None):
    user = user or current_user()
    if not user or user.role != 'student':
        return None
    completed = {r.item_key for r in FinalItemProgress.query.filter_by(
        user_id=user.id, lesson_id=lesson.id, completed=True
    ).all()}
    for key in final_item_keys(lesson):
        if key not in completed:
            return key
    return None


def persist_final_result(lesson, status='dokončeno', focus_lost=None):
    """Učitelovi se uchovává nejlepší dosažený výsledek závěrečného úkolu.

    Student může závěrečný úkol opakovat libovolně. Každý nový pokus začíná
    od nuly, ale horší pozdější pokus nikdy nepřepíše lepší dřívější výsledek.
    """
    user = current_user()
    if not user or user.role != 'student':
        return final_progress_snapshot(lesson, user)
    progress = final_progress_snapshot(lesson, user)
    if focus_lost is None:
        focus_lost = get_focus_count('html', lesson.id)

    rows = Result.query.filter_by(user_id=user.id, lesson_id=lesson.id).order_by(Result.created_at.desc()).all()
    row = rows[0] if rows else None
    for old in rows[1:]:
        db.session.delete(old)

    # První ukončený pokus vždy uložíme. Další pokus přepíše učitelův
    # výsledek pouze tehdy, pokud je lepší než dosavadní maximum.
    should_store = row is None or int(progress['percent'] or 0) > int(row.percent or 0)
    if row is None:
        row = Result(user_id=user.id, lesson_id=lesson.id)
        db.session.add(row)

    if should_store:
        row.percent = progress['percent']
        row.grade = progress['grade']
        row.score = progress['completed']
        row.total = progress['total']
        row.focus_lost = int(focus_lost or 0)
        row.status = status
        row.created_at = datetime.utcnow()

    db.session.commit()
    touch_progress(lesson.id, 1000 if status == 'dokončeno' else 999, status)
    return progress


def update_final_result(lesson):
    # Průběžný stav závěrečného pokusu je pracovní stav.
    # Do učitelovy databáze se zapíše až ukončený pokus.
    progress = final_progress_snapshot(lesson)
    if progress.get('done'):
        persist_final_result(lesson, status='dokončeno', focus_lost=get_focus_count('html', lesson.id))
        final_attempt_mark_finished(lesson.id)
        end_focus_attempt('html', lesson.id)
    return progress

def mark_final_item(lesson_id, item_key, answer, ok):
    user = current_user()
    if not user or user.role != 'student':
        return
    row = FinalItemProgress.query.filter_by(user_id=user.id, lesson_id=lesson_id, item_key=item_key).first()
    if not row:
        row = FinalItemProgress(user_id=user.id, lesson_id=lesson_id, item_key=item_key)
        db.session.add(row)
    row.attempts = int(row.attempts or 0) + 1
    row.answer_json = json.dumps(answer, ensure_ascii=False)
    row.completed = bool(row.completed or ok)
    row.updated_at = datetime.utcnow()
    db.session.commit()


@app.route('/api/activity-check', methods=['POST'])
def api_activity_check():
    r = require_login()
    if r: return jsonify({'ok': False, 'error': 'login'}), 401
    d = request.get_json(silent=True) or {}
    activity = db.session.get(PracticalActivity, int(d.get('activity_id', 0) or 0))
    if not activity:
        return jsonify({'ok': False, 'error': 'activity'}), 404
    context = 'final' if d.get('context') == 'final' else 'study'
    answer = d.get('answer') or {}
    if context == 'final' and current_user().role == 'student':
        expected = final_next_item_key(activity.lesson)
        current_key = f'a:{activity.id}'
        if expected != current_key:
            return jsonify({'ok': False, 'blocked': True, 'message': 'Nejdřív správně dokonči předchozí otázku nebo úkol.'}), 409
    ok = check_practical_activity(activity, answer)
    _upsert_activity_progress(activity, context, answer, ok)
    progress = None
    if context == 'final':
        mark_final_item(activity.lesson_id, f'a:{activity.id}', answer, ok)
        progress = update_final_result(activity.lesson)
    return jsonify({'ok': ok, 'message': 'Správně, můžeš pokračovat.' if ok else 'Ještě ne. Zkus to znovu – vše potřebné najdeš ve studijním materiálu.', 'progress': progress})

@app.route('/api/final-question-check', methods=['POST'])
def api_final_question_check():
    r = require_login()
    if r: return jsonify({'ok': False, 'error': 'login'}), 401
    d = request.get_json(silent=True) or {}
    q = db.session.get(Question, int(d.get('question_id', 0) or 0))
    if not q:
        return jsonify({'ok': False, 'error': 'question'}), 404
    lesson = q.lesson
    if current_user().role == 'student':
        expected = final_next_item_key(lesson)
        current_key = f'q:{q.id}'
        if expected != current_key:
            return jsonify({'ok': False, 'blocked': True, 'message': 'Nejdřív správně dokonči předchozí otázku nebo úkol.'}), 409
    answer = d.get('answer', '')
    ok = check_question(q_to_dict(q), answer)
    mark_final_item(lesson.id, f'q:{q.id}', answer, ok)
    progress = update_final_result(lesson)
    return jsonify({'ok': ok, 'message': 'Správně. Úkol je splněný.' if ok else 'Ještě ne. Můžeš to zkusit znovu nebo se podívat do studijního materiálu.', 'progress': progress})


@app.route('/api/final-status/<int:lesson_id>')
def api_final_status(lesson_id):
    r = require_login()
    if r: return jsonify({'ok': False, 'error': 'login'}), 401
    lesson = db.session.get(Lesson, lesson_id)
    if not lesson: return jsonify({'ok': False, 'error': 'lesson'}), 404
    progress = update_final_result(lesson)
    completed = {r.item_key for r in FinalItemProgress.query.filter_by(user_id=current_user().id, lesson_id=lesson.id, completed=True).all()} if current_user().role == 'student' else set()
    return jsonify({'ok': True, 'progress': progress, 'completed': sorted(completed)})


@app.route('/api/section-read', methods=['POST'])
def api_section_read():
    r = require_login()
    if r: return jsonify({'ok': False, 'error': 'login'}), 401
    d = request.get_json(silent=True) or {}
    lesson = db.session.get(Lesson, int(d.get('lesson_id', 0) or 0))
    step = int(d.get('step', 0) or 0)
    if not lesson: return jsonify({'ok': False, 'error': 'lesson'}), 404
    sections = sorted(lesson.sections, key=lambda x: x.order)
    if step < 0 or step >= len(sections): return jsonify({'ok': False, 'error': 'step'}), 400
    user = current_user()
    if user.role == 'student':
        sec = sections[step]
        row = StudentSectionProgress.query.filter_by(user_id=user.id, lesson_id=lesson.id, section_id=sec.id).first()
        if not row:
            row = StudentSectionProgress(user_id=user.id, lesson_id=lesson.id, section_id=sec.id)
            db.session.add(row)
        row.read_complete = True
        row.updated_at = datetime.utcnow()
        db.session.commit()
    return jsonify({'ok': True})


def save_html_partial_result(lesson, status='rozpracováno'):
    user = current_user()
    if user and user.role == 'student':
        consume_pending_lesson_reset(lesson.id)
        user = current_user()
    if not user or user.role != 'student':
        return None
    data = lesson_to_dict(lesson)
    progress_map = session.get('html_partial_progress', {})
    lesson_map = progress_map.get(str(lesson.id), {})
    total = 0
    score = 0
    for idx, section in enumerate(data.get('sections', [])):
        units = len(section.get('questions', [])) + 1  # + aktivita
        total += units
        saved = lesson_map.get(str(idx), {})
        score += min(units, int(saved.get('questions', 0)) + (1 if saved.get('activity') else 0))
    final_total = len(data.get('final_test', []))
    total += final_total
    final_saved = lesson_map.get('_final_test', {})
    score += min(final_total, int(final_saved.get('answered', 0)))
    percent = round(score / max(total, 1) * 100)
    row = Result.query.filter_by(user_id=user.id, lesson_id=lesson.id).filter(Result.status != 'dokončeno').order_by(Result.created_at.desc()).first()
    if not row:
        row = Result(user_id=user.id, lesson_id=lesson.id, created_at=datetime.utcnow())
        db.session.add(row)
    row.percent = percent
    row.grade = grade_from_percent(percent)
    row.score = score
    row.total = total
    row.focus_lost = get_focus_count('html', lesson.id)
    row.status = status
    row.created_at = datetime.utcnow()
    db.session.commit()
    touch_progress(lesson.id, 0, status)
    return row

@app.route('/api/html-progress', methods=['POST'])
def api_html_progress():
    r = require_login()
    if r:
        return jsonify({'ok': False, 'error': 'login'}), 401
    d = request.get_json(silent=True) or {}
    lesson_id = int(d.get('lesson_id', 0))
    step = int(d.get('step', 0))
    lesson = db.session.get(Lesson, lesson_id)
    if not lesson:
        return jsonify({'ok': False, 'error': 'lesson'}), 404
    data = lesson_to_dict(lesson)
    if step < 0 or step >= len(data.get('sections', [])):
        return jsonify({'ok': False, 'error': 'step'}), 400
    q_total = len(data['sections'][step].get('questions', []))
    q_done = max(0, min(q_total, int(d.get('questions', 0))))
    activity = bool(d.get('activity', False))
    progress_map = session.get('html_partial_progress', {})
    lesson_map = progress_map.setdefault(str(lesson_id), {})
    old = lesson_map.get(str(step), {})
    lesson_map[str(step)] = {
        'questions': max(int(old.get('questions', 0)), q_done),
        'activity': bool(old.get('activity')) or activity,
    }
    session['html_partial_progress'] = progress_map
    session.modified = True
    row = save_html_partial_result(lesson, str(d.get('status') or 'rozpracováno'))
    return jsonify({'ok': True, 'percent': row.percent if row else 0})

@app.route('/api/html-test-progress', methods=['POST'])
def api_html_test_progress():
    r = require_login()
    if r:
        return jsonify({'ok': False, 'error': 'login'}), 401
    d = request.get_json(silent=True) or {}
    lesson_id = int(d.get('lesson_id', 0))
    lesson = db.session.get(Lesson, lesson_id)
    if not lesson:
        return jsonify({'ok': False, 'error': 'lesson'}), 404
    data = lesson_to_dict(lesson)
    total = len(data.get('final_test', []))
    answered = max(0, min(total, int(d.get('answered', 0))))
    progress_map = session.get('html_partial_progress', {})
    lesson_map = progress_map.setdefault(str(lesson_id), {})
    old = lesson_map.get('_final_test', {})
    lesson_map['_final_test'] = {
        'answered': max(int(old.get('answered', 0)), answered)
    }
    session['html_partial_progress'] = progress_map
    session.modified = True
    row = save_html_partial_result(lesson, str(d.get('status') or 'závěrečný test – rozpracováno'))
    return jsonify({'ok': True, 'percent': row.percent if row else 0})

@app.route('/api/section-complete', methods=['POST'])
def api_section_complete():
    r = require_login()
    if r: return jsonify({'ok': False, 'error': 'login'}), 401
    d = request.get_json(force=True)
    lesson_id = int(d.get('lesson_id', 0))
    step = int(d.get('step', 0))
    lesson = db.session.get(Lesson, lesson_id)
    if not lesson:
        return jsonify({'ok': False, 'error': 'lesson'}), 404
    data = lesson_to_dict(lesson)
    if step < 0 or step >= len(data['sections']):
        return jsonify({'ok': False, 'error': 'step'}), 400
    user = current_user()
    sec = sorted(lesson.sections, key=lambda x: x.order)[step]
    if user and user.role == 'student':
        read_row = StudentSectionProgress.query.filter_by(user_id=user.id, lesson_id=lesson.id, section_id=sec.id).first()
        if not read_row or not read_row.read_complete:
            return jsonify({'ok': False, 'error': 'read', 'message': 'Nejdřív projdi celý studijní materiál.'}), 400
        q_ids = {q.id for q in sec.questions if q.area == 'study'}
        done_q = {r.question_id for r in StudyQuestionProgress.query.filter_by(user_id=user.id, lesson_id=lesson.id, completed=True).all()}
        if not q_ids.issubset(done_q):
            return jsonify({'ok': False, 'error': 'questions', 'message': 'Ještě nejsou splněné všechny otázky.'}), 400
        act_ids = {a.id for a in sec.practical_activities}
        done_a = {r.activity_id for r in StudentActivityProgress.query.filter_by(user_id=user.id, context='study', completed=True).all()}
        if not act_ids.issubset(done_a):
            return jsonify({'ok': False, 'error': 'activities', 'message': 'Ještě nejsou splněné všechny praktické aktivity.'}), 400
    mark_step_complete(lesson_id, step)
    touch_progress(lesson_id, step, 'splněná část lekce')
    return jsonify({'ok': True, 'ready_for_test': lesson_ready_for_test(lesson)})


# ============================================================
# INFORMATIKA – UNIVERZÁLNÍ ENGINE
# Učitel nahraje hotový referenční soubor pouze pro analýzu.
# Student jej nikdy nedostane a vytváří vlastní nový soubor od nuly.
# ============================================================

INFORMATICS_SOURCE_DIR = DATA_DIR / 'informatics_sources'
INFORMATICS_SUBMISSION_DIR = DATA_DIR / 'informatics_submissions'
INFORMATICS_SOURCE_DIR.mkdir(parents=True, exist_ok=True)
INFORMATICS_SUBMISSION_DIR.mkdir(parents=True, exist_ok=True)


def informatics_grade_from_percent(percent):
    percent = int(percent or 0)
    if percent >= 90: return 1
    if percent >= 75: return 2
    if percent >= 60: return 3
    if percent >= 40: return 4
    return 5


def render_informatics_html(filename):
    if not filename:
        return ''
    path = INFORMATICS_SOURCE_DIR / filename
    if not path.exists():
        return ''
    try:
        return path.read_text(encoding='utf-8', errors='replace')
    except Exception:
        return ''


def _safe_json(value, default):
    try:
        return json.loads(value or '')
    except Exception:
        return default


def _save_uploaded_file(file_storage, folder, prefix='file'):
    ext = Path(file_storage.filename or '').suffix.lower()
    name = f'{prefix}_{uuid.uuid4().hex}{ext}'
    folder.mkdir(parents=True, exist_ok=True)
    file_storage.save(folder / name)
    return name


def analyze_informatics_file(path, original_name):
    """Vrátí strukturu souboru a seznam kontrol, které lze studentovi nabídnout."""
    ext = Path(original_name).suffix.lower()
    info = {'extension': ext, 'name': original_name}
    checks = []

    if ext in ('.xlsx', '.xlsm'):
        info['type'] = 'Excel'
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=False)
            info['sheets'] = wb.sheetnames
            total_nonempty = 0
            total_formulas = 0
            functions = set()
            chart_count = 0
            sheet_specs = []
            for ws in wb.worksheets:
                nonempty = 0
                formulas = []
                max_row = min(ws.max_row or 1, 500)
                max_col = min(ws.max_column or 1, 50)
                headers = []
                for c in ws[1][:max_col]:
                    if c.value not in (None, ''):
                        headers.append(str(c.value))
                for row in ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col):
                    for cell in row:
                        if cell.value not in (None, ''):
                            nonempty += 1
                        if isinstance(cell.value, str) and cell.value.startswith('='):
                            total_formulas += 1
                            formulas.append({'cell': cell.coordinate, 'formula': cell.value})
                            for fn in re.findall(r'([A-ZÁ-Ž][A-ZÁ-Ž0-9_.]*)\s*\(', cell.value.upper()):
                                functions.add(fn)
                total_nonempty += nonempty
                chart_count += len(getattr(ws, '_charts', []) or [])
                sheet_specs.append({
                    'name': ws.title,
                    'rows': ws.max_row or 0,
                    'cols': ws.max_column or 0,
                    'headers': headers,
                    'nonempty': nonempty,
                    'formulas': formulas[:40]
                })
            info.update({
                'sheet_specs': sheet_specs,
                'nonempty_count': total_nonempty,
                'formula_count': total_formulas,
                'formula_functions': sorted(functions),
                'chart_count': chart_count,
            })
        except Exception as exc:
            info['analysis_error'] = str(exc)

        checks = [
            {'code':'excel_sheets','label':'Počet a názvy listů'},
            {'code':'excel_headers','label':'Záhlaví tabulky'},
            {'code':'excel_size','label':'Rozsah tabulky / počet řádků a sloupců'},
            {'code':'excel_filled','label':'Vyplněné části tabulky'},
            {'code':'excel_formulas','label':'Použití vzorců'},
            {'code':'excel_functions','label':'Použití konkrétních funkcí (např. SUM, IF, AVERAGE)'},
            {'code':'excel_chart','label':'Graf, pokud je v učitelském souboru'},
        ]

    elif ext == '.docx':
        info['type'] = 'Word'
        try:
            from docx import Document
            doc = Document(path)
            paragraphs = [p for p in doc.paragraphs if p.text.strip()]
            headings = [p.text.strip() for p in paragraphs if str(p.style.name).lower().startswith(('heading','nadpis'))]
            words = [w for p in paragraphs for w in p.text.split()]
            info.update({
                'paragraph_count': len(paragraphs),
                'word_count': len(words),
                'heading_count': len(headings),
                'headings': headings[:30],
                'table_count': len(doc.tables),
                'image_count': len(doc.inline_shapes),
            })
        except Exception as exc:
            info['analysis_error'] = str(exc)

        checks = [
            {'code':'word_length','label':'Rozsah dokumentu'},
            {'code':'word_headings','label':'Nadpisy'},
            {'code':'word_table','label':'Tabulky'},
            {'code':'word_images','label':'Obrázky'},
        ]

    elif ext == '.pptx':
        info['type'] = 'PowerPoint'
        try:
            from pptx import Presentation
            prs = Presentation(path)
            titles = []
            image_count = 0
            for slide in prs.slides:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    titles.append(slide.shapes.title.text.strip())
                for shape in slide.shapes:
                    if getattr(shape, 'shape_type', None) == 13:  # picture
                        image_count += 1
            info.update({
                'slide_count': len(prs.slides),
                'titles': titles[:50],
                'title_count': len(titles),
                'image_count': image_count,
            })
        except Exception as exc:
            info['analysis_error'] = str(exc)

        checks = [
            {'code':'ppt_slides','label':'Počet snímků'},
            {'code':'ppt_titles','label':'Nadpisy snímků'},
            {'code':'ppt_images','label':'Obrázky'},
        ]

    elif ext == '.py':
        info['type'] = 'Python'
        try:
            import ast
            source = Path(path).read_text(encoding='utf-8', errors='replace')
            tree = ast.parse(source)
            functions = [n.name for n in ast.walk(tree) if isinstance(n, ast.FunctionDef)]
            variables = sorted({n.id for n in ast.walk(tree) if isinstance(n, ast.Name) and isinstance(n.ctx, ast.Store)})
            imports = []
            for n in ast.walk(tree):
                if isinstance(n, ast.Import):
                    imports.extend(a.name for a in n.names)
                elif isinstance(n, ast.ImportFrom):
                    imports.append(n.module or '')
            calls = [getattr(n.func, 'id', '') for n in ast.walk(tree) if isinstance(n, ast.Call)]
            info.update({
                'functions': functions,
                'variables': variables[:80],
                'imports': sorted(set(x for x in imports if x)),
                'has_if': any(isinstance(n, ast.If) for n in ast.walk(tree)),
                'has_loop': any(isinstance(n, (ast.For, ast.While)) for n in ast.walk(tree)),
                'has_input': 'input' in calls,
                'has_print': 'print' in calls,
                'line_count': len(source.splitlines()),
            })
        except Exception as exc:
            info['syntax_error'] = str(exc)

        checks = [
            {'code':'py_syntax','label':'Program bez syntaktické chyby'},
            {'code':'py_function','label':'Vlastní funkce'},
            {'code':'py_condition','label':'Podmínka if'},
            {'code':'py_loop','label':'Cyklus for / while'},
            {'code':'py_io','label':'Vstup a výstup programu'},
            {'code':'py_imports','label':'Použité knihovny'},
            {'code':'py_required_names','label':'Názvy funkcí podle učitelského řešení'},
        ]
    else:
        info['type'] = 'Soubor'
        checks = [{'code':'file_type','label':'Správný typ souboru'}]

    return info, checks


def generated_assignment(info):
    """Návrh instrukcí. Učitel ho před zveřejněním může libovolně přepsat."""
    ext = info.get('extension','')
    lines = []

    if ext in ('.xlsx','.xlsm'):
        specs = info.get('sheet_specs') or []
        lines.append('Vytvoř nový prázdný sešit v Excelu a zpracuj jej podle následujících požadavků:')
        if info.get('sheets'):
            if len(info['sheets']) == 1:
                lines.append(f'• Pracuj v jednom listu. List pojmenuj „{info["sheets"][0]}“.')
            else:
                lines.append('• Vytvoř listy: ' + ', '.join(info['sheets']) + '.')
        for sp in specs[:5]:
            headers = [x for x in sp.get('headers',[]) if x]
            if headers:
                lines.append(f'• V listu „{sp["name"]}“ vytvoř záhlaví: ' + ' | '.join(headers) + '.')
            if sp.get('rows',0) > 1:
                lines.append(f'• Tabulka v listu „{sp["name"]}“ má mít přibližně {sp["rows"]} řádků včetně záhlaví.')
        funcs = info.get('formula_functions') or []
        if funcs:
            lines.append('• Pro výpočty použij funkce: ' + ', '.join(funcs) + '.')
        elif info.get('formula_count',0):
            lines.append(f'• Použij alespoň {info["formula_count"]} vzorec/vzorce.')
        if info.get('chart_count',0):
            lines.append(f'• Vytvoř {info["chart_count"]} graf/grafy.')
        lines.append('• Soubor ulož ve formátu Excel a nahraj ho do lekce ke kontrole.')

    elif ext == '.docx':
        lines.append('Vytvoř nový prázdný dokument ve Wordu a zpracuj jej podle následujících požadavků:')
        if info.get('heading_count',0):
            lines.append(f'• Použij alespoň {info["heading_count"]} nadpis/nadpisy.')
        if info.get('paragraph_count',0):
            lines.append(f'• Dokument rozděl přibližně do {info["paragraph_count"]} odstavců.')
        if info.get('word_count',0):
            lines.append(f'• Rozsah dokumentu má být přibližně {info["word_count"]} slov.')
        if info.get('table_count',0):
            lines.append(f'• Vlož {info["table_count"]} tabulku/tabulky.')
        if info.get('image_count',0):
            lines.append(f'• Vlož alespoň {info["image_count"]} obrázek/obrázky.')
        lines.append('• Hotový dokument ulož jako DOCX a nahraj ho do lekce.')

    elif ext == '.pptx':
        lines.append('Vytvoř novou prázdnou prezentaci v PowerPointu:')
        if info.get('slide_count',0):
            lines.append(f'• Prezentace má mít {info["slide_count"]} snímků.')
        if info.get('title_count',0):
            lines.append(f'• Použij nadpisy alespoň na {info["title_count"]} snímcích.')
        if info.get('image_count',0):
            lines.append(f'• Vlož alespoň {info["image_count"]} obrázek/obrázky.')
        lines.append('• Hotovou prezentaci ulož jako PPTX a nahraj ji do lekce.')

    elif ext == '.py':
        lines.append('Vytvoř nový Python soubor od nuly a naprogramuj řešení podle těchto požadavků:')
        funcs = info.get('functions') or []
        if funcs:
            lines.append('• Vytvoř funkci/funkce: ' + ', '.join(funcs) + '.')
        if info.get('has_if'):
            lines.append('• Použij podmínku if/else.')
        if info.get('has_loop'):
            lines.append('• Použij cyklus for nebo while.')
        if info.get('has_input'):
            lines.append('• Program má načíst vstup od uživatele.')
        if info.get('has_print'):
            lines.append('• Program má zobrazit výsledek uživateli.')
        if info.get('imports'):
            lines.append('• Použij knihovnu/knihovny: ' + ', '.join(info['imports']) + '.')
        lines.append('• Program ulož jako .py a nahraj ho do lekce.')

    else:
        lines.append('Vytvoř nový soubor podle zadání učitele a nahraj ho do lekce.')

    return '\n'.join(lines)


def check_hint(code):
    return {
        'excel_sheets':'Podívej se dole v Excelu na názvy a počet záložek listů.',
        'excel_headers':'Zkontroluj první řádek tabulky. Názvy sloupců musí odpovídat zadání.',
        'excel_size':'Zkontroluj, jestli má tabulka požadovaný počet řádků a sloupců.',
        'excel_filled':'Zkontroluj, zda v tabulce nezůstala místa, která mají být vyplněná.',
        'excel_formulas':'Výpočet v Excelu musí být vzorec – nezačínej výsledkem, ale znakem =.',
        'excel_functions':'Podívej se na funkci ve vzorci. Může jít například o SUM, AVERAGE nebo IF.',
        'excel_chart':'Označ vhodná data a zkus kartu Vložení → Graf.',
        'word_length':'Zkontroluj, zda dokument není kratší než požadovaný rozsah.',
        'word_headings':'Použij styl Nadpis, ne pouze větší nebo tučné písmo.',
        'word_table':'Tabulku vložíš přes Vložit → Tabulka.',
        'word_images':'Zkontroluj, zda jsou v dokumentu vložené požadované obrázky.',
        'ppt_slides':'Spočítej snímky v levém panelu PowerPointu.',
        'ppt_titles':'Zkontroluj, zda mají požadované snímky skutečný nadpis.',
        'ppt_images':'Zkontroluj, zda jsi vložil požadované obrázky.',
        'py_syntax':'Spusť program. Python ti v první chybové hlášce ukáže řádek, kde je problém.',
        'py_function':'Vlastní funkce začíná klíčovým slovem def.',
        'py_condition':'Podmínku vytvoříš pomocí if, případně elif a else.',
        'py_loop':'Pro opakování použij for nebo while.',
        'py_io':'Zkontroluj práci se vstupem a výstupem, například input() a print().',
        'py_imports':'Zkontroluj importy na začátku programu.',
        'py_required_names':'Porovnej názvy svých funkcí se zadáním.',
    }.get(code, 'Vrať se k zadání a zkontroluj tuto část práce krok po kroku.')


def evaluate_informatics_file(student_path, student_name, task):
    teacher = _safe_json(task.analysis_json, {})
    raw_checks = _safe_json(task.checks_json, [])
    checks = []
    for item in raw_checks:
        if isinstance(item, str):
            checks.append({'code': item, 'question': ''})
        elif isinstance(item, dict) and item.get('code'):
            checks.append(item)
    student, _ = analyze_informatics_file(student_path, student_name)
    results = []

    for check in checks:
        code = check.get('code','')
        question = check.get('question','')
        custom_hint = check.get('hint','')
        ok = True
        label = code

        if code == 'excel_sheets':
            want = teacher.get('sheets') or []
            have = student.get('sheets') or []
            ok = want == have
            label = 'Počet a názvy listů'
        elif code == 'excel_headers':
            want = teacher.get('sheet_specs') or []
            have = student.get('sheet_specs') or []
            if len(have) < len(want):
                ok = False
            else:
                for i, sp in enumerate(want):
                    if (sp.get('headers') or []) != (have[i].get('headers') or []):
                        ok = False; break
            label = 'Záhlaví tabulky'
        elif code == 'excel_size':
            want = teacher.get('sheet_specs') or []
            have = student.get('sheet_specs') or []
            ok = len(have) >= len(want)
            if ok:
                for i, sp in enumerate(want):
                    if have[i].get('rows',0) < sp.get('rows',0) or have[i].get('cols',0) < sp.get('cols',0):
                        ok=False; break
            label = 'Rozsah tabulky'
        elif code == 'excel_filled':
            ok = int(student.get('nonempty_count',0) or 0) >= int(teacher.get('nonempty_count',0) or 0)
            label = 'Vyplněné části tabulky'
        elif code == 'excel_formulas':
            ok = int(student.get('formula_count',0) or 0) >= int(teacher.get('formula_count',0) or 0)
            label = 'Použití vzorců'
        elif code == 'excel_functions':
            ok = set(teacher.get('formula_functions') or []).issubset(set(student.get('formula_functions') or []))
            label = 'Požadované funkce'
        elif code == 'excel_chart':
            ok = int(student.get('chart_count',0) or 0) >= int(teacher.get('chart_count',0) or 0)
            label = 'Graf'
        elif code == 'word_length':
            ok = int(student.get('word_count',0) or 0) >= int(teacher.get('word_count',0) or 0)
            label = 'Rozsah dokumentu'
        elif code == 'word_headings':
            ok = int(student.get('heading_count',0) or 0) >= int(teacher.get('heading_count',0) or 0)
            label = 'Nadpisy'
        elif code == 'word_table':
            ok = int(student.get('table_count',0) or 0) >= int(teacher.get('table_count',0) or 0)
            label = 'Tabulky'
        elif code == 'word_images':
            ok = int(student.get('image_count',0) or 0) >= int(teacher.get('image_count',0) or 0)
            label = 'Obrázky'
        elif code == 'ppt_slides':
            ok = int(student.get('slide_count',0) or 0) >= int(teacher.get('slide_count',0) or 0)
            label = 'Počet snímků'
        elif code == 'ppt_titles':
            ok = int(student.get('title_count',0) or 0) >= int(teacher.get('title_count',0) or 0)
            label = 'Nadpisy snímků'
        elif code == 'ppt_images':
            ok = int(student.get('image_count',0) or 0) >= int(teacher.get('image_count',0) or 0)
            label = 'Obrázky'
        elif code == 'py_syntax':
            ok = not student.get('syntax_error')
            label = 'Program bez syntaktické chyby'
        elif code == 'py_function':
            ok = len(student.get('functions') or []) >= len(teacher.get('functions') or [])
            label = 'Vlastní funkce'
        elif code == 'py_condition':
            ok = (not teacher.get('has_if')) or bool(student.get('has_if'))
            label = 'Podmínka if'
        elif code == 'py_loop':
            ok = (not teacher.get('has_loop')) or bool(student.get('has_loop'))
            label = 'Cyklus'
        elif code == 'py_io':
            ok = ((not teacher.get('has_input')) or student.get('has_input')) and ((not teacher.get('has_print')) or student.get('has_print'))
            label = 'Vstup a výstup'
        elif code == 'py_imports':
            ok = set(teacher.get('imports') or []).issubset(set(student.get('imports') or []))
            label = 'Použité knihovny'
        elif code == 'py_required_names':
            ok = set(teacher.get('functions') or []).issubset(set(student.get('functions') or []))
            label = 'Názvy funkcí'

        results.append({
            'code':code, 'label':label, 'question':question,
            'ok':bool(ok), 'hint':custom_hint or check_hint(code)
        })

    return results


def informatics_preview(path, original_name):
    ext = Path(original_name).suffix.lower()
    if ext in ('.xlsx','.xlsm'):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=False)
            ws = wb[wb.sheetnames[0]]
            rows = []
            for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row,25), max_col=min(ws.max_column,12)):
                rows.append([c.value for c in row])
            return {'kind':'table','title':ws.title,'rows':rows}
        except Exception as exc:
            return {'kind':'text','text':f'Náhled se nepodařilo vytvořit: {exc}'}
    if ext == '.docx':
        try:
            from docx import Document
            doc = Document(path)
            text_value = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
            return {'kind':'text','text':text_value[:18000]}
        except Exception as exc:
            return {'kind':'text','text':f'Náhled se nepodařilo vytvořit: {exc}'}
    if ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(path)
            lines = []
            for i, slide in enumerate(prs.slides[:25],1):
                texts = [sh.text.strip() for sh in slide.shapes if hasattr(sh,'text') and sh.text.strip()]
                lines.append(f'Snímek {i}: ' + ' | '.join(texts))
            return {'kind':'text','text':'\n'.join(lines)}
        except Exception as exc:
            return {'kind':'text','text':f'Náhled se nepodařilo vytvořit: {exc}'}
    if ext == '.py':
        return {'kind':'code','text':Path(path).read_text(encoding='utf-8', errors='replace')[:18000]}
    return {'kind':'text','text':'Soubor byl nahrán. Pro tento typ souboru zatím není náhled.'}


def _informatics_signature(user_id, task_id, nonce):
    secret = str(app.config.get('SECRET_KEY') or 'change-me').encode('utf-8')
    msg = f'{user_id}:{task_id}:{nonce}'.encode('utf-8')
    return hmac.new(secret, msg, hashlib.sha256).hexdigest()


def _informatics_make_token(user_id, task_id):
    nonce = uuid.uuid4().hex
    sig = _informatics_signature(user_id, task_id, nonce)
    return f'UCEBNICE2:{user_id}:{task_id}:{nonce}:{sig}'


def _informatics_verify_token(token, user_id, task_id):
    try:
        prefix, uid, tid, nonce, sig = str(token or '').strip().split(':', 4)
        if prefix != 'UCEBNICE2' or int(uid) != int(user_id) or int(tid) != int(task_id):
            return False
        return hmac.compare_digest(sig, _informatics_signature(uid, tid, nonce))
    except Exception:
        return False


def _embed_work_token(path, ext, token):
    ext = ext.lower()
    if ext in ('.xlsx', '.xlsm'):
        import openpyxl
        wb = openpyxl.load_workbook(path, keep_vba=(ext == '.xlsm'))
        ws = wb['__UCEBNICE_ID__'] if '__UCEBNICE_ID__' in wb.sheetnames else wb.create_sheet('__UCEBNICE_ID__')
        ws['A1'] = token
        ws.sheet_state = 'veryHidden'
        wb.save(path)
        return
    if ext == '.docx':
        from docx import Document
        doc = Document(path)
        doc.core_properties.keywords = token
        doc.save(path)
        return
    if ext == '.pptx':
        # PowerPoint: identita studenta je na prvním SKRYTÉM snímku.
        # Kontrolní engine tento snímek pozná podle podepsaného UCEBNICE2 tokenu
        # a úplně jej vynechá z počtu snímků i ze všech kontrol prezentace.
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation(path)
        if len(prs.slides) == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
        else:
            slide = prs.slides[0]
        box = slide.shapes.add_textbox(Inches(0.1), Inches(0.1), Inches(12.0), Inches(0.5))
        box.text_frame.clear()
        run = box.text_frame.paragraphs[0].add_run()
        run.text = token
        run.font.size = Pt(1)
        # p:sld show="0" = skrytý snímek v OOXML.
        slide._element.set('show', '0')
        # Druhá, záložní kopie tokenu kvůli kompatibilitě se staršími pracovními soubory.
        prs.core_properties.keywords = token
        prs.save(path)
        return
    if ext == '.py':
        text = Path(path).read_text(encoding='utf-8', errors='replace')
        Path(path).write_text(f'# UCEBNICE2-ID: {token}\n' + text, encoding='utf-8')


def _extract_work_token(path, ext):
    ext = ext.lower()
    try:
        if ext in ('.xlsx', '.xlsm'):
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=False, data_only=False, keep_vba=(ext == '.xlsm'))
            if '__UCEBNICE_ID__' not in wb.sheetnames:
                return ''
            return str(wb['__UCEBNICE_ID__']['A1'].value or '').strip()
        if ext == '.docx':
            from docx import Document
            return str(Document(path).core_properties.keywords or '').strip()
        if ext == '.pptx':
            # Primárně čteme podepsanou identitu přímo z OOXML skrytého snímku.
            try:
                import zipfile
                from xml.etree import ElementTree as ET
                _P='http://schemas.openxmlformats.org/presentationml/2006/main'
                _A='http://schemas.openxmlformats.org/drawingml/2006/main'
                with zipfile.ZipFile(path) as z:
                    names=sorted((n for n in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml$',n)),
                                 key=lambda n:int(re.search(r'slide(\d+)',n).group(1)))
                    for n in names[:2]:
                        root=ET.fromstring(z.read(n))
                        text=''.join((t.text or '') for t in root.iter('{%s}t' % _A))
                        m=re.search(r'(UCEBNICE2:[^\s]+)', text)
                        if m:
                            return m.group(1)
            except Exception:
                pass
            # Kompatibilita se starší verzí, která ukládala token do keywords.
            from pptx import Presentation
            return str(Presentation(path).core_properties.keywords or '').strip()
        if ext == '.py':
            head = Path(path).read_text(encoding='utf-8', errors='replace')[:1000]
            m = re.search(r'UCEBNICE2-ID:\s*(UCEBNICE2:[^\s]+)', head)
            return m.group(1) if m else ''
    except Exception:
        return ''
    return ''


def _create_student_work_file(task, user):
    ext = Path(task.source_original).suffix.lower()
    user_dir = INFORMATICS_SUBMISSION_DIR / str(user.id) / 'workfiles'
    user_dir.mkdir(parents=True, exist_ok=True)
    safe_base = secure_filename(Path(task.source_original).stem) or f'ukol_{task.id}'
    filename = f'{safe_base}_student_{user.id}{ext}'
    path = user_dir / filename

    # Každý student má pro daný úkol jeden stálý podepsaný originál.
    # Opakované stažení už nezneplatní dříve rozpracovanou kopii stejného studenta.
    row = InformaticsWorkFile.query.filter_by(user_id=user.id, task_id=task.id).first()
    token = row.token if row and row.token and _informatics_verify_token(row.token, user.id, task.id) else _informatics_make_token(user.id, task.id)

    # Vytvoříme čistý pracovní soubor stejného typu, nikoli učitelovo řešení.
    if ext in ('.xlsx', '.xlsm'):
        import openpyxl
        wb = openpyxl.Workbook()
        wb.active.title = 'List1'
        wb.save(path)
    elif ext == '.docx':
        from docx import Document
        Document().save(path)
    elif ext == '.pptx':
        from pptx import Presentation
        Presentation().save(path)
    elif ext == '.py':
        path.write_text('# Pracovní soubor – vypracuj zadání níže.\n', encoding='utf-8')
    else:
        raise ValueError('Nepodporovaný typ pracovního souboru.')

    _embed_work_token(path, ext, token)
    if not row:
        row = InformaticsWorkFile(user_id=user.id, task_id=task.id, token=token, original_name=filename, stored_name=filename)
        db.session.add(row)
    else:
        row.token = token
        row.original_name = filename
        row.stored_name = filename
        row.downloaded_at = datetime.utcnow()
    db.session.commit()
    return path, filename


# --- Informatika Engine 2.1: rozšířené kontroly Office + funkční testy Pythonu ---
from informatics_engine_v2 import (
    analyze_file as _inf2_analyze_file,
    generated_assignment as _inf2_generated_assignment,
    check_hint as _inf2_check_hint,
    evaluate as _inf2_evaluate,
    preview as _inf2_preview,
)

# Přepíšeme původní univerzální funkce novou verzí. Starší implementace výše
# zůstává v souboru jen kvůli snadnému návratu při testování.
def analyze_informatics_file(path, original_name):
    return _inf2_analyze_file(path, original_name)

def generated_assignment(info):
    return _inf2_generated_assignment(info)

def check_hint(code):
    return _inf2_check_hint(code)

# České názvy kontrol v učitelském rozhraní. Technické kódy student ani učitel vidět nemusí.
INFORMATICS_CHECK_LABELS = {
    'word_sections':'Počet oddílů', 'word_unlinked':'Vypnuté „Propojit s předchozím“',
    'word_page_numbering':'Číslování stránek', 'word_alignment':'Vodorovné a svislé zarovnání',
    'word_images':'Vložené obrázky', 'word_image_size':'Velikost obrázků', 'word_image_center':'Obrázky zarovnané na střed',
    'word_title_first':'Titul / nadpis na první stránce', 'word_toc':'Automatický obsah',
    'word_heading_styles':'Styly nadpisů', 'word_heading_numbering':'Číslování nadpisů',
    'word_captions':'Titulky obrázků', 'word_citations':'Vložené citace', 'word_sources':'Vložené zdroje citací',
    'word_pages':'Počet stran', 'word_list_figures':'Seznam obrázků', 'word_bibliography':'Závěrečná bibliografie / seznam literatury',
}
WORD_CHECK_CODES = list(INFORMATICS_CHECK_LABELS)
INFORMATICS_CHECK_LABELS.update({
    'ppt_slides':'Počet snímků',
    'ppt_titles':'Nadpisy – font a velikost písma',
    'ppt_images':'Obrázky – počet a velikost',
    'ppt_tables':'Tabulky', 'ppt_charts':'Grafy', 'ppt_shapes':'Tvary a textová pole',
    'ppt_bullets':'Odrážky', 'ppt_animations':'Animace – počet a stejné typy',
    'ppt_transition':'Přechody – stejné typy',
    'ppt_background':'Pozadí snímků – musí nějaké být',
})

def evaluate_informatics_file(student_path, student_name, task):
    teacher = _safe_json(task.analysis_json, {})
    raw_checks = _safe_json(task.checks_json, [])

    # Starší úlohy mají v DB uloženou analýzu vytvořenou před kontrolou
    # skutečného pole PAGE. Pokud je u nich zapnutá kontrola číslování,
    # znovu načteme učitelský DOCX ze zdroje, aby nebylo nutné lekci mazat
    # nebo znovu nahrávat jen kvůli této opravě.
    check_codes = [x if isinstance(x,str) else x.get('code','') for x in raw_checks]
    if 'word_page_numbering' in check_codes and 'page_field_count' not in teacher:
        try:
            source = INFORMATICS_SOURCE_DIR / task.source_stored
            if source.exists():
                refreshed, _ = _inf2_analyze_file(source, task.source_original or source.name)
                if str(refreshed.get('type','')).lower() == 'word':
                    teacher = refreshed
        except Exception:
            pass

    # Starší již vytvořené Excelové úlohy nemají v analysis_json uložené
    # operátory vzorců. Při první kontrole je znovu načteme z učitelského XLSX,
    # takže kvůli této opravě není nutné úlohu mazat ani znovu vytvářet.
    if 'excel_functions' in check_codes and 'formula_operators' not in teacher:
        try:
            source = INFORMATICS_SOURCE_DIR / task.source_stored
            if source.exists():
                refreshed, _ = _inf2_analyze_file(source, task.source_original or source.name)
                if str(refreshed.get('type','')).lower() == 'excel':
                    teacher = refreshed
        except Exception:
            pass

    # PowerPoint kontroly nyní vycházejí přímo z OOXML (font/velikost nadpisů,
    # rozměry obrázků, typy animací a přechodů, pozadí). Staré analysis_json
    # proto při první kontrole transparentně obnovíme z učitelského PPTX.
    ppt_codes={'ppt_slides','ppt_titles','ppt_images','ppt_tables','ppt_charts','ppt_shapes','ppt_bullets','ppt_animations','ppt_transition','ppt_background'}
    # Jednodušší a bezpečnější: u PowerPointu s aktivními kontrolami načti učitelský
    # zdroj vždy znovu. Je malý a máme jistotu, že používáme nejnovější XML profil.
    if bool(ppt_codes.intersection(check_codes)) and str(teacher.get('type','')).lower() == 'powerpoint':
        try:
            source = INFORMATICS_SOURCE_DIR / task.source_stored
            if source.exists():
                refreshed, _ = _inf2_analyze_file(source, task.source_original or source.name)
                if str(refreshed.get('type','')).lower() == 'powerpoint':
                    teacher = refreshed
        except Exception:
            pass

    return _inf2_evaluate(student_path, student_name, teacher, raw_checks)

def informatics_preview(path, original_name, teacher=False):
    return _inf2_preview(path, original_name, teacher=teacher)

@app.route('/informatics-task/<int:task_id>/download-workfile')
def download_informatics_workfile(task_id):
    r = require_login()
    if r: return r
    user = current_user()
    if not user or user.role != 'student':
        return 'Pracovní soubor stahuje student.', 403
    task = db.session.get(InformaticsTask, task_id)
    if not task or not informatics_task_unlocked(task):
        return 'Úkol není dostupný.', 404
    path, filename = _create_student_work_file(task, user)
    return send_file(path, as_attachment=True, download_name=filename)


def informatics_task_unlocked(task):
    u = current_user()
    if not u or u.role == 'teacher' or task.order <= 1:
        return True
    previous = InformaticsTask.query.filter_by(lesson_id=task.lesson_id, order=task.order-1).first()
    if not previous:
        return True
    last = InformaticsSubmission.query.filter_by(
        user_id=u.id, task_id=previous.id, status='odevzdáno'
    ).order_by(InformaticsSubmission.created_at.desc()).first()
    return bool(last and last.percent == 100)


@app.route('/teacher/informatics/new', methods=['GET','POST'])
def new_informatics_lesson():
    r = require_teacher()
    if r: return r

    if request.method == 'POST':
        school = request.form.get('school','').strip()
        grade_name = request.form.get('grade_name','').strip()
        topic = request.form.get('topic','').strip()
        title = request.form.get('title','').strip()
        intro = request.form.get('intro','').strip()
        html_file = request.files.get('lesson_html')
        html_original = ''
        html_stored = ''
        if html_file and html_file.filename:
            if Path(html_file.filename).suffix.lower() not in ('.html', '.htm'):
                flash('Společný výklad musí být HTML soubor.')
                return redirect(url_for('new_informatics_lesson'))
            html_original = html_file.filename
            html_stored = _save_uploaded_file(html_file, INFORMATICS_SOURCE_DIR, 'lesson_html')

        if not all((school, grade_name, topic, title)):
            flash('Vyplň školu/třídu, ročník, téma a název lekce.')
            return redirect(url_for('new_informatics_lesson'))

        source_files = request.files.getlist('task_files')
        preview_files = request.files.getlist('task_preview_files')
        task_titles = request.form.getlist('task_titles')
        valid = [(i,f) for i,f in enumerate(source_files) if f and f.filename]
        if not valid:
            flash('Přidej alespoň jeden hotový učitelský soubor k analýze.')
            return redirect(url_for('new_informatics_lesson'))

        allowed = {'.xlsx','.xlsm','.docx','.pptx','.py'}
        prepared = []
        for i, f in valid:
            ext = Path(f.filename).suffix.lower()
            if ext not in allowed:
                flash(f'{f.filename}: podporovaný je Excel, Word, PowerPoint nebo Python.')
                return redirect(url_for('new_informatics_lesson'))
            stored = _save_uploaded_file(f, INFORMATICS_SOURCE_DIR, 'teacher')
            preview_stored = ''
            # Přesný vizuální vzor dodává učitel jako PDF pro libovolný typ úkolu
            # (Word, Excel, PowerPoint i Python). Zdrojový Office/Python soubor se
            # používá pro hodnocení; PDF pouze pro zobrazení studentovi v okně.
            if i < len(preview_files):
                pf = preview_files[i]
                if pf and pf.filename:
                    if Path(pf.filename).suffix.lower() != '.pdf':
                        flash(f'{pf.filename}: náhled učitelského vzoru musí být PDF.')
                        return redirect(url_for('new_informatics_lesson'))
                    preview_stored = _save_uploaded_file(pf, INFORMATICS_SOURCE_DIR, 'teacher_preview')
            analysis, suggested = analyze_informatics_file(INFORMATICS_SOURCE_DIR/stored, f.filename)
            prepared.append({
                'title': (task_titles[i].strip() if i < len(task_titles) else '') or f'Úkol {len(prepared)+1}',
                'source_original': f.filename,
                'source_stored': stored,
                'preview_stored': preview_stored,
                'analysis': analysis,
                'suggested': suggested,
                'assignment': generated_assignment(analysis),
            })

        session['informatics_builder'] = {
            'school':school,'grade_name':grade_name,'topic':topic,'title':title,'intro':intro,
            'html_original':html_original,'html_stored':html_stored,
            'tasks':prepared
        }
        session.modified = True
        return redirect(url_for('review_informatics_lesson'))

    lessons = InformaticsLesson.query.order_by(InformaticsLesson.created_at.desc()).all()
    return render_template('informatics_new.html', course=course_from_lesson(None), lesson=None, lessons=lessons)


@app.route('/teacher/informatics/review', methods=['GET','POST'])
def review_informatics_lesson():
    r = require_teacher()
    if r: return r
    data = session.get('informatics_builder')
    if not data:
        flash('Nejdřív nahraj soubory k analýze.')
        return redirect(url_for('new_informatics_lesson'))

    if request.method == 'POST':
        lesson_row = InformaticsLesson(
            school=data['school'], grade_name=data['grade_name'], topic=data['topic'],
            title=data['title'], intro=data.get('intro',''),
            html_original=data.get('html_original',''), html_stored=data.get('html_stored',''),
            is_published=True
        )
        db.session.add(lesson_row); db.session.flush()

        for i, t in enumerate(data['tasks']):
            selected_codes = request.form.getlist(f'checks_{i}')
            check_items = []
            for code in selected_codes:
                q = request.form.get(f'question_{i}_{code}', '').strip()
                hint = request.form.get(f'hint_{i}_{code}', '').strip()
                if not q:
                    label = next((c.get('label','') for c in t.get('suggested',[]) if c.get('code') == code), code)
                    q = f'Splň požadavek: {label}.'
                if not hint:
                    hint = check_hint(code)
                check_items.append({'code': code, 'question': q, 'hint': hint})
            assignment = request.form.get(f'assignment_{i}', t['assignment']).strip()
            task_title = request.form.get(f'task_title_{i}', t['title']).strip() or t['title']
            db.session.add(InformaticsTask(
                lesson_id=lesson_row.id, order=i+1, title=task_title, assignment=assignment,
                source_original=t['source_original'], source_stored=t['source_stored'],
                file_type=t['analysis'].get('type',''),
                analysis_json=json.dumps(t['analysis'], ensure_ascii=False),
                checks_json=json.dumps(check_items, ensure_ascii=False),
                image_file=t.get('preview_stored','')
            ))

        db.session.commit()
        session.pop('informatics_builder', None)
        flash('Informatická lekce byla vytvořena a zveřejněna.')
        return redirect(url_for('informatics_lesson', lesson_id=lesson_row.id))

    for t in data.get('tasks', []):
        source_path = INFORMATICS_SOURCE_DIR / t.get('source_stored','')
        t['preview'] = informatics_preview(source_path, t.get('source_original',''), teacher=True) if source_path.exists() else None
    return render_template('informatics_review.html', course=course_from_lesson(None), lesson=None, data=data)



@app.route('/teacher/informatics/<int:lesson_id>/edit', methods=['GET','POST'])
def edit_informatics_lesson(lesson_id):
    r = require_teacher()
    if r: return r
    item = db.session.get(InformaticsLesson, lesson_id)
    if not item:
        flash('Informatická lekce nebyla nalezena.')
        return redirect(url_for('teacher_home'))

    tasks = InformaticsTask.query.filter_by(lesson_id=item.id).order_by(InformaticsTask.order).all()

    if request.method == 'POST':
        item.school = request.form.get('school','').strip()
        item.grade_name = request.form.get('grade_name','').strip()
        item.topic = request.form.get('topic','').strip()
        item.title = request.form.get('title','').strip()
        item.intro = request.form.get('intro','').strip()
        if not all((item.school, item.grade_name, item.topic, item.title)):
            flash('Vyplň školu/třídu, ročník, téma a název lekce.')
            return redirect(url_for('edit_informatics_lesson', lesson_id=item.id))

        html_file = request.files.get('lesson_html')
        if html_file and html_file.filename:
            if Path(html_file.filename).suffix.lower() not in ('.html','.htm'):
                flash('Společný výklad musí být HTML soubor.')
                return redirect(url_for('edit_informatics_lesson', lesson_id=item.id))
            item.html_original = html_file.filename
            item.html_stored = _save_uploaded_file(html_file, INFORMATICS_SOURCE_DIR, 'lesson_html')

        allowed={'.xlsx','.xlsm','.docx','.pptx','.py'}
        for task in tasks:
            task.title = request.form.get(f'task_title_{task.id}', task.title).strip() or task.title
            task.assignment = request.form.get(f'assignment_{task.id}', task.assignment).strip()

            # Učitel může při editaci kontroly nejen přepsat, ale také přidat/odebrat.
            old_checks={ch.get('code',''):ch for ch in _safe_json(task.checks_json, []) if ch.get('code')}
            selected_codes=request.form.getlist(f'checks_{task.id}')
            edited=[]
            for code in selected_codes:
                old=old_checks.get(code,{})
                q=request.form.get(f'question_{task.id}_{code}', old.get('question','')).strip()
                hint=request.form.get(f'hint_{task.id}_{code}', old.get('hint','')).strip()
                if not q:
                    q='Splň požadavek: '+INFORMATICS_CHECK_LABELS.get(code, code)+'.'
                if not hint:
                    hint=check_hint(code)
                edited.append({'code':code,'question':q,'hint':hint})
            task.checks_json=json.dumps(edited, ensure_ascii=False)

            replacement=request.files.get(f'source_file_{task.id}')
            if replacement and replacement.filename:
                ext=Path(replacement.filename).suffix.lower()
                if ext not in allowed:
                    flash(f'{replacement.filename}: podporovaný je Excel, Word, PowerPoint nebo Python.')
                    return redirect(url_for('edit_informatics_lesson', lesson_id=item.id))
                stored=_save_uploaded_file(replacement, INFORMATICS_SOURCE_DIR, 'teacher')
                analysis, suggested=analyze_informatics_file(INFORMATICS_SOURCE_DIR/stored, replacement.filename)
                task.source_original=replacement.filename
                task.source_stored=stored
                task.file_type=analysis.get('type','')
                task.analysis_json=json.dumps(analysis, ensure_ascii=False)
                # Smažeme starý vygenerovaný PDF náhled, aby se po změně souboru nepoužil cache.
                try:
                    cached=INFORMATICS_SOURCE_DIR/'_visual_previews'/f'task_{task.id}.pdf'
                    if cached.exists(): cached.unlink()
                except Exception:
                    pass

            # Přesný učitelský náhled pro Word, Excel, PowerPoint i Python.
            # Je to pouze vizuální PDF v okně; hodnocení stále vychází ze zdrojového souboru.
            visual_preview = request.files.get(f'preview_file_{task.id}')
            if visual_preview and visual_preview.filename:
                pext = Path(visual_preview.filename).suffix.lower()
                if pext != '.pdf':
                    flash(f'{visual_preview.filename}: náhled učitelského vzoru musí být PDF.')
                    return redirect(url_for('edit_informatics_lesson', lesson_id=item.id))
                task.image_file = _save_uploaded_file(visual_preview, INFORMATICS_SOURCE_DIR, 'teacher_preview')
            if request.form.get(f'remove_preview_{task.id}') == '1':
                task.image_file = ''

        db.session.commit()
        flash('Informatická lekce byla upravena.')
        return redirect(url_for('teacher_home'))

    task_data=[]
    for task in tasks:
        current=_safe_json(task.checks_json, [])
        current_map={ch.get('code',''):ch for ch in current if ch.get('code')}
        analysis=_safe_json(task.analysis_json,{})
        options=[]
        if str(task.file_type or analysis.get('type','')).lower() == 'word':
            # U Wordu zobrazíme všechny podporované kontroly. Učitel tak může doplnit i obsah,
            # citace, zdroje a bibliografii, i když je starší analyzátor při prvním importu nenavrhl.
            codes=WORD_CHECK_CODES
        else:
            source=INFORMATICS_SOURCE_DIR/task.source_stored
            try:
                _, suggested=analyze_informatics_file(source, task.source_original) if source.exists() else ({},[])
                codes=list(dict.fromkeys([c.get('code') for c in suggested if c.get('code')] + list(current_map)))
            except Exception:
                codes=list(current_map)
        for code in codes:
            ch=current_map.get(code,{})
            options.append({'code':code,'label':INFORMATICS_CHECK_LABELS.get(code, code),
                            'selected':code in current_map,'question':ch.get('question',''),'hint':ch.get('hint',check_hint(code))})
        task_data.append({'task':task,'checks':current,'options':options})
    return render_template('informatics_edit.html', item=item, task_data=task_data,
                           course=course_from_lesson(None), lesson=None)


@app.route('/teacher/informatics/<int:lesson_id>/delete', methods=['POST'])
def delete_informatics_lesson(lesson_id):
    r = require_teacher()
    if r:
        return r

    item = db.session.get(InformaticsLesson, lesson_id)
    if not item:
        flash('Informatická lekce už neexistuje.')
        return redirect(url_for('teacher_home'))

    try:
        # Načteme si úkoly jako ORM objekty a mažeme je postupně.
        tasks = list(
            InformaticsTask.query.filter_by(lesson_id=item.id)
            .order_by(InformaticsTask.order)
            .all()
        )

        for task in tasks:
            # Nejprve všechny studentské výsledky k úkolu.
            submissions = list(
                InformaticsSubmission.query.filter_by(task_id=task.id).all()
            )
            for submission in submissions:
                db.session.delete(submission)

            db.session.flush()
            db.session.delete(task)

        db.session.flush()
        db.session.delete(item)
        db.session.commit()

        flash('Informatická lekce byla smazána.')
    except Exception:
        db.session.rollback()
        raise

    return redirect(url_for('teacher_home'))


@app.route('/informatics-lesson/<int:lesson_id>')
def informatics_lesson(lesson_id):
    r = require_login()
    if r: return r
    item = db.session.get(InformaticsLesson, lesson_id)
    if not item or (not item.is_published and current_user().role != 'teacher'):
        return 'Lekce nebyla nalezena.', 404
    tasks = sorted(item.tasks, key=lambda x:x.order)
    states = []
    for task in tasks:
        last = None
        if current_user().role == 'student':
            last = InformaticsSubmission.query.filter_by(user_id=current_user().id, task_id=task.id).order_by(InformaticsSubmission.created_at.desc()).first()
        states.append({'task':task,'unlocked':informatics_task_unlocked(task),'last':last})
    course = {'subject':'Informatika','grade':item.grade_name,'block':item.topic,'icon':'💻'}
    lesson_obj = {'title':item.title,'block':item.topic}
    html_content = render_informatics_html(item.html_stored)
    return render_template('informatics_lesson.html', course=course, lesson=lesson_obj, item=item, states=states, html_content=html_content)



def _office_visual_html_fallback(source, ext, title='Vzorový soubor'):
    """Vizuální Office náhled bez systémového LibreOffice."""
    esc = html.escape
    shell_head = '''<!doctype html><html lang="cs"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>__TITLE__</title><style>
*{box-sizing:border-box} body{margin:0;background:#e9edf2;color:#111827;font-family:Arial,Segoe UI,sans-serif;padding:28px}
.word-page{position:relative;width:min(794px,100%);height:1123px;margin:0 auto 24px;background:#fff;box-shadow:0 8px 30px rgba(15,23,42,.18);overflow:hidden}
.word-page-inner{height:100%;display:flex;flex-direction:column}.word-content{width:100%}
.word-page.vcenter .word-page-inner{justify-content:center}.word-page.vbottom .word-page-inner{justify-content:flex-end}
.word-page p{line-height:1.35;white-space:pre-wrap;overflow-wrap:anywhere}
.word-page table{border-collapse:collapse;width:100%;margin:12px 0}.word-page td,.word-page th{border:1px solid #cbd5e1;padding:7px;vertical-align:top}
.doc-image{display:block;max-width:100%;height:auto;margin:10px auto}
.sheet{max-width:1200px;margin:auto;background:#fff;padding:20px;box-shadow:0 8px 30px rgba(15,23,42,.18);overflow:auto}
.slide{position:relative;width:min(960px,100%);aspect-ratio:16/9;margin:0 auto 24px;background:#fff;box-shadow:0 8px 30px rgba(15,23,42,.18);padding:48px;overflow:hidden}
.slide-num{position:absolute;right:12px;bottom:8px;color:#64748b;font-size:12px}
@media(max-width:840px){body{padding:10px}.word-page{height:auto;min-height:calc((100vw - 20px)*1.414)}}
</style></head><body>'''.replace('__TITLE__', esc(title))
    end='</body></html>'
    try:
        if ext == '.docx':
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml.ns import qn
            import base64 as _base64
            doc = Document(source)

            def twips_to_px(v, default=0):
                try: return max(0, round(int(v) / 15))
                except Exception: return default

            def sect_layout(sectPr):
                top = right = bottom = left = 95
                valign = 'top'
                if sectPr is not None:
                    pgmar = sectPr.find(qn('w:pgMar'))
                    if pgmar is not None:
                        top = twips_to_px(pgmar.get(qn('w:top')), top)
                        right = twips_to_px(pgmar.get(qn('w:right')), right)
                        bottom = twips_to_px(pgmar.get(qn('w:bottom')), bottom)
                        left = twips_to_px(pgmar.get(qn('w:left')), left)
                    va = sectPr.find(qn('w:vAlign'))
                    if va is not None:
                        valign = (va.get(qn('w:val')) or 'top').lower()
                return top,right,bottom,left,valign

            def image_html_from_paragraph(par):
                imgs=[]
                try:
                    for blip in par._p.xpath('.//a:blip'):
                        rid = blip.get(qn('r:embed'))
                        if not rid or rid not in doc.part.rels: continue
                        part = doc.part.rels[rid].target_part
                        b64 = _base64.b64encode(part.blob).decode('ascii')
                        mime = getattr(part,'content_type',None) or 'image/png'
                        width_style=''
                        for inline in par._p.xpath('.//wp:inline'):
                            extent = inline.find(qn('wp:extent'))
                            if extent is not None and extent.get('cx'):
                                width=max(1, round(int(extent.get('cx')) / 9525))
                                width_style=f'width:{width}px;'
                                break
                        imgs.append(f'<img class="doc-image" style="{width_style}" src="data:{mime};base64,{b64}" alt="Obrázek z dokumentu">')
                except Exception:
                    pass
                return ''.join(imgs)

            def paragraph_html(par):
                align={WD_ALIGN_PARAGRAPH.CENTER:'center',WD_ALIGN_PARAGRAPH.RIGHT:'right',WD_ALIGN_PARAGRAPH.JUSTIFY:'justify'}.get(par.alignment,'left')
                pf=par.paragraph_format
                pst=[f'text-align:{align}', 'margin-top:0', 'margin-bottom:10px']
                if pf.space_before: pst.append(f'margin-top:{max(0,pf.space_before.pt):.1f}pt')
                if pf.space_after: pst.append(f'margin-bottom:{max(0,pf.space_after.pt):.1f}pt')
                if pf.left_indent: pst.append(f'margin-left:{pf.left_indent.pt:.1f}pt')
                if pf.right_indent: pst.append(f'margin-right:{pf.right_indent.pt:.1f}pt')
                if pf.first_line_indent: pst.append(f'text-indent:{pf.first_line_indent.pt:.1f}pt')
                nm=(par.style.name.lower() if par.style and par.style.name else '')
                if 'title' in nm: pst += ['font-size:24pt','font-weight:700','margin-bottom:18px']
                elif 'heading 1' in nm: pst += ['font-size:18pt','font-weight:700','margin-top:18px']
                elif 'heading 2' in nm: pst += ['font-size:15pt','font-weight:700','margin-top:15px']
                elif 'caption' in nm: pst += ['font-size:9pt','font-style:italic','text-align:center']
                runs=[]
                for run in par.runs:
                    txt=esc(run.text)
                    if not txt: continue
                    st=[]
                    if run.bold: st.append('font-weight:700')
                    if run.italic: st.append('font-style:italic')
                    if run.underline: st.append('text-decoration:underline')
                    if run.font.size: st.append(f'font-size:{run.font.size.pt:.1f}pt')
                    if run.font.name: st.append('font-family:'+esc(run.font.name)+',Arial,sans-serif')
                    color=getattr(getattr(run.font,'color',None),'rgb',None)
                    if color: st.append('color:#'+str(color))
                    runs.append('<span style="'+esc(';'.join(st),quote=True)+'">'+txt+'</span>')
                body=''.join(runs) if runs else ('&nbsp;' if not par.text else esc(par.text))
                return '<p style="'+esc(';'.join(pst),quote=True)+'">'+body+'</p>'+image_html_from_paragraph(par)

            def table_html(tbl):
                rows=[]
                for row in tbl.rows:
                    cells=''.join('<td>'+esc(cell.text).replace('\\n','<br>')+'</td>' for cell in row.cells)
                    rows.append('<tr>'+cells+'</tr>')
                return '<table>'+''.join(rows)+'</table>'

            pages=[]; current=[]
            body=doc._element.body
            final_sect=body.sectPr
            def finish_page(sectPr=None):
                nonlocal current
                if not current: return
                top,right,bottom,left,valign=sect_layout(sectPr if sectPr is not None else final_sect)
                cls='word-page'
                if valign=='center': cls+=' vcenter'
                elif valign in ('bottom','both'): cls+=' vbottom'
                inner=f'<div class="word-page-inner" style="padding:{top}px {right}px {bottom}px {left}px"><div class="word-content">'+''.join(current)+'</div></div>'
                pages.append(f'<section class="{cls}">{inner}</section>')
                current=[]

            from docx.text.paragraph import Paragraph
            from docx.table import Table
            for child in body.iterchildren():
                if child.tag == qn('w:p'):
                    par=Paragraph(child, doc._body)
                    current.append(paragraph_html(par))
                    ppr=child.find(qn('w:pPr'))
                    sect=ppr.find(qn('w:sectPr')) if ppr is not None else None
                    if sect is not None:
                        finish_page(sect)
                    elif child.xpath('.//w:br[@w:type="page"]'):
                        finish_page(final_sect)
                elif child.tag == qn('w:tbl'):
                    current.append(table_html(Table(child, doc._body)))
            finish_page(final_sect)
            return shell_head+''.join(pages)+end

        if ext in ('.xlsx','.xlsm'):
            import openpyxl
            wb=openpyxl.load_workbook(source,data_only=False)
            sections=[]
            for ws in wb.worksheets:
                if ws.title=='__UCEBNICE_ID__': continue
                maxr=min(ws.max_row or 1,80); maxc=min(ws.max_column or 1,24)
                trs=[]
                for row in ws.iter_rows(min_row=1,max_row=maxr,max_col=maxc):
                    tds=[]
                    for c in row:
                        val='' if c.value is None else esc(str(c.value))
                        sty=[]
                        if c.font and c.font.bold: sty.append('font-weight:700')
                        if c.alignment and c.alignment.horizontal: sty.append('text-align:'+c.alignment.horizontal)
                        if c.fill and getattr(c.fill,'fgColor',None) and c.fill.fgColor.type=='rgb' and c.fill.fgColor.rgb:
                            sty.append('background:#'+str(c.fill.fgColor.rgb)[-6:])
                        tds.append('<td style="'+esc(';'.join(sty),quote=True)+'">'+val+'</td>')
                    trs.append('<tr>'+''.join(tds)+'</tr>')
                sections.append('<h2>'+esc(ws.title)+'</h2><table>'+''.join(trs)+'</table>')
            return shell_head+'<div class="sheet">'+''.join(sections)+'</div>'+end

        if ext == '.pptx':
            from pptx import Presentation
            prs=Presentation(source); slides=[]
            for i,slide in enumerate(prs.slides,1):
                texts=[]
                for sh in slide.shapes:
                    if hasattr(sh,'text') and sh.text.strip(): texts.append('<div style="margin:8px 0;white-space:pre-wrap">'+esc(sh.text)+'</div>')
                slides.append('<div class="slide">'+''.join(texts)+f'<span class="slide-num">Snímek {i}</span></div>')
            return shell_head+''.join(slides)+end
    except Exception as exc:
        return shell_head+'<div class="word-page"><div class="word-page-inner" style="padding:95px"><div class="word-content"><h2>Náhled souboru</h2><p>'+esc(str(exc))+'</p></div></div></div>'+end
    return shell_head+'<div class="word-page"><div class="word-page-inner" style="padding:95px"><div class="word-content"><p>Náhled tohoto souboru není dostupný.</p></div></div></div>'+end


def _make_office_preview_pdf(source, ext, out_pdf):
    """Vytvoří skutečný PDF náhled přes LibreOffice/soffice. Žádný HTML/Python převod."""
    source = Path(source); out_pdf = Path(out_pdf)
    if out_pdf.exists() and out_pdf.stat().st_mtime >= source.stat().st_mtime:
        return out_pdf
    import subprocess, tempfile, shutil
    office = shutil.which('libreoffice') or shutil.which('soffice')
    if not office:
        raise RuntimeError('Na serveru chybí LibreOffice pro převod dokumentu do PDF.')
    out_pdf.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix='infpreview_') as td:
        td=Path(td); local=td/source.name; shutil.copy2(source, local)
        proc=subprocess.run([office,'--headless','--convert-to','pdf','--outdir',str(td),str(local)],
                            stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=90, text=True)
        made=td/(local.stem+'.pdf')
        if proc.returncode != 0 or not made.exists():
            raise RuntimeError('Převod dokumentu do PDF se nezdařil.')
        shutil.copy2(made,out_pdf)
    return out_pdf


def _docx_preview_pdf_python(source, out_pdf):
    """Nouzový PDF náhled DOCX v čistém Pythonu; zachová text včetně TOC/bibliografie a obrázky."""
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak, Table, TableStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
    from reportlab.lib import colors
    from reportlab.lib.units import mm
    from docx import Document
    from docx.oxml.ns import qn
    from xml.sax.saxutils import escape
    import io, zipfile

    doc=Document(str(source)); styles=getSampleStyleSheet()
    normal=ParagraphStyle('WordNormal', parent=styles['BodyText'], fontName='Helvetica', fontSize=10.5,
                          leading=14, spaceAfter=5)
    title=ParagraphStyle('WordTitle', parent=normal, fontSize=18, leading=22, alignment=TA_CENTER, spaceAfter=10)
    h1=ParagraphStyle('WordH1', parent=normal, fontSize=15, leading=19, spaceBefore=8, spaceAfter=6, textColor=colors.HexColor('#17365d'))
    h2=ParagraphStyle('WordH2', parent=normal, fontSize=12.5, leading=16, spaceBefore=6, spaceAfter=5, textColor=colors.HexColor('#17365d'))
    story=[]
    body=doc.element.body
    media=[]
    try:
        with zipfile.ZipFile(source) as z:
            for n in z.namelist():
                if n.startswith('word/media/'):
                    media.append((n,z.read(n)))
    except Exception:
        pass
    media_i=0
    section_idx=0
    for child in body.iterchildren():
        tag=child.tag.rsplit('}',1)[-1]
        if tag=='p':
            texts=[t.text or '' for t in child.iter(qn('w:t'))]
            txt=''.join(texts).strip()
            pstyle=''
            ppr=child.find(qn('w:pPr'))
            if ppr is not None:
                ps=ppr.find(qn('w:pStyle'))
                if ps is not None: pstyle=ps.get(qn('w:val')) or ''
            if not txt:
                # vložené obrázky přibližně na místě výskytu
                drawings=list(child.iter(qn('w:drawing')))
                if drawings and media_i < len(media):
                    try:
                        img=RLImage(io.BytesIO(media[media_i][1])); media_i+=1
                        maxw=150*mm; maxh=180*mm
                        scale=min(maxw/img.imageWidth,maxh/img.imageHeight,1)
                        img.drawWidth*=scale; img.drawHeight*=scale
                        story += [Spacer(1,4*mm),img,Spacer(1,4*mm)]
                    except Exception: pass
                continue
            st=normal
            low=pstyle.lower()
            if 'title' in low or 'titul' in low: st=title
            elif 'heading1' in low or 'nadpis1' in low or low.endswith('1'): st=h1
            elif 'heading2' in low or 'nadpis2' in low or low.endswith('2'): st=h2
            # zarovnání
            align=None
            if ppr is not None:
                jc=ppr.find(qn('w:jc'))
                if jc is not None: align=jc.get(qn('w:val'))
            if align:
                amap={'center':TA_CENTER,'right':TA_RIGHT,'both':TA_JUSTIFY,'left':TA_LEFT}
                if align in amap:
                    st=ParagraphStyle('dyn'+str(len(story)), parent=st, alignment=amap[align])
            story.append(Paragraph(escape(txt), st))
            drawings=list(child.iter(qn('w:drawing')))
            if drawings and media_i < len(media):
                try:
                    img=RLImage(io.BytesIO(media[media_i][1])); media_i+=1
                    maxw=150*mm; maxh=180*mm
                    scale=min(maxw/img.imageWidth,maxh/img.imageHeight,1)
                    img.drawWidth*=scale; img.drawHeight*=scale
                    story += [Spacer(1,3*mm),img,Spacer(1,3*mm)]
                except Exception: pass
        elif tag=='tbl':
            rows=[]
            for tr in child.iter(qn('w:tr')):
                row=[]
                for tc in tr.findall(qn('w:tc')):
                    cell=''.join((t.text or '') for t in tc.iter(qn('w:t'))).strip()
                    row.append(Paragraph(escape(cell), normal))
                if row: rows.append(row)
            if rows:
                table=Table(rows, repeatRows=1)
                table.setStyle(TableStyle([('GRID',(0,0),(-1,-1),0.4,colors.grey),('VALIGN',(0,0),(-1,-1),'TOP'),('LEFTPADDING',(0,0),(-1,-1),4),('RIGHTPADDING',(0,0),(-1,-1),4)]))
                story += [table,Spacer(1,4*mm)]
        elif tag=='sectPr':
            section_idx += 1
    # přidej nevyužité obrázky (aby v náhledu nic nezmizelo)
    while media_i < len(media):
        try:
            img=RLImage(io.BytesIO(media[media_i][1])); media_i+=1
            maxw=150*mm; maxh=180*mm; scale=min(maxw/img.imageWidth,maxh/img.imageHeight,1)
            img.drawWidth*=scale; img.drawHeight*=scale; story += [Spacer(1,3*mm),img]
        except Exception: media_i+=1
    pdf=SimpleDocTemplate(str(out_pdf), pagesize=A4, rightMargin=18*mm,leftMargin=18*mm,topMargin=18*mm,bottomMargin=18*mm)
    pdf.build(story)


@app.route('/informatics-task/<int:task_id>/teacher-preview.pdf')
def informatics_teacher_preview_pdf(task_id):
    """PDF/obrázek učitelského vzoru zobrazený přímo v iframe, bez stahování."""
    r = require_login()
    if r: return r
    task = db.session.get(InformaticsTask, task_id)
    if not task:
        return 'Úkol nebyl nalezen.', 404

    # Nejvyšší priorita: učitelův přesný PDF export / screenshot.
    if task.image_file:
        manual_preview = INFORMATICS_SOURCE_DIR / task.image_file
        if manual_preview.exists():
            suffix = manual_preview.suffix.lower()
            mimetype = {'.pdf':'application/pdf','.png':'image/png','.jpg':'image/jpeg',
                        '.jpeg':'image/jpeg','.webp':'image/webp'}.get(suffix,'application/octet-stream')
            resp = send_file(manual_preview, mimetype=mimetype, as_attachment=False,
                             download_name=f'nahled_{task.id}{suffix}')
            resp.headers['Content-Disposition'] = 'inline'
            resp.headers['Cache-Control'] = 'no-store'
            return resp

    source = INFORMATICS_SOURCE_DIR / task.source_stored
    if not source.exists():
        return 'Vzorový soubor nebyl nalezen.', 404
    ext = Path(task.source_original or '').suffix.lower()
    if ext not in ('.docx','.xlsx','.xlsm','.pptx'):
        return 'Pro tento typ souboru není PDF náhled.', 404

    preview_dir = INFORMATICS_SOURCE_DIR / '_visual_previews'
    preview_dir.mkdir(parents=True, exist_ok=True)
    out_pdf = preview_dir / f'task_{task.id}.pdf'
    try:
        _make_office_preview_pdf(source, ext, out_pdf)
        resp = send_file(out_pdf, mimetype='application/pdf', as_attachment=False,
                         download_name=f'nahled_{task.id}.pdf')
        resp.headers['Content-Disposition'] = 'inline'
        resp.headers['Cache-Control'] = 'no-store'
        return resp
    except Exception as exc:
        return ('<div style="font-family:Arial;padding:24px"><b>PDF náhled se nepodařilo vytvořit.</b><br>'
                'Na serveru musí být nainstalovaný LibreOffice. Kontrola souboru tím není ovlivněna.</div>'), 503, {'Content-Type':'text/html; charset=utf-8','Cache-Control':'no-store'}


@app.route('/informatics-submission/<int:submission_id>/preview.pdf')
def informatics_submission_preview_pdf(submission_id):
    """Náhled studentovy právě zkontrolované práce přímo v okně lekce."""
    r = require_login()
    if r: return r
    submission = db.session.get(InformaticsSubmission, submission_id)
    if not submission:
        return 'Odevzdaná práce nebyla nalezena.', 404
    user = current_user()
    if user.role == 'student' and submission.user_id != user.id:
        return 'K této práci nemáš přístup.', 403
    task = db.session.get(InformaticsTask, submission.task_id)
    if not task:
        return 'Úkol nebyl nalezen.', 404
    source = INFORMATICS_SUBMISSION_DIR / str(submission.user_id) / submission.stored_name
    if not source.exists():
        return 'Soubor práce nebyl nalezen.', 404
    ext = Path(submission.original_name or '').suffix.lower()
    if ext not in ('.docx','.xlsx','.xlsm','.pptx'):
        return 'Pro tento typ souboru není PDF náhled.', 404
    preview_dir = INFORMATICS_SUBMISSION_DIR / str(submission.user_id) / '_visual_previews'
    preview_dir.mkdir(parents=True, exist_ok=True)
    out_pdf = preview_dir / f'submission_{submission.id}.pdf'
    try:
        _make_office_preview_pdf(source, ext, out_pdf)
        resp = send_file(out_pdf, mimetype='application/pdf', as_attachment=False,
                         download_name=f'prace_{submission.id}.pdf')
        resp.headers['Content-Disposition'] = 'inline'
        resp.headers['Cache-Control'] = 'no-store'
        return resp
    except Exception as exc:
        return ('<div style="font-family:Arial;padding:24px"><b>PDF náhled práce se nepodařilo vytvořit.</b><br>'
                'Na serveru musí být nainstalovaný LibreOffice. Kontrola práce tím není ovlivněna.</div>'), 503, {'Content-Type':'text/html; charset=utf-8','Cache-Control':'no-store'}


@app.route('/informatics-task/<int:task_id>', methods=['GET','POST'])
def informatics_task(task_id):
    r = require_login()
    if r: return r
    task = db.session.get(InformaticsTask, task_id)
    if not task:
        return 'Úkol nebyl nalezen.', 404
    if not informatics_task_unlocked(task):
        flash('Nejdřív dokonči předchozí úkol na 100 %.')
        return redirect(url_for('informatics_lesson', lesson_id=task.lesson_id))

    if current_user().role == 'student':
        begin_focus_attempt('informatics', task.id)

    feedback = None
    preview = None
    percent = None
    grade = None
    checked_submission = None

    if request.method == 'POST':
        if current_user().role != 'student':
            flash('Soubor odevzdává student.')
            return redirect(url_for('informatics_task', task_id=task.id))

        action = request.form.get('action','check')

        if action == 'submit_existing':
            sid = int(request.form.get('submission_id',0) or 0)
            row = InformaticsSubmission.query.filter_by(
                id=sid, user_id=current_user().id, task_id=task.id, status='kontrola'
            ).first()
            if not row:
                flash('Nejdřív svůj soubor zkontroluj.')
                return redirect(url_for('informatics_task', task_id=task.id))
            row.status = 'odevzdáno'
            row.grade = informatics_grade_from_percent(row.percent)
            row.focus_lost = consume_focus_count('informatics', task.id)
            row.created_at = datetime.utcnow()
            db.session.commit()
            flash(f'Úkol byl odevzdán: {row.percent} %, známka {row.grade}.')
            return redirect(url_for('informatics_lesson', lesson_id=task.lesson_id))

        f = request.files.get('student_file')
        if not f or not f.filename:
            flash('Vyber svůj hotový soubor.')
        else:
            expected = Path(task.source_original).suffix.lower()
            got = Path(f.filename).suffix.lower()
            if got != expected:
                flash(f'Nahraj soubor typu {expected}. Vybral jsi {got or "soubor bez přípony"}.')
            else:
                user_dir = INFORMATICS_SUBMISSION_DIR / str(current_user().id)
                stored = _save_uploaded_file(f, user_dir, f'task{task.id}')
                path = user_dir / stored
                token = _extract_work_token(path, got)
                work_row = InformaticsWorkFile.query.filter_by(user_id=current_user().id, task_id=task.id).first()
                if not token or not work_row or token != work_row.token or not _informatics_verify_token(token, current_user().id, task.id):
                    try: path.unlink(missing_ok=True)
                    except Exception: pass
                    flash('Tento soubor nepochází z tvého pracovního souboru staženého z této stránky. Stáhni si svůj soubor tlačítkem „Stáhnout můj pracovní soubor“ a pracuj v něm.')
                    return redirect(url_for('informatics_task', task_id=task.id))
                feedback = evaluate_informatics_file(path, f.filename, task)
                ok_count = sum(1 for x in feedback if x['ok'])
                percent = round(ok_count / max(len(feedback),1) * 100)
                grade = informatics_grade_from_percent(percent)
                preview = informatics_preview(path, f.filename)
                checked_submission = InformaticsSubmission(
                    user_id=current_user().id, task_id=task.id, original_name=f.filename,
                    stored_name=stored, feedback_json=json.dumps(feedback, ensure_ascii=False),
                    percent=percent, grade=grade, status='kontrola',
                    focus_lost=get_focus_count('informatics', task.id)
                )
                db.session.add(checked_submission)
                db.session.commit()

    if current_user().role == 'student' and feedback is None:
        last = InformaticsSubmission.query.filter_by(
            user_id=current_user().id, task_id=task.id, status='kontrola'
        ).order_by(InformaticsSubmission.created_at.desc()).first()
        if last:
            feedback = _safe_json(last.feedback_json, [])
            percent = last.percent
            grade = informatics_grade_from_percent(percent)
            checked_submission = last
            p = INFORMATICS_SUBMISSION_DIR / str(current_user().id) / last.stored_name
            if p.exists():
                preview = informatics_preview(p, last.original_name)

    item = task.lesson
    course = {'subject':'Informatika','grade':item.grade_name,'block':item.topic,'icon':'💻'}
    lesson_obj = {'title':item.title,'block':item.topic}
    teacher_preview = None
    teacher_source = INFORMATICS_SOURCE_DIR / task.source_stored
    if teacher_source.exists():
        teacher_preview = informatics_preview(teacher_source, task.source_original, teacher=True)
    html_content = render_informatics_html(item.html_stored)
    return render_template('informatics_task.html', course=course, lesson=lesson_obj,
                           item=item, task=task, feedback=feedback, preview=preview,
                           teacher_preview=teacher_preview, percent=percent, grade=grade,
                           checked_submission=checked_submission, html_content=html_content)



# ============================================================
# MATEMATIKA – KROKOVÝ ENGINE
# ============================================================

def normalize_math_answer(value):
    value = str(value or '').strip().lower()
    value = value.replace(' ', '').replace(',', '.')
    value = value.replace('−','-').replace('–','-').replace(':','/')

    # Mocniny: přijímáme x^2 i x**2.
    supers = str.maketrans({'⁰':'0','¹':'1','²':'2','³':'3','⁴':'4','⁵':'5','⁶':'6','⁷':'7','⁸':'8','⁹':'9'})
    value = re.sub(r'([a-zA-Z0-9_)])([⁰¹²³⁴⁵⁶⁷⁸⁹]+)', lambda m: m.group(1) + '**' + m.group(2).translate(supers), value)
    value = value.replace('^', '**')

    # Odmocniny:
    # √x      -> sqrt(x)
    # √(x+1)  -> sqrt(x+1)
    # sqrt(x) zůstává beze změny.
    import re as _re

    # √(něco)
    value = _re.sub(r'√\(([^()]*)\)', r'sqrt(\1)', value)

    # √číslo nebo √proměnná
    value = _re.sub(r'√([a-zA-Z0-9_.]+)', r'sqrt(\1)', value)

    return value


def validate_math_expression(value):
    """Ověří, že učitelův matematický zápis umíme přečíst.

    V38: podporuje i více matematických řádků (např. soustavu rovnic).
    Jednotlivé řádky lze oddělit novým řádkem nebo středníkem.
    """
    raw = str(value or '').strip()
    if not raw:
        return False, 'zápis je prázdný'

    # V38: více rovnic / výrazů v jednom matematickém vzoru.
    # Každý řádek validujeme stejným parserem jako samostatný příklad.
    multi_parts = [p.strip() for p in re.split(r'[\r\n;]+', raw) if p.strip()]
    if len(multi_parts) > 1:
        for idx, part in enumerate(multi_parts, 1):
            ok, why = validate_math_expression(part)
            if not ok:
                return False, f'řádek {idx}: {why}'
        return True, ''

    txt = normalize_math_answer(raw)

    # V35: matematický vzor generátoru může být i seznam přiřazení,
    # např. a=3,b=4,c=5. Je to praktické pro slovní úlohy; student tento
    # technický vzor nemusí vidět.
    if ',' in txt:
        assignments=[p.strip() for p in txt.split(',') if p.strip()]
        if assignments and all(re.fullmatch(r'[A-Za-z][A-Za-z0-9_]*\s*=\s*.+', p) for p in assignments):
            try:
                import sympy as sp
                from sympy.parsing.sympy_parser import parse_expr, standard_transformations, implicit_multiplication_application
                transformations = standard_transformations + (implicit_multiplication_application,)
                local_dict={'sqrt':sp.sqrt,'log':sp.log,'ln':sp.log,'sin':sp.sin,'cos':sp.cos,'tan':sp.tan,'abs':sp.Abs,'pi':sp.pi}
                for assignment in assignments:
                    rhs=assignment.split('=',1)[1].strip()
                    parse_expr(rhs, transformations=transformations, evaluate=False, local_dict=local_dict)
                return True, ''
            except Exception:
                return False, 'tomuto matematickému vzoru aplikace nerozumí'

    # V29: učitel může zapisovat stupně jako 30deg, 45deg, ...
    # Pro validaci je převedeme na radiánový zápis, kterému SymPy rozumí.
    # Původní text zůstává beze změny pro studentský renderer, kde je celý
    # úhel (včetně znaku °) pevná konstrukce.
    txt = re.sub(r'(\d+(?:\.\d+)?)deg\b', r'(\1*pi/180)', txt, flags=re.IGNORECASE)
    try:
        import sympy as sp
        from sympy.parsing.sympy_parser import parse_expr, standard_transformations, implicit_multiplication_application
        transformations = standard_transformations + (implicit_multiplication_application,)
        local_dict = {
            'sqrt': sp.sqrt, 'log': sp.log, 'ln': sp.log,
            'sin': sp.sin, 'cos': sp.cos, 'tan': sp.tan,
            'asin': sp.asin, 'acos': sp.acos, 'atan': sp.atan,
            'abs': sp.Abs, 'pi': sp.pi,
            'diff': sp.diff, 'integrate': sp.integrate,
        }
        # Samostatné '=' bereme jako rovnici; <=, >=, != necháme SymPy jako relaci.
        eq = re.search(r'(?<![<>=!])=(?!=)', txt)
        parts = [txt[:eq.start()], txt[eq.end():]] if eq else [txt]
        if any(not part for part in parts):
            return False, 'chybí část výrazu před nebo za ='
        for part in parts:
            parse_expr(part, transformations=transformations, evaluate=False, local_dict=local_dict)
        return True, ''
    except Exception:
        return False, 'tomuto matematickému zápisu aplikace nerozumí'


def format_math_display(value):
    """Vizuální převod učitelského matematického zápisu pro studentskou stránku."""
    text = str(value or '')
    text = re.sub(r'(\d+(?:[.,]\d+)?)\s*deg\b', lambda m: m.group(1) + '°', text, flags=re.IGNORECASE)
    text = re.sub(r'\balpha\b', 'α', text, flags=re.IGNORECASE)
    text = re.sub(r'\b(asin|acos|atan)\(([^()]+)\)', lambda m: {'asin':'sin⁻¹','acos':'cos⁻¹','atan':'tan⁻¹'}[m.group(1).lower()]+' '+m.group(2), text, flags=re.IGNORECASE)
    text = re.sub(r'\b(sin|cos|tan)\(([^()]+)\)', r'\1 \2', text, flags=re.IGNORECASE)
    return text

app.jinja_env.filters['math_display'] = format_math_display

def math_input_layout(expected):
    """Vytvoří strom studentských políček.

    Student píše pouze číslice, proměnné a běžné operátory. Speciální
    matematické konstrukce (odmocniny, funkce, integrály, derivace, | |,
    pí, závorky, zlomková čára) vykresluje aplikace napevno.
    """
    raw = str(expected or '').replace('−', '-').replace('–', '-')
    # V40: správná odpověď může mít více řádků (např. x=2 a y=3).
    # Každý řádek se vykreslí zvlášť, ale všechna pole zůstávají součástí jednoho kroku.
    raw_lines = [p.strip() for p in re.split(r'[\r\n;]+', raw) if p.strip()] or ['']
    compact_lines = [''.join(ch for ch in line if not ch.isspace()) for line in raw_lines]
    compact_lines = [re.sub(r'([A-Za-z0-9_.]+)\*\*\(1/2\)', r'sqrt(\1)', line) for line in compact_lines]
    # Učitelský zápis úhlů: 30deg -> 30°. Úhel i znak stupně jsou pro studenta pevné.
    field_index = 0

    def field(ch, super_=False):
        nonlocal field_index
        node={'kind':'field','index':field_index,'expected':ch,'super':super_}
        field_index += 1
        return node

    def split_args(text):
        args=[]; depth=0; last=0
        for j,ch in enumerate(text):
            if ch=='(': depth+=1
            elif ch==')': depth-=1
            elif ch==',' and depth==0:
                args.append(text[last:j]); last=j+1
        args.append(text[last:])
        return args

    def matching_paren(text, pos):
        depth=0
        for j in range(pos,len(text)):
            if text[j]=='(': depth+=1
            elif text[j]==')':
                depth-=1
                if depth==0: return j
        return -1

    special={'sqrt','log','ln','sin','cos','tan','asin','acos','atan','abs','integrate','diff'}

    def parse(text, super_mode=False):
        nodes=[]; i=0
        while i < len(text):
            # Speciální funkce: název a konstrukce jsou pevné, argumenty se dále rozloží.
            m=re.match(r'([A-Za-z]+)\(', text[i:])
            if m and m.group(1).lower() in special:
                name=m.group(1).lower(); op=i+len(m.group(1)); close=matching_paren(text,op)
                if close!=-1:
                    args=split_args(text[op+1:close])
                    if name=='sqrt' and args:
                        nodes.append({'kind':'sqrt','body':parse(args[0])})
                    elif name=='abs' and args:
                        nodes.append({'kind':'abs','body':parse(args[0])})
                    elif name=='integrate' and args:
                        nodes.append({'kind':'integral','body':parse(args[0]),'var':parse(args[1]) if len(args)>1 else []})
                    elif name=='diff' and args:
                        nodes.append({'kind':'derivative','body':parse(args[0]),'var':parse(args[1]) if len(args)>1 else [],'order':parse(args[2], True) if len(args)>2 else []})
                    else:
                        display_name={'asin':'sin⁻¹','acos':'cos⁻¹','atan':'tan⁻¹'}.get(name,name)
                        nodes.append({'kind':'function','name':name,'display_name':display_name,'args':[parse(a) for a in args]})
                    i=close+1; continue
            if text.startswith('pi',i) and (i+2==len(text) or not text[i+2].isalpha()):
                nodes.append({'kind':'fixed','display':'π','answer':'pi'}); i+=2; continue
            # Řecký úhel alpha je speciální pevný symbol.
            if text.startswith('alpha',i) and (i+5==len(text) or not text[i+5].isalpha()):
                nodes.append({'kind':'fixed','display':'α','answer':'alpha'}); i+=5; continue
            # Stupně ve SPRÁVNÉ ODPOVĚDI: číslo vyplní student, znak ° je pevný.
            # Např. 30deg -> [3][0]° a 117deg -> [1][1][7]°.
            dm=re.match(r'(\d+(?:[.,]\d+)?)\s*deg\b', text[i:], re.IGNORECASE)
            if dm:
                number=dm.group(1)
                for digit in number:
                    if digit in '.,':
                        # desetinný oddělovač zůstává součástí odpovědi jako běžné políčko
                        nodes.append(field(digit, super_mode))
                    else:
                        nodes.append(field(digit, super_mode))
                nodes.append({'kind':'fixed','display':'°','answer':'deg'})
                i+=len(dm.group(0)); continue
            # Mocnina: ** je pevná konstrukce a exponent je horní index s inputy.
            if text.startswith('**',i):
                i+=2
                if i<len(text) and text[i]=='(':
                    close=matching_paren(text,i)
                    exp=text[i+1:close] if close!=-1 else text[i+1:]
                    i=(close+1 if close!=-1 else len(text))
                    nodes.append({'kind':'power','exp':parse(exp, True),'paren':True})
                else:
                    mexp=re.match(r'[A-Za-z0-9.+\-]+',text[i:])
                    exp=mexp.group(0) if mexp else (text[i:i+1] or '')
                    i+=len(exp)
                    nodes.append({'kind':'power','exp':parse(exp, True),'paren':False})
                continue
            # Jednoduchý zlomek: čára je pevná.
            fm=re.match(r'([A-Za-z0-9]+)/(?!/)([A-Za-z0-9]+)',text[i:])
            if fm:
                nodes.append({'kind':'fraction','num':parse(fm.group(1)),'den':parse(fm.group(2))})
                i+=len(fm.group(0)); continue
            ch=text[i]
            if ch in '(),[]{}':
                nodes.append({'kind':'fixed','display':ch,'answer':ch}); i+=1; continue
            # Operátory a relační znaky student doplňuje sám.
            if ch in '+-*/:<>=!':
                nodes.append(field(ch,super_mode)); i+=1; continue
            if ch=='√':
                nodes.append({'kind':'fixed','display':'√','answer':'√'}); i+=1; continue
            # Každá číslice a každé písmeno/proměnná = vlastní input.
            nodes.append(field(ch,super_mode)); i+=1
        return nodes

    def is_simple_result_assignment(line):
        m=re.match(r'^(alpha|[A-Za-z])=(.+)$', line, re.IGNORECASE)
        if not m:
            return None
        lhs, rhs=m.group(1), m.group(2)
        probe=re.sub(r'(?i)(sqrt|log|ln|sin|cos|tan|asin|acos|atan|abs|pi|deg)', '', rhs)
        if re.search(r'[A-Za-z]', probe):
            return None
        return lhs, rhs

    tokens=[]
    for line_no, compact in enumerate(compact_lines):
        if line_no:
            tokens.append({'kind':'linebreak'})
        assignment=is_simple_result_assignment(compact)
        if assignment:
            lhs,rhs=assignment
            tokens.append({'kind':'fixed','display':'α' if lhs.lower()=='alpha' else lhs,'answer':lhs})
            tokens.append({'kind':'fixed','display':'=','answer':'='})
            tokens.extend(parse(rhs))
        else:
            tokens.extend(parse(compact))
    return {'tokens':tokens,'field_count':field_index}


def _math_answer_nodes(nodes, form):
    out=[]
    for tok in nodes:
        k=tok['kind']
        if k=='field':
            val=(form.get(f'math_char_{tok["index"]}') or '').strip()
            out.append(val[:1])
        elif k=='fixed': out.append(tok.get('answer',''))
        elif k=='fraction': out.append(_math_answer_nodes(tok['num'],form)+'/'+_math_answer_nodes(tok['den'],form))
        elif k=='sqrt': out.append('sqrt('+_math_answer_nodes(tok['body'],form)+')')
        elif k=='abs': out.append('abs('+_math_answer_nodes(tok['body'],form)+')')
        elif k=='function': out.append(tok['name']+'('+','.join(_math_answer_nodes(a,form) for a in tok['args'])+')')
        elif k=='integral': out.append('integrate('+_math_answer_nodes(tok['body'],form)+','+_math_answer_nodes(tok['var'],form)+')')
        elif k=='derivative':
            args=[_math_answer_nodes(tok['body'],form),_math_answer_nodes(tok['var'],form)]
            if tok.get('order'): args.append(_math_answer_nodes(tok['order'],form))
            out.append('diff('+','.join(args)+')')
        elif k=='power':
            exp=_math_answer_nodes(tok['exp'],form)
            out.append('**'+(('('+exp+')') if tok.get('paren') else exp))
        elif k=='linebreak':
            out.append('\n')
    return ''.join(out)


def math_answer_from_fields(expected, form):
    """Složí studentská políčka zpět do matematického zápisu."""
    layout=math_input_layout(expected)
    return _math_answer_nodes(layout['tokens'],form)


def render_math_input_layout(layout):
    """Vykreslí matematiku: speciální konstrukce pevně, obsah jako jednotlivé inputy."""
    def field_html(tok):
        cls='math-char math-char-super' if tok.get('super') else 'math-char'
        return f'<input class="{cls}" name="math_char_{tok["index"]}" maxlength="1" autocomplete="off" required aria-label="matematický znak {tok["index"]+1}">'
    def render(nodes):
        parts=[]
        for tok in nodes:
            k=tok['kind']
            if k=='field': parts.append(field_html(tok))
            elif k=='fixed': parts.append(f'<span class="math-fixed">{escape(tok.get("display",""))}</span>')
            elif k=='fraction': parts.append(f'<span class="math-fraction"><span class="math-num">{render(tok["num"])}</span><span class="math-den">{render(tok["den"])}</span></span>')
            elif k=='sqrt': parts.append(f'<span class="math-root"><span class="math-root-sign">√</span><span class="math-radicand">{render(tok["body"])}</span></span>')
            elif k=='abs': parts.append(f'<span class="math-abs">|&nbsp;{render(tok["body"])}&nbsp;|</span>')
            elif k=='function':
                args='<span class="math-fixed">,</span>'.join(render(a) for a in tok['args'])
                parts.append(f'<span class="math-function"><span class="math-fixed">{escape(tok.get("display_name",tok["name"]))}</span><span class="math-fixed">(</span>{args}<span class="math-fixed">)</span></span>')
            elif k=='integral': parts.append(f'<span class="math-integral"><span class="math-special">∫</span>{render(tok["body"])}<span class="math-fixed">d</span>{render(tok["var"])}</span>')
            elif k=='derivative':
                var=render(tok['var']); order=render(tok.get('order',[]))
                top='d'+(f'<sup>{order}</sup>' if order else '')
                bottom='d'+var+(f'<sup>{order}</sup>' if order else '')
                parts.append(f'<span class="math-derivative"><span class="math-deriv-frac"><span>{top}</span><span>{bottom}</span></span><span class="math-fixed">(</span>{render(tok["body"])}<span class="math-fixed">)</span></span>')
            elif k=='power': parts.append(f'<sup class="math-power">{render(tok["exp"])}</sup>')
            elif k=='linebreak': parts.append('</div><div class="math-input-line">')
        return ''.join(parts)
    return Markup('<div class="math-input-line">'+render(layout.get('tokens',[]))+'</div>')

def math_answers_equivalent(student_value, expected_value):
    """Porovnává matematický význam, ne přesný text.

    Příklady, které uzná jako ekvivalentní:
    x + 3 = 4   ↔   3 + x = 4
    x = 1       ↔   -1 = -x
    2*x = 2     ↔   x = 1

    Když SymPy výraz neumí bezpečně zpracovat, použije se původní
    normalizované textové porovnání.
    """
    # V40: více výsledků v jednom kroku porovnáváme po řádcích.
    # Pořadí řádků odpovídá pořadí výsledkových vzorců z učitelského editoru.
    a_parts=[p.strip() for p in re.split(r'[\r\n;]+', str(student_value or '')) if p.strip()]
    b_parts=[p.strip() for p in re.split(r'[\r\n;]+', str(expected_value or '')) if p.strip()]
    if len(a_parts)>1 or len(b_parts)>1:
        if len(a_parts) != len(b_parts):
            return False
        return all(math_answers_equivalent(x, y) for x,y in zip(a_parts,b_parts))

    a = normalize_math_answer(student_value)
    b = normalize_math_answer(expected_value)

    if a == b:
        return True

    try:
        import sympy as sp
        from sympy.parsing.sympy_parser import (
            parse_expr,
            standard_transformations,
            implicit_multiplication_application,
        )

        transformations = standard_transformations + (implicit_multiplication_application,)

        def parse_side(txt):
            local_dict = {
                'sqrt': sp.sqrt, 'log': sp.log, 'ln': sp.log,
                'sin': sp.sin, 'cos': sp.cos, 'tan': sp.tan,
                'abs': sp.Abs, 'pi': sp.pi, 'diff': sp.diff, 'integrate': sp.integrate,
            }
            return parse_expr(
                txt,
                transformations=transformations,
                evaluate=True,
                local_dict=local_dict
            )

        def relation_to_expr(txt):
            # Rovnici převedeme na výraz levá - pravá = 0.
            if '=' in txt:
                left, right = txt.split('=', 1)
                return sp.simplify(parse_side(left) - parse_side(right)), True
            return sp.simplify(parse_side(txt)), False

        expr_a, is_eq_a = relation_to_expr(a)
        expr_b, is_eq_b = relation_to_expr(b)

        # Rovnici s výrazem nemícháme.
        if is_eq_a != is_eq_b:
            return False

        # Běžné algebraické přeuspořádání.
        if sp.simplify(expr_a - expr_b) == 0:
            return True

        if not is_eq_a:
            return bool(sp.simplify(expr_a - expr_b) == 0)

        # U rovnic uznáme i násobek celé rovnice nenulovou konstantou:
        # 2x=2 je totéž jako x=1.
        if expr_b != 0:
            ratio = sp.simplify(expr_a / expr_b)
            if ratio != 0 and not getattr(ratio, 'free_symbols', set()):
                return True

        # Poslední kontrola: porovnáme množinu řešení pro všechny proměnné.
        symbols = sorted(expr_a.free_symbols | expr_b.free_symbols, key=lambda x: x.name)
        if len(symbols) == 1:
            x = symbols[0]
            sol_a = sp.solveset(expr_a, x, domain=sp.S.Reals)
            sol_b = sp.solveset(expr_b, x, domain=sp.S.Reals)
            return sol_a == sol_b

    except Exception:
        pass

    return a == b


def _variant_number_text(value):
    """Pěkný zápis náhodné hodnoty bez zbytečných .0."""
    try:
        f = float(value)
        if abs(f - round(f)) < 1e-10:
            return str(int(round(f)))
        return (f'{f:.10f}').rstrip('0').rstrip('.')
    except Exception:
        return str(value)


def _parse_variant_values(raw):
    """Hodnoty učitelského vzoru oddělené středníkem nebo novým řádkem."""
    parts = [x.strip().replace(',', '.') for x in re.split(r'[;\n]+', str(raw or '')) if x.strip()]
    vals=[]
    for x in parts:
        try:
            vals.append(float(x))
        except Exception:
            raise ValueError(f'„{x}“ není číslo')
    return vals


def _parse_variant_variable_names(raw):
    """Nový režim: učitel může místo vzorových čísel zadat n1;n2;n3;..."""
    parts=[x.strip() for x in re.split(r'[;\n]+', str(raw or '')) if x.strip()]
    if not parts or not all(re.fullmatch(r'n[1-9]\d*', x) for x in parts):
        return []
    expected=[f'n{i+1}' for i in range(len(parts))]
    if parts != expected:
        raise ValueError('proměnné musí být postupně n1; n2; n3; ...')
    return parts




def _is_variable_declaration(text, variable_names):
    parts=[x.strip() for x in re.split(r'[;\n]+', str(text or '')) if x.strip()]
    return bool(variable_names) and parts == list(variable_names)

def _replace_variant_variables(text_value, env):
    """Dosadí n1, n2, ... jako samostatné matematické tokeny a také {n1} v textu."""
    out=str(text_value or '')
    for name,value in sorted((env or {}).items(), key=lambda kv: len(str(kv[0])), reverse=True):
        val=_variant_number_text(value)
        out=out.replace('{'+str(name)+'}', val)
        out=re.sub(rf'(?<![A-Za-z0-9_]){re.escape(str(name))}(?![A-Za-z0-9_])', val, out)
    return out


def _formula_assignment(formula):
    """Vrátí (alias, výraz). Podporuje např. h=n1/n3 i původní degrees(...)."""
    raw=str(formula or '').strip()
    m=re.fullmatch(r'([A-Za-z][A-Za-z0-9_]*)\s*=\s*(.+)', raw, flags=re.S)
    return (m.group(1), m.group(2).strip()) if m else (None, raw)


def _formula_returns_degrees(expr):
    """Pozná pomocný výsledek typu degrees(...), aby student dostal pevný znak °."""
    return bool(re.match(r'^\s*degrees\s*\(', str(expr or ''), flags=re.IGNORECASE))


def _render_alias_in_text(text, alias, value, expr, decimals):
    """Dosadí pomocnou proměnnou do učitelova správného stavu.

    U d=degrees(...) se automaticky vytvoří např. 36.87deg. Studentský renderer
    pak ponechá α, = a ° napevno a student doplňuje jen číslice výsledku.
    """
    rendered=_format_result_number(value, decimals)
    if _formula_returns_degrees(expr):
        rendered += 'deg'
    return re.sub(rf'(?<![A-Za-z0-9_]){re.escape(alias)}(?![A-Za-z0-9_])', rendered, str(text or ''))


def _evaluate_variable_steps(example, base_env, render=False):
    """Vyhodnotí kroky postupně a dovolí navazující pomocné proměnné.

    Např. h=n1/n3 v kroku 1 vytvoří číselné h pro další kroky a
    d=degrees(asin(h)) v kroku 2 na něj může přímo navázat.
    """
    env=dict(base_env or {})
    values=[]
    rendered_steps={}
    for st in sorted(example.steps, key=lambda x: x.order):
        expected=_replace_variant_variables(st.expected, env)
        instruction=_replace_variant_variables(st.instruction, env)
        hint=_replace_variant_variables(st.hint, env)
        for formula in _split_result_formulas(getattr(st, 'result_formula', '') or ''):
            alias,expr=_formula_assignment(formula)
            value=_safe_variant_formula(expr, env)
            values.append(value)
            if alias:
                env[alias]=value
                expected=_render_alias_in_text(expected, alias, value, expr, getattr(st, 'result_decimals', 2))
                instruction=_render_alias_in_text(instruction, alias, value, expr, getattr(st, 'result_decimals', 2))
                hint=_render_alias_in_text(hint, alias, value, expr, getattr(st, 'result_decimals', 2))
        # Po výpočtu pomocných proměnných ještě jednou dosaď všechna n1,n2,...
        expected=_replace_variant_variables(expected, env)
        instruction=_replace_variant_variables(instruction, env)
        hint=_replace_variant_variables(hint, env)
        if render:
            rendered_steps[st.id]={'instruction':instruction,'expected':expected,'hint':hint}
    return values, env, rendered_steps


def _render_variable_mode_step(st, env):
    """Zachováno pro kompatibilitu; nový režim používá sekvenční pipeline."""
    expected=_replace_variant_variables(st.expected, env)
    return {
        'instruction': _replace_variant_variables(st.instruction, env),
        'expected': expected,
        'hint': _replace_variant_variables(st.hint, env),
    }


def _safe_variant_condition(condition, env):
    """Bezpečně vyhodnotí obecnou matematickou podmínku n1, n2, ..."""
    cond = str(condition or '').strip()
    if not cond:
        return True
    tree = ast.parse(cond, mode='eval')
    allowed_nodes=(ast.Expression, ast.BoolOp, ast.BinOp, ast.UnaryOp, ast.Compare, ast.Name, ast.Load,
                   ast.Constant, ast.And, ast.Or, ast.Not, ast.Add, ast.Sub, ast.Mult, ast.Div,
                   ast.FloorDiv, ast.Mod, ast.Pow, ast.USub, ast.UAdd, ast.Eq, ast.NotEq,
                   ast.Lt, ast.LtE, ast.Gt, ast.GtE, ast.Call)
    allowed_funcs={'abs':abs,'min':min,'max':max,'round':round,'sqrt':math.sqrt}
    for node in ast.walk(tree):
        if not isinstance(node, allowed_nodes):
            raise ValueError('nepovolená konstrukce v podmínce')
        if isinstance(node, ast.Name) and node.id not in env and node.id not in allowed_funcs:
            raise ValueError(f'neznámá proměnná {node.id}')
        if isinstance(node, ast.Call):
            if not isinstance(node.func, ast.Name) or node.func.id not in allowed_funcs:
                raise ValueError('nepovolená funkce v podmínce')
    return bool(eval(compile(tree, '<math-condition>', 'eval'), {'__builtins__':{}, **allowed_funcs}, env))


def _safe_variant_formula(formula, env):
    """Bezpečně spočítá číselný výsledek z n1, n2, ... bez eval přístupu k Pythonu."""
    expr = str(formula or '').strip()
    _alias, expr = _formula_assignment(expr)
    if not expr:
        raise ValueError('vzorec pro výpočet výsledku je prázdný')
    tree = ast.parse(expr, mode='eval')
    allowed_nodes=(ast.Expression, ast.BinOp, ast.UnaryOp, ast.Name, ast.Load, ast.Constant,
                   ast.Add, ast.Sub, ast.Mult, ast.Div, ast.FloorDiv, ast.Mod, ast.Pow,
                   ast.USub, ast.UAdd, ast.Call)
    allowed_funcs={
        'abs':abs, 'min':min, 'max':max, 'round':round, 'sqrt':math.sqrt,
        'sin':math.sin, 'cos':math.cos, 'tan':math.tan,
        'asin':math.asin, 'acos':math.acos, 'atan':math.atan,
        'degrees':math.degrees, 'radians':math.radians,
        'log':math.log, 'ln':math.log, 'exp':math.exp
    }
    constants={'pi':math.pi, 'e':math.e}
    for node in ast.walk(tree):
        if not isinstance(node, allowed_nodes):
            raise ValueError('nepovolená konstrukce ve vzorci výsledku')
        if isinstance(node, ast.Name) and node.id not in env and node.id not in allowed_funcs and node.id not in constants:
            raise ValueError(f'neznámá proměnná {node.id}')
        if isinstance(node, ast.Call):
            if not isinstance(node.func, ast.Name) or node.func.id not in allowed_funcs:
                raise ValueError('nepovolená funkce ve vzorci výsledku')
    value=eval(compile(tree, '<math-result-formula>', 'eval'), {'__builtins__':{}, **allowed_funcs, **constants}, env)
    try:
        value=float(value)
    except Exception:
        raise ValueError('vzorec výsledku musí vrátit jedno číslo')
    if not math.isfinite(value):
        raise ValueError('vzorec výsledku nevrátil konečné číslo')
    return value


def _format_result_number(value, decimals=2):
    try:
        decimals=max(0, min(10, int(decimals)))
    except Exception:
        decimals=2
    rounded=round(float(value), decimals)
    if decimals == 0:
        return str(int(round(rounded)))
    return f'{rounded:.{decimals}f}'.rstrip('0').rstrip('.')


def _split_result_formulas(raw):
    """Jeden nebo více obecných vzorců, oddělených středníkem nebo novým řádkem."""
    return [x.strip() for x in re.split(r'[;\n]+', str(raw or '')) if x.strip()]


def _computed_result_markers(expected_text, source_env, target_env, formulas, decimals):
    """Nahradí vypočtené vzorové výsledky dočasnými značkami.

    Díky značkám se následná výměna vstupních náhodných čísel nemůže omylem
    dotknout právě dopočítaného výsledku. Podporuje libovolný počet výsledků.
    """
    text=str(expected_text or '')
    replacements=[]
    for idx, formula in enumerate(_split_result_formulas(formulas), 1):
        base=_format_result_number(_safe_variant_formula(formula, source_env), decimals)
        target=_format_result_number(_safe_variant_formula(formula, target_env), decimals)
        candidates=[base]
        if '.' in base:
            candidates.append(base.replace('.', ','))
        # V42: interní značka nesmí obsahovat číslice. Obecný generátor totiž
        # nahrazuje číselné tokeny ze vzoru a značka RESULT_2 se mohla omylem
        # změnit, pokud mezi měněnými hodnotami bylo číslo 2. Index proto
        # kódujeme pouze písmeny (A, B, ... AA, AB...).
        n=idx
        letters=''
        while n:
            n, rem = divmod(n-1, 26)
            letters = chr(65+rem) + letters
        marker=f'__UCEBNICE_RESULT_{letters}__'
        found=False
        for token in candidates:
            pattern=rf'(?<![\d.]){re.escape(token)}(?![\d.])'
            if re.search(pattern, text):
                text=re.sub(pattern, marker, text, count=1)
                found=True
                break
        if not found:
            raise ValueError(f've správné odpovědi nebyl nalezen vypočtený vzorový výsledek {base}')
        replacements.append((marker,target))
    return text, replacements


def _restore_computed_result_markers(text, replacements):
    out=str(text or '')
    for marker,target in replacements:
        out=out.replace(marker,target)
    return out

def _variant_slot_name(index):
    """Bezpečná interní značka bez číslic, aby se nikdy nestřetla s matematickými hodnotami."""
    n=int(index)+1
    letters=''
    while n:
        n, rem=divmod(n-1, 26)
        letters=chr(65+rem)+letters
    return f'__UCEBNICE_VARIANT_{letters}__'


def _mark_variant_numbers(text_value, source_values):
    """Převede učitelovy číselné hodnoty na stabilní interní sloty.

    V45: nejprve se vytvoří šablona a až potom se do slotů vloží nová čísla.
    Díky tomu se žádná nově vložená hodnota nemůže znovu přepsat a znaménka,
    desetinné tečky i struktura zápisu zůstávají přesně zachované.
    """
    out=str(text_value or '')
    # Stejná zdrojová hodnota reprezentuje stejný alias všude v textu.
    value_to_marker={}
    for idx, old in enumerate(source_values):
        old_txt=_variant_number_text(old)
        value_to_marker.setdefault(old_txt, _variant_slot_name(idx))
    alternatives='|'.join(re.escape(k) for k in sorted(value_to_marker, key=len, reverse=True))
    if not alternatives:
        return out, value_to_marker
    pattern=re.compile(rf'(?<![\d.])(?:{alternatives})(?![\d.])')
    return pattern.sub(lambda m:value_to_marker[m.group(0)], out), value_to_marker


def _fill_variant_slots(marked_text, source_values, new_values):
    out=str(marked_text or '')
    marker_to_value={}
    seen={}
    for idx,(old,new) in enumerate(zip(source_values,new_values)):
        old_txt=_variant_number_text(old)
        marker=seen.setdefault(old_txt, _variant_slot_name(idx))
        marker_to_value.setdefault(marker, _variant_number_text(new))
    for marker,value in marker_to_value.items():
        out=out.replace(marker,value)
    return out


def _replace_variant_numbers(text_value, source_values, new_values):
    """Obecná bezpečná náhrada přes interní šablonové sloty (V45)."""
    marked,_=_mark_variant_numbers(text_value, source_values)
    return _fill_variant_slots(marked, source_values, new_values)


def _render_variant_prose(text_value, env):
    """Slovní zadání je oddělené od matematického vzoru.

    V46: náhodné hodnoty se do běžného textu vkládají pouze explicitními
    značkami {n1}, {n2}, ... . Tím se nemůže stát, že staré číslo nebo rovnice
    ve slovním zadání vytvoří jinou úlohu než matematický vzor kontrolovaný enginem.
    Čistě matematické řádky ve starších slovních zadáních skryjeme, protože
    matematický vzor se studentovi nově zobrazuje samostatně.
    """
    text=str(text_value or '')
    for name,value in (env or {}).items():
        text=text.replace('{'+str(name)+'}', _variant_number_text(value))
    kept=[]
    for line in text.splitlines():
        stripped=line.strip()
        if stripped and any(ch in stripped for ch in '=<>'):
            try:
                ok,_=validate_math_expression(stripped)
            except Exception:
                ok=False
            if ok:
                continue
        kept.append(line)
    return '\n'.join(kept).strip()

def _variant_computed_values(example, env):
    """Dopočítá všechny pomocné výsledky v pořadí kroků včetně návazností."""
    values, _env, _steps = _evaluate_variable_steps(example, env, render=False)
    return values


def _variant_filter_settings(example):
    """Normalizované obecné nastavení kvality výsledků."""
    kind=str(getattr(example, 'variant_result_kind', 'any') or 'any').strip().lower()
    # toleruj i starší/lokalizované hodnoty, aby filtr nikdy omylem nespadl na 'any'
    if kind in ('integer','int','whole','cela','celá','cela_cisla','celá čísla','cela cisla') or 'cel' in kind:
        kind='integer'
    sign=str(getattr(example, 'variant_result_sign', 'any') or 'any').strip().lower()
    min_v=getattr(example, 'variant_result_min', None)
    max_v=getattr(example, 'variant_result_max', None)
    try:
        decimals=int(getattr(example, 'variant_result_decimals', -1))
    except Exception:
        decimals=-1
    active=(kind!='any' or sign!='any' or min_v is not None or max_v is not None or decimals>=0)
    return kind, sign, min_v, max_v, decimals, active


def _variant_values_match_filter(example, values):
    """Přísná obecná kontrola již vypočtených číselných výsledků."""
    kind, sign, min_v, max_v, decimals, active=_variant_filter_settings(example)
    if not active:
        return True
    if not values:
        raise ValueError('je nastaven filtr výsledků, ale žádný krok nemá vzorec pro přepočet výsledku')
    eps=1e-9
    for raw in values:
        value=float(raw)
        if not math.isfinite(value):
            return False
        if kind=='integer' and abs(value-round(value))>eps:
            return False
        if sign=='nonnegative' and value < -eps:
            return False
        if sign=='positive' and value <= eps:
            return False
        if min_v is not None and value < float(min_v)-eps:
            return False
        if max_v is not None and value > float(max_v)+eps:
            return False
        if decimals >= 0 and abs(value-round(value, decimals))>eps:
            return False
    return True


def _variant_results_ok(example, env):
    """Obecný filtr kvality výsledků; není svázaný s žádným matematickým tématem."""
    _, _, _, _, _, active=_variant_filter_settings(example)
    if not active:
        return True
    return _variant_values_match_filter(example, _variant_computed_values(example, env))


def _generic_math_variant(example, user_id):
    """V47: přísný obecný generátor typu generate -> compute -> filter -> render.

    Kandidát se NESMÍ vykreslit dřív, než jsou z jeho n-hodnot vypočteny všechny
    result_formula a ty projdou obecným filtrem výsledků. Není zde žádná logika
    specifická pro rovnice/soustavy; stejný postup platí pro každý matematický typ.
    """
    variable_names = _parse_variant_variable_names(example.variant_values)
    variable_mode = bool(variable_names)
    source = ([1.0] * len(variable_names)) if variable_mode else _parse_variant_values(example.variant_values)
    if not source:
        raise ValueError('nejsou zadané proměnné ani čísla ze vzoru')

    lo = float(example.variant_min if example.variant_min is not None else 1)
    hi = float(example.variant_max if example.variant_max is not None else 30)
    step = float(example.variant_step if example.variant_step is not None else 1)
    if step <= 0 or hi < lo:
        raise ValueError('neplatný rozsah nebo krok')
    count = int(math.floor((hi-lo)/step + 1e-9)) + 1
    if count < 1 or count > 10000:
        raise ValueError('rozsah generování je příliš velký')

    # Je-li zapnutý filtr výsledků, musí existovat alespoň jeden vzorec výsledku.
    _, _, _, _, _, filter_active = _variant_filter_settings(example)
    formulas_present = any(str(getattr(st, 'result_formula', '') or '').strip() for st in example.steps)
    if filter_active and not formulas_present:
        raise ValueError('je nastaven filtr výsledků, ale chybí vzorec pro přepočet výsledku')

    # V50: skutečné náhodné losování. Varianta už není odvozena deterministicky
    # z user_id. Po prvním vylosování ji uložíme k pokusu žáka, takže při obnově
    # stránky zůstane stejná. Současně se pokud možno vyhneme variantám, které už
    # mají jiní žáci u stejného příkladu.
    attempt = MathAttempt.query.filter_by(user_id=user_id, lesson_id=example.lesson_id).first()
    if not attempt:
        attempt = MathAttempt(user_id=user_id, lesson_id=example.lesson_id)
        db.session.add(attempt)
        db.session.flush()

    saved_variants = _safe_json(getattr(attempt, 'variant_json', '{}') or '{}', {})
    if not isinstance(saved_variants, dict):
        saved_variants = {}
    saved = saved_variants.get(str(example.id))

    chosen = None
    chosen_env = None
    chosen_results = None

    # Jestli už žák variantu jednou dostal a stále vyhovuje aktuálním pravidlům,
    # použijeme ji znovu. Učitel tak může bezpečně obnovit stránku bez změny zadání.
    if isinstance(saved, list) and len(saved) == len(source):
        try:
            vals = [float(v) for v in saved]
            env = {f'n{i+1}': v for i, v in enumerate(vals)}
            results = _variant_computed_values(example, env)
            if (_safe_variant_condition(example.variant_condition, env)
                    and _variant_values_match_filter(example, results)):
                chosen, chosen_env, chosen_results = vals, env, results
        except (TypeError, ZeroDivisionError, ValueError, OverflowError):
            pass

    # Varianty už přidělené jiným žákům. Nejsou absolutním zákazem: když je
    # možných variant méně než žáků, po delším hledání dovolíme opakování.
    used = set()
    if chosen is None:
        for other in MathAttempt.query.filter_by(lesson_id=example.lesson_id).all():
            if other.user_id == user_id:
                continue
            data = _safe_json(getattr(other, 'variant_json', '{}') or '{}', {})
            vals = data.get(str(example.id)) if isinstance(data, dict) else None
            if isinstance(vals, list) and len(vals) == len(source):
                try:
                    used.add(tuple(round(float(v), 12) for v in vals))
                except (TypeError, ValueError):
                    pass

        rng = random.SystemRandom()
        duplicate_candidate = None
        # Prvních 50 000 losů hledá platnou a dosud nepoužitou variantu.
        for _ in range(50000):
            vals = [lo + rng.randrange(count)*step for _ in source]
            env = {f'n{i+1}': v for i, v in enumerate(vals)}
            try:
                if not _safe_variant_condition(example.variant_condition, env):
                    continue
                results = _variant_computed_values(example, env)
                if not _variant_values_match_filter(example, results):
                    continue
                key = tuple(round(float(v), 12) for v in vals)
                if key in used:
                    if duplicate_candidate is None:
                        duplicate_candidate = (vals, env, results)
                    continue
                chosen, chosen_env, chosen_results = vals, env, results
                break
            except (ZeroDivisionError, ValueError, OverflowError):
                continue

        # Pokud jsou všechny snadno dosažitelné platné varianty už rozdané,
        # použijeme náhodně nalezenou duplicitní variantu místo chyby.
        if chosen is None and duplicate_candidate is not None:
            chosen, chosen_env, chosen_results = duplicate_candidate

        if chosen is not None:
            saved_variants[str(example.id)] = [float(v) for v in chosen]
            attempt.variant_json = json.dumps(saved_variants, ensure_ascii=False)
            db.session.commit()

    if chosen is None:
        # Fail closed: nikdy neposílat studentovi neověřenou náhodnou variantu.
        source_env = {f'n{i+1}': v for i, v in enumerate(source)}
        try:
            source_results = _variant_computed_values(example, source_env)
            if (_safe_variant_condition(example.variant_condition, source_env)
                    and _variant_values_match_filter(example, source_results)):
                chosen = list(source)
                chosen_env = source_env
                chosen_results = source_results
            else:
                raise ValueError
        except Exception:
            raise ValueError('v daném rozsahu se nepodařilo najít bezpečnou variantu splňující filtr výsledků')

    # 4. až nyní vykreslit texty a správné odpovědi
    if variable_mode:
        verify_results, final_env, steps = _evaluate_variable_steps(example, chosen_env, render=True)
        rendered_problem=_replace_variant_variables(example.problem, final_env)
        rendered_prose=_replace_variant_variables(getattr(example, 'prose_problem', '') or '', final_env)
        # Pokud je matematický vzor jen technický seznam vstupních proměnných,
        # studentovi ho vůbec neukazujeme.
        if _is_variable_declaration(example.problem, variable_names):
            rendered_problem=''
        if not _variant_values_match_filter(example, verify_results):
            raise ValueError('bezpečnostní kontrola odmítla finální variantu')
        return {
            'problem': rendered_problem, 'prose_problem': rendered_prose, 'steps': steps,
            'variant_values': chosen, 'variant_results': chosen_results, 'variable_mode': True,
            'computed_env': final_env,
        }

    source_env = {f'n{i+1}': v for i, v in enumerate(source)}
    steps = {}
    for st in example.steps:
        expected_text = str(st.expected or '')
        replacements = []
        if str(getattr(st, 'result_formula', '') or '').strip():
            expected_text, replacements = _computed_result_markers(
                expected_text, source_env, chosen_env, st.result_formula,
                getattr(st, 'result_decimals', 2)
            )
        expected_text = _replace_variant_numbers(expected_text, source, chosen)
        expected_text = _restore_computed_result_markers(expected_text, replacements)
        steps[st.id] = {
            'instruction': _replace_variant_numbers(st.instruction, source, chosen),
            'expected': expected_text,
            'hint': _replace_variant_numbers(st.hint, source, chosen)
        }

    rendered_problem = _replace_variant_numbers(example.problem, source, chosen)
    rendered_prose = _replace_variant_numbers(getattr(example, 'prose_problem', '') or '', source, chosen)

    # 5. nezávislá poslední pojistka těsně před returnem.
    verify_results = _variant_computed_values(example, chosen_env)
    if not _variant_values_match_filter(example, verify_results):
        raise ValueError('V47 bezpečnostní kontrola odmítla finální variantu')

    return {
        'problem': rendered_problem,
        'prose_problem': rendered_prose,
        'steps': steps,
        'variant_values': chosen,
        'variant_results': chosen_results,
    }

def generated_math_variant(example, user_id):
    """Stabilní varianta pro žáka. Nový obecný režim má přednost; původní lineární generátor zůstává jako fallback."""
    if getattr(example, 'variant_enabled', False):
        try:
            return _generic_math_variant(example, user_id)
        except Exception as e:
            # V49: proměnný režim se nikdy nesmí vrátit k parseru starých čísel.
            # Chybná / nemožná podmínka se hlásí přímo, místo pádu na 500 při parsování n1;n2;...
            if _parse_variant_variable_names(getattr(example, 'variant_values', '') or ''):
                raise ValueError(f'Náhodnou variantu nelze vytvořit: {e}')
            _, _, _, _, _, filter_active = _variant_filter_settings(example)
            if filter_active:
                raise ValueError(f'Náhodnou variantu nelze bezpečně vytvořit: {e}')
            legacy_vals=_parse_variant_values(getattr(example,'variant_values','') or '')
            legacy_env={f'n{i+1}':v for i,v in enumerate(legacy_vals)}
            return {'problem':example.problem, 'prose_problem':_render_variant_prose(getattr(example, 'prose_problem', '') or '', legacy_env), 'steps':{st.id:{'instruction':st.instruction,'expected':st.expected,'hint':st.hint} for st in example.steps}, 'variant_fallback': True}

    problem = str(example.problem or '')
    compact = normalize_math_answer(problem).replace('**', '^')
    m = re.fullmatch(r'([+-]?\d*)\*?x([+-]\d+)=([+-]?\d+)', compact)
    if not m:
        return {'problem': problem, 'prose_problem':getattr(example, 'prose_problem', '') or '', 'steps': {st.id:{'instruction':st.instruction,'expected':st.expected,'hint':st.hint} for st in example.steps}}
    a_txt, b_txt, c_txt = m.groups()
    if a_txt in ('', '+'): a0 = 1
    elif a_txt == '-': a0 = -1
    else: a0 = int(a_txt)
    b0, c0 = int(b_txt), int(c_txt)
    if a0 == 0 or (c0 - b0) % a0 != 0:
        return {'problem': problem, 'prose_problem':getattr(example, 'prose_problem', '') or '', 'steps': {st.id:{'instruction':st.instruction,'expected':st.expected,'hint':st.hint} for st in example.steps}}
    x0 = (c0 - b0) // a0
    seed_src = f'{app.config.get("SECRET_KEY")}:{user_id}:{example.id}:math-variant'
    seed = int(hashlib.sha256(seed_src.encode()).hexdigest()[:16], 16)
    rng = random.Random(seed)
    a = rng.randint(2, 9) * (-1 if a0 < 0 else 1)
    x = rng.randint(2, 12)
    b_abs = rng.randint(1, 15)
    b = b_abs if b0 >= 0 else -b_abs
    c = a*x + b
    ax = a*x
    mapping = {a0:a, b0:b, c0:c, a0*x0:ax, x0:x}
    def transform(text_value):
        out = str(text_value or '')
        for old, new in sorted(mapping.items(), key=lambda kv: len(str(abs(kv[0]))), reverse=True):
            out = re.sub(rf'(?<![\d.]){re.escape(str(old))}(?![\d.])', str(new), out)
        return out
    sign = '+' if b >= 0 else '-'
    aa = '' if a == 1 else ('-' if a == -1 else str(a))
    new_problem = f'{aa}x {sign} {abs(b)} = {c}'
    steps = {st.id:{'instruction':transform(st.instruction),'expected':transform(st.expected),'hint':transform(st.hint)} for st in example.steps}
    return {'problem':new_problem, 'prose_problem':_replace_variant_numbers(getattr(example, 'prose_problem', '') or '', list(mapping.keys()), list(mapping.values())), 'steps':steps}

def validate_math_payload(payload):
    """Vrací None, nebo dict {message, field}; field dovolí UI označit jen chybný input."""
    for ei, ex in enumerate(payload, 1):
        if not (ex.get('steps') or []):
            return {'message': f'Příklad {ei} nemá žádné kroky.', 'field': f'examples.{ei-1}.steps'}
        problem_text=str(ex.get("problem", "") or '').strip()
        variable_mode=False
        vals=[]
        source_env={}
        helper_env={}
        if ex.get("variant_enabled"):
            try:
                variable_names=_parse_variant_variable_names(ex.get("variant_values", ""))
                variable_mode=bool(variable_names)
                if variable_mode:
                    # Čísla učitel nezadává. Pro kontrolu syntaxe použijeme pouze dočasné
                    # bezpečné hodnoty; pravdivost podmínky se ověřuje až generátorem.
                    source_env={name: float(i+2) for i,name in enumerate(variable_names)}
                else:
                    vals=_parse_variant_values(ex.get("variant_values", ""))
                    if not vals:
                        return {'message': f"Příklad {ei}: zadej vstupní proměnné, např. n1;n2;n3.", 'field': f'examples.{ei-1}.variant_values'}
                    source_env={f'n{i+1}':v for i,v in enumerate(vals)}
                _safe_variant_condition(ex.get("variant_condition", ""), source_env)
                helper_env=dict(source_env)
            except Exception as e:
                return {'message': f"Příklad {ei}: chyba ve vstupních proměnných nebo podmínce – {e}.", 'field': f'examples.{ei-1}.variant_condition'}
        # U slovní úlohy v novém proměnném režimu může být matematický vzor prázdný.
        if problem_text:
            ok, why = validate_math_expression(problem_text)
            if not ok:
                return {'message': f"Příklad {ei}: {why}. Oprav matematický zápis.", 'field': f'examples.{ei-1}.problem'}
        elif not variable_mode:
            return {'message': f"Příklad {ei}: matematický zápis je prázdný.", 'field': f'examples.{ei-1}.problem'}
            try:
                lo=float(ex.get("variant_min",1)); hi=float(ex.get("variant_max",30)); step=float(ex.get("variant_step",1))
                if step<=0 or hi<lo:
                    raise ValueError('minimum musí být menší nebo rovno maximu a krok musí být kladný')
            except Exception as e:
                return {'message': f"Příklad {ei}: zkontroluj rozsah a krok generování – {e}.", 'field': f'examples.{ei-1}.variant_range'}
            try:
                kind=str(ex.get('variant_result_kind') or 'any')
                sign=str(ex.get('variant_result_sign') or 'any')
                if kind not in {'any','integer'}:
                    raise ValueError('neplatný typ výsledku')
                if sign not in {'any','nonnegative','positive'}:
                    raise ValueError('neplatné znaménkové omezení')
                dec=int(ex.get('variant_result_decimals',-1) if str(ex.get('variant_result_decimals','')).strip() else -1)
                if dec < -1 or dec > 10:
                    raise ValueError('počet desetinných míst musí být -1 až 10')
                rmin=ex.get('variant_result_min'); rmax=ex.get('variant_result_max')
                if rmin not in (None,''): float(rmin)
                if rmax not in (None,''): float(rmax)
                if rmin not in (None,'') and rmax not in (None,'') and float(rmin)>float(rmax):
                    raise ValueError('minimum výsledku nesmí být větší než maximum')
            except Exception as e:
                return {'message': f"Příklad {ei}: zkontroluj pravidla pro výsledky – {e}.", 'field': f'examples.{ei-1}.variant_result_rule'}
        for si, st in enumerate(ex.get("steps") or [], 1):
            formula=str(st.get('result_formula') or '').strip()
            decimals=st.get('result_decimals',2)
            expected=str(st.get("expected", "") or '')
            if formula:
                if not ex.get('variant_enabled'):
                    return {'message': f"Příklad {ei}, krok {si}: vzorec pro přepočet výsledku vyžaduje zapnuté obecné náhodné varianty.", 'field': f'examples.{ei-1}.steps.{si-1}.result_formula'}
                try:
                    d=int(decimals)
                    if d < 0 or d > 10:
                        raise ValueError('zaokrouhlení musí být 0 až 10 desetinných míst')
                    # Ověříme všechny vzorce a že jejich vzorové výsledky umíme
                    # ve správné odpovědi najít. Vzorců může být více (středník / nový řádek).
                    formulas=_split_result_formulas(formula)
                    if not formulas:
                        raise ValueError('zadej alespoň jeden vzorec')
                    if variable_mode:
                        for f in formulas:
                            alias, expr = _formula_assignment(f)
                            value=_safe_variant_formula(expr, helper_env)
                            if alias:
                                helper_env[alias]=value
                    else:
                        _computed_result_markers(expected, source_env, source_env, formula, d)
                except Exception as e:
                    return {'message': f"Příklad {ei}, krok {si}: chyba ve vzorci výsledku – {e}.", 'field': f'examples.{ei-1}.steps.{si-1}.result_formula'}
            ok, why = validate_math_expression(expected)
            if not ok:
                return {'message': f"Příklad {ei}, krok {si}: {why}. Oprav správný stav.", 'field': f'examples.{ei-1}.steps.{si-1}.expected'}
    return None


def math_lesson_total_steps(lesson_id):
    return MathStep.query.join(MathExample).filter(MathExample.lesson_id == lesson_id).count()


def get_math_attempt(user_id, lesson_id, create=False):
    row = MathAttempt.query.filter_by(user_id=user_id, lesson_id=lesson_id).first()
    if not row and create:
        row = MathAttempt(user_id=user_id, lesson_id=lesson_id)
        db.session.add(row)
        db.session.flush()
    return row


def math_current_position(lesson, user_id):
    examples = sorted(lesson.examples, key=lambda x:x.order)
    flat = []
    for ex in examples:
        for st in sorted(ex.steps, key=lambda x:x.order):
            flat.append((ex, st))
    attempt = get_math_attempt(user_id, lesson.id, create=True)
    idx = min(int(attempt.completed_steps or 0), len(flat))
    return flat, attempt, idx


@app.route('/teacher/math/new', methods=['GET','POST'])
def new_math_lesson():
    r = require_teacher()
    if r: return r
    if request.method == 'POST':
        school=request.form.get('school','').strip()
        grade_name=request.form.get('grade_name','').strip()
        topic=request.form.get('topic','').strip()
        title=request.form.get('title','').strip()
        if not all((school,grade_name,topic,title)):
            lessons=MathLesson.query.order_by(MathLesson.created_at.desc()).all()
            missing = 'school' if not school else ('grade_name' if not grade_name else ('topic' if not topic else 'title'))
            return render_template('math_new.html', course=course_from_lesson(None), lesson=None, lessons=lessons,
                form_data={'school':school,'grade_name':grade_name,'topic':topic,'title':title},
                examples_json=request.form.get('examples_json','[]'),
                math_errors={'message':'Vyplň všechna povinná pole zařazení lekce.','field':missing})

        html_file=request.files.get('lesson_html')
        html_original=''; html_stored=''
        if html_file and html_file.filename:
            if Path(html_file.filename).suffix.lower() not in ('.html','.htm'):
                flash('Výklad musí být HTML soubor.')
                return redirect(url_for('new_math_lesson'))
            html_original=html_file.filename
            html_stored=_save_uploaded_file(html_file, INFORMATICS_SOURCE_DIR, 'math_html')

        lesson_row=MathLesson(
            school=school, grade_name=grade_name, topic=topic, title=title,
            html_original=html_original, html_stored=html_stored, is_published=True
        )
        db.session.add(lesson_row); db.session.flush()

        # Data přicházejí jako JSON z dynamického editoru.
        try:
            payload=json.loads(request.form.get('examples_json','[]') or '[]')
        except Exception:
            payload=[]
        if not payload:
            db.session.rollback()
            lessons=MathLesson.query.order_by(MathLesson.created_at.desc()).all()
            return render_template('math_new.html', course=course_from_lesson(None), lesson=None, lessons=lessons,
                form_data={'school':school,'grade_name':grade_name,'topic':topic,'title':title},
                examples_json=request.form.get('examples_json','[]'),
                math_errors={'message':'Přidej alespoň jeden příklad.','field':'examples'})
        math_error = validate_math_payload(payload)
        if math_error:
            db.session.rollback()
            lessons=MathLesson.query.order_by(MathLesson.created_at.desc()).all()
            return render_template('math_new.html', course=course_from_lesson(None), lesson=None, lessons=lessons,
                form_data={'school':school,'grade_name':grade_name,'topic':topic,'title':title},
                examples_json=json.dumps(payload, ensure_ascii=False), math_errors=math_error)

        previous_steps=[]
        for ei, ex in enumerate(payload,1):
            image_file = request.files.get(f'example_image_{ei-1}')
            image_stored = save_math_example_image(image_file)
            if image_file and image_file.filename and not image_stored:
                db.session.rollback()
                flash(f'Příklad {ei}: obrázek musí být PNG, JPG, JPEG, WEBP, GIF nebo SVG.')
                return redirect(url_for('new_math_lesson'))
            ex_row=MathExample(
                lesson_id=lesson_row.id, order=ei,
                title=str(ex.get('title') or f'Příklad {ei}'),
                problem=(str(ex.get('problem') or '').strip() or (str(ex.get('variant_values') or '').strip() if _parse_variant_variable_names(ex.get('variant_values','')) else '')),
                prose_problem=str(ex.get('prose_problem') or '').strip(),
                image_stored=image_stored,
                variant_enabled=bool(ex.get('variant_enabled')),
                variant_values=str(ex.get('variant_values') or '').strip(),
                variant_condition=str(ex.get('variant_condition') or '').strip(),
                variant_min=float(ex.get('variant_min') or 1),
                variant_max=float(ex.get('variant_max') or 30),
                variant_step=float(ex.get('variant_step') or 1),
                variant_result_kind=str(ex.get('variant_result_kind') or 'any'),
                variant_result_sign=str(ex.get('variant_result_sign') or 'any'),
                variant_result_min=(float(ex.get('variant_result_min')) if ex.get('variant_result_min') not in (None,'') else None),
                variant_result_max=(float(ex.get('variant_result_max')) if ex.get('variant_result_max') not in (None,'') else None),
                variant_result_decimals=int(ex.get('variant_result_decimals',-1) if str(ex.get('variant_result_decimals','')).strip() else -1)
            )
            db.session.add(ex_row); db.session.flush()
            steps=ex.get('steps') or []
            if ex.get('copy_previous') and previous_steps and not steps:
                steps=[dict(x) for x in previous_steps]
            if not steps:
                db.session.rollback()
                flash(f'Příklad {ei} nemá žádné kroky.')
                return redirect(url_for('new_math_lesson'))
            previous_steps=[]
            for si, st in enumerate(steps,1):
                item={
                    'instruction':str(st.get('instruction') or '').strip(),
                    'expected':str(st.get('expected') or '').strip(),
                    'hint':str(st.get('hint') or '').strip(),
                    'result_formula':str(st.get('result_formula') or '').strip(),
                    'result_decimals':int(st.get('result_decimals') if str(st.get('result_decimals','')).strip() else 2)
                }
                previous_steps.append(item)
                db.session.add(MathStep(
                    example_id=ex_row.id, order=si,
                    instruction=item['instruction'], expected=item['expected'], hint=item['hint'],
                    result_formula=item['result_formula'], result_decimals=item['result_decimals']
                ))
        db.session.commit()
        flash('Matematická lekce byla vytvořena.')
        return redirect(url_for('math_lesson', lesson_id=lesson_row.id))

    lessons=MathLesson.query.order_by(MathLesson.created_at.desc()).all()
    return render_template('math_new.html', course=course_from_lesson(None), lesson=None, lessons=lessons, form_data={}, examples_json='', math_errors=None)


@app.route('/math-lesson/<int:lesson_id>', methods=['GET','POST'])
def math_lesson(lesson_id):
    r=require_login()
    if r:return r
    item=db.session.get(MathLesson,lesson_id)
    if not item:return 'Matematická lekce nebyla nalezena.',404

    if current_user().role=='student':
        begin_focus_attempt('math', item.id)

    examples=sorted(item.examples,key=lambda x:x.order)
    math_variants = {ex.id: generated_math_variant(ex, current_user().id) for ex in examples} if current_user().role=='student' else {}
    html_content=render_informatics_html(item.html_stored)
    current=None; attempt=None; idx=0; total=0; message=None; hint=None
    if current_user().role=='student':
        flat,attempt,idx=math_current_position(item,current_user().id)
        total=len(flat)
        if idx < total:
            current=flat[idx]
        if request.method=='POST':
            action=request.form.get('action','check')
            if action=='submit':
                percent=round((attempt.completed_steps or 0)/max(total,1)*100)
                attempt.total_steps=total; attempt.percent=percent
                attempt.grade=informatics_grade_from_percent(percent)
                attempt.status='odevzdáno'
                attempt.focus_lost=consume_focus_count('math',item.id)
                attempt.updated_at=datetime.utcnow()
                db.session.commit()
                flash(f'Lekce byla odevzdána: {percent} %, známka {attempt.grade}.')
                return redirect(url_for('subject_catalog',kind='matematika'))

            if current:
                ex,step=current
                expected_now = math_variants.get(ex.id, {}).get('steps', {}).get(step.id, {}).get('expected', step.expected)
                answer = math_answer_from_fields(expected_now, request.form)
                if math_answers_equivalent(answer, expected_now):
                    answers = _safe_json(attempt.answers_json, [])
                    answers.append({
                        'example_id': ex.id,
                        'step_id': step.id,
                        'answer': str(answer).strip()
                    })
                    attempt.answers_json = json.dumps(answers, ensure_ascii=False)
                    attempt.completed_steps=int(attempt.completed_steps or 0)+1
                    attempt.total_steps=total
                    attempt.percent=round(attempt.completed_steps/max(total,1)*100)
                    attempt.grade=informatics_grade_from_percent(attempt.percent)
                    attempt.status='rozpracováno'
                    attempt.focus_lost=get_focus_count('math',item.id)
                    attempt.updated_at=datetime.utcnow()
                    db.session.commit()
                    message='✅ Správně. Odemkl se další krok.'
                    flat,attempt,idx=math_current_position(item,current_user().id)
                    current=flat[idx] if idx < len(flat) else None
                else:
                    message='❌ Tento krok ještě není správně.'
                    hint=step.hint
        percent=round((attempt.completed_steps or 0)/max(total,1)*100) if attempt else 0
        grade=informatics_grade_from_percent(percent)
    else:
        percent=0;grade=None

    history_by_example = {}
    if attempt:
        for saved in _safe_json(attempt.answers_json, []):
            history_by_example.setdefault(str(saved.get('example_id')), []).append(saved.get('answer',''))

    current_example_number = current[0].order if current else (len(examples) if examples else 0)
    current_step_number = current[1].order if current else None

    current_input_html = None
    if current and current_user().role == 'student':
        ex_now, st_now = current
        expected_now = math_variants.get(ex_now.id, {}).get('steps', {}).get(st_now.id, {}).get('expected', st_now.expected)
        current_input_html = render_math_input_layout(math_input_layout(expected_now))

    course={'subject':'Matematika','grade':item.grade_name,'block':item.topic,'icon':'➗'}
    lesson_obj={'title':item.title,'block':item.topic}
    return render_template('math_lesson_engine.html',course=course,lesson=lesson_obj,item=item,
        examples=examples,current=current,attempt=attempt,total=total,idx=idx,
        percent=percent,grade=grade,message=message,hint=hint,html_content=html_content,
        history_by_example=history_by_example,
        current_example_number=current_example_number,
        current_step_number=current_step_number,
        math_variants=math_variants,current_input_html=current_input_html)




@app.route('/teacher/math/<int:lesson_id>/edit', methods=['GET','POST'])
def edit_math_lesson(lesson_id):
    r = require_teacher()
    if r: return r

    item = db.session.get(MathLesson, lesson_id)
    if not item:
        return 'Matematická lekce nebyla nalezena.', 404

    if request.method == 'POST':
        item.school = request.form.get('school', item.school).strip()
        item.grade_name = request.form.get('grade_name', item.grade_name).strip()
        item.topic = request.form.get('topic', item.topic).strip()
        item.title = request.form.get('title', item.title).strip()

        html_file = request.files.get('lesson_html')
        if html_file and html_file.filename:
            if Path(html_file.filename).suffix.lower() not in ('.html','.htm'):
                flash('Výklad musí být HTML soubor.')
                return redirect(url_for('edit_math_lesson', lesson_id=item.id))
            item.html_original = html_file.filename
            item.html_stored = _save_uploaded_file(html_file, INFORMATICS_SOURCE_DIR, 'math_html')

        try:
            payload = json.loads(request.form.get('examples_json','[]') or '[]')
        except Exception:
            payload = []

        if not payload:
            db.session.rollback()
            return render_template('math_edit.html', course=course_from_lesson(None), lesson=None, item=item,
                form_data={'school':request.form.get('school',''),'grade_name':request.form.get('grade_name',''),'topic':request.form.get('topic',''),'title':request.form.get('title','')},
                examples_json='[]', math_errors={'message':'Lekce musí obsahovat alespoň jeden příklad.','field':'examples'})
        math_error = validate_math_payload(payload)
        if math_error:
            db.session.rollback()
            return render_template('math_edit.html', course=course_from_lesson(None), lesson=None, item=item,
                form_data={'school':request.form.get('school',''),'grade_name':request.form.get('grade_name',''),'topic':request.form.get('topic',''),'title':request.form.get('title','')},
                examples_json=json.dumps(payload, ensure_ascii=False), math_errors=math_error)

        old_example_ids = [e.id for e in item.examples]
        if old_example_ids:
            MathStep.query.filter(MathStep.example_id.in_(old_example_ids)).delete(synchronize_session=False)
            MathExample.query.filter(MathExample.id.in_(old_example_ids)).delete(synchronize_session=False)
        db.session.flush()

        previous_steps = []
        for ei, ex in enumerate(payload, 1):
            existing_image = str(ex.get('image_stored') or '').strip()
            remove_image = bool(ex.get('remove_image'))
            image_file = request.files.get(f'example_image_{ei-1}')
            new_image = save_math_example_image(image_file)
            if image_file and image_file.filename and not new_image:
                db.session.rollback()
                flash(f'Příklad {ei}: obrázek musí být PNG, JPG, JPEG, WEBP, GIF nebo SVG.')
                return redirect(url_for('edit_math_lesson', lesson_id=item.id))
            image_stored = '' if remove_image else (new_image or existing_image)
            ex_row = MathExample(
                lesson_id=item.id,
                order=ei,
                title=str(ex.get('title') or f'Příklad {ei}'),
                problem=(str(ex.get('problem') or '').strip() or (str(ex.get('variant_values') or '').strip() if _parse_variant_variable_names(ex.get('variant_values','')) else '')),
                prose_problem=str(ex.get('prose_problem') or '').strip(),
                image_stored=image_stored,
                variant_enabled=bool(ex.get('variant_enabled')),
                variant_values=str(ex.get('variant_values') or '').strip(),
                variant_condition=str(ex.get('variant_condition') or '').strip(),
                variant_min=float(ex.get('variant_min') or 1),
                variant_max=float(ex.get('variant_max') or 30),
                variant_step=float(ex.get('variant_step') or 1),
                variant_result_kind=str(ex.get('variant_result_kind') or 'any'),
                variant_result_sign=str(ex.get('variant_result_sign') or 'any'),
                variant_result_min=(float(ex.get('variant_result_min')) if ex.get('variant_result_min') not in (None,'') else None),
                variant_result_max=(float(ex.get('variant_result_max')) if ex.get('variant_result_max') not in (None,'') else None),
                variant_result_decimals=int(ex.get('variant_result_decimals',-1) if str(ex.get('variant_result_decimals','')).strip() else -1)
            )
            db.session.add(ex_row)
            db.session.flush()

            steps = ex.get('steps') or []
            if ex.get('copy_previous') and previous_steps and not steps:
                steps = [dict(x) for x in previous_steps]

            if not steps:
                db.session.rollback()
                flash(f'Příklad {ei} nemá žádné kroky.')
                return redirect(url_for('edit_math_lesson', lesson_id=item.id))

            previous_steps = []
            for si, st in enumerate(steps, 1):
                step_data = {
                    'instruction': str(st.get('instruction') or '').strip(),
                    'expected': str(st.get('expected') or '').strip(),
                    'hint': str(st.get('hint') or '').strip(),
                    'result_formula': str(st.get('result_formula') or '').strip(),
                    'result_decimals': int(st.get('result_decimals') if str(st.get('result_decimals','')).strip() else 2),
                }
                previous_steps.append(step_data)
                db.session.add(MathStep(
                    example_id=ex_row.id,
                    order=si,
                    instruction=step_data['instruction'],
                    expected=step_data['expected'],
                    hint=step_data['hint'],
                    result_formula=step_data['result_formula'],
                    result_decimals=step_data['result_decimals'],
                ))

        # Při změně struktury lekce vynulujeme staré rozpracované pokusy,
        # aby jejich počet kroků neodkazoval na starou strukturu.
        MathAttempt.query.filter_by(lesson_id=item.id).delete()
        db.session.commit()
        flash('Matematická lekce byla upravena.')
        return redirect(url_for('math_lesson', lesson_id=item.id))

    examples_data = []
    for ex in sorted(item.examples, key=lambda x:x.order):
        examples_data.append({
            'title': ex.title,
            'problem': ('' if _is_variable_declaration(ex.problem, _parse_variant_variable_names(ex.variant_values or '')) else ex.problem),
            'prose_problem': ex.prose_problem or '',
            'image_stored': ex.image_stored or '',
            'remove_image': False,
            'variant_enabled': bool(ex.variant_enabled),
            'variant_values': ex.variant_values or '',
            'variant_condition': ex.variant_condition or '',
            'variant_min': ex.variant_min if ex.variant_min is not None else 1,
            'variant_max': ex.variant_max if ex.variant_max is not None else 30,
            'variant_step': ex.variant_step if ex.variant_step is not None else 1,
            'variant_result_kind': getattr(ex, 'variant_result_kind', 'any') or 'any',
            'variant_result_sign': getattr(ex, 'variant_result_sign', 'any') or 'any',
            'variant_result_min': getattr(ex, 'variant_result_min', None),
            'variant_result_max': getattr(ex, 'variant_result_max', None),
            'variant_result_decimals': getattr(ex, 'variant_result_decimals', -1) if getattr(ex, 'variant_result_decimals', None) is not None else -1,
            'steps': [
                {
                    'instruction': st.instruction,
                    'expected': st.expected,
                    'hint': st.hint,
                    'result_formula': getattr(st, 'result_formula', '') or '',
                    'result_decimals': getattr(st, 'result_decimals', 2) if getattr(st, 'result_decimals', None) is not None else 2,
                }
                for st in sorted(ex.steps, key=lambda x:x.order)
            ]
        })

    return render_template(
        'math_edit.html',
        course=course_from_lesson(None),
        lesson=None,
        item=item,
        examples_json=json.dumps(examples_data, ensure_ascii=False),
        form_data={}, math_errors=None
    )

@app.route('/teacher/math/<int:lesson_id>/delete', methods=['POST'])
def delete_math_lesson(lesson_id):
    r = require_teacher()
    if r:
        return r

    item = db.session.get(MathLesson, lesson_id)
    if not item:
        flash('Matematická lekce už neexistuje.')
        return redirect(url_for('teacher_home'))

    try:
        # Nejdřív smažeme výsledky/pokusy studentů.
        MathAttempt.query.filter_by(lesson_id=item.id).delete(synchronize_session=False)

        # Kroky a příklady mažeme jako načtené ORM objekty.
        # Tím zabráníme StaleDataError při následném smazání lekce.
        examples = list(MathExample.query.filter_by(lesson_id=item.id).order_by(MathExample.order).all())

        for example in examples:
            steps = list(MathStep.query.filter_by(example_id=example.id).all())
            for step in steps:
                db.session.delete(step)
            db.session.flush()
            db.session.delete(example)

        db.session.flush()
        db.session.delete(item)
        db.session.commit()

        flash('Matematická lekce byla smazána.')
    except Exception:
        db.session.rollback()
        raise

    return redirect(url_for('teacher_home'))


@app.route('/teacher')
def teacher_home():
    r=require_teacher();
    if r: return r
    subjects = Subject.query.order_by(Subject.name).all()
    students = User.query.filter_by(role='student').order_by(User.name).all()
    interactive_lessons = InteractiveLesson.query.order_by(
        InteractiveLesson.subject,
        InteractiveLesson.school,
        InteractiveLesson.grade_name,
        InteractiveLesson.topic,
        InteractiveLesson.title
    ).all()
    informatics_lessons = InformaticsLesson.query.order_by(
        InformaticsLesson.school,
        InformaticsLesson.grade_name,
        InformaticsLesson.topic,
        InformaticsLesson.title
    ).all()
    math_lessons = MathLesson.query.order_by(
        MathLesson.school, MathLesson.grade_name, MathLesson.topic, MathLesson.title
    ).all()
    return render_template(
        'teacher.html',
        course=course_from_lesson(None),
        subjects=subjects,
        students=students,
        student_rows=student_overview_rows(),
        interactive_lessons=interactive_lessons,
        informatics_lessons=informatics_lessons,
        math_lessons=math_lessons
    )

@app.route('/teacher/students', methods=['GET','POST'])
def teacher_students():
    r=require_teacher();
    if r: return r
    if request.method == 'POST':
        username = request.form.get('username','').strip().lower()
        password = request.form.get('password','').strip()
        username = strip_accents(username).replace(' ', '.')
        username = re.sub(r'[^a-z0-9._-]+', '', username)
        username = re.sub(r'\.+', '.', username).strip('.')
        name = ' '.join(part.capitalize() for part in username.split('.') if part)
        if not username or '.' not in username or not password:
            flash('Zadej studenta ve tvaru jmeno.prijmeni a vyplň heslo.')
        elif User.query.filter_by(username=username).first():
            flash('Toto uživatelské jméno už existuje.')
        else:
            db.session.add(User(username=username, name=name or username, role='student', password_hash=generate_password_hash(password)))
            db.session.commit()
            flash('Student byl vytvořen. Může se přihlásit vlastním jménem a heslem.')
            return redirect(url_for('teacher_students'))
    students = User.query.filter_by(role='student').order_by(User.name).all()
    return render_template('students.html', course=course_from_lesson(None), students=students, student_rows=student_overview_rows())



def _answer_text(q):
    """Čitelná správná odpověď pro učitelský klíč."""
    try:
        opts = json.loads(q.options_json or '[]')
    except Exception:
        opts = []
    if q.qtype in ('choice', 'image_choice'):
        try:
            idx = int(json.loads(q.correct_json or '0'))
        except Exception:
            idx = 0
        if isinstance(opts, list) and 0 <= idx < len(opts):
            val = opts[idx]
            return str(val) if val else f'možnost {idx + 1}'
        return f'možnost {idx + 1}'
    try:
        roots = json.loads(q.roots_json or '[]')
    except Exception:
        roots = []
    if roots:
        return ' / '.join(str(x) for x in roots)
    return '—'


def _activity_solution(a, english=False):
    raw = (a.config_en_json or '') if english else (a.config_json or '')
    if english and not raw.strip():
        raw = a.config_json or '{}'
    try:
        cfg = json.loads(raw or '{}')
    except Exception:
        cfg = {}
    typ = a.activity_type or ''
    if typ == 'video_observe':
        ops = cfg.get('options') or []
        idx = int(cfg.get('correct', 0) or 0)
        return str(ops[idx]) if 0 <= idx < len(ops) else f"{'option' if english else 'možnost'} {idx+1}"
    if typ in ('find_image', 'video_find'):
        return 'Correct marked area in the image/video' if english else 'Správně označená oblast v obrázku/videu'
    if typ == 'real_world':
        concepts = cfg.get('concepts') or []
        vals = []
        for group in concepts:
            if isinstance(group, list) and group:
                vals.append(str(group[0]))
            elif group:
                vals.append(str(group))
        minimum = cfg.get('min_items', '')
        prefix = f"At least {minimum}: " if english and minimum else (f"Alespoň {minimum}: " if minimum else '')
        return prefix + ', '.join(vals)
    if typ == 'cards':
        # Kartičky se vyhodnocují podle uloženého pořadí/přiřazení.
        return 'Correct card matching/order as defined by the teacher' if english else 'Správné přiřazení/pořadí kartiček podle zadání učitele'
    if typ == 'sort':
        cats = cfg.get('categories') or []
        items = cfg.get('items') or []
        out=[]
        for item in items:
            if not isinstance(item, dict):
                continue
            ci = int(item.get('category', 0) or 0)
            cat = cats[ci] if 0 <= ci < len(cats) else str(ci+1)
            out.append(f"{item.get('label','')} → {cat}")
        return '; '.join(out) or '—'
    return 'Correct completion of the activity' if english else 'Správné splnění aktivity'


def _lesson_solution_text(lesson, bilingual=False):
    lines = []
    subject = lesson.block.grade.subject.name
    lines += ['UČITELSKÉ ŘEŠENÍ / TEACHER ANSWER KEY' if bilingual else 'UČITELSKÉ ŘEŠENÍ',
              '=' * 58,
              f"Předmět: {subject}", f"Ročník: {lesson.block.grade.name}", f"Téma: {lesson.block.title}",
              f"Lekce: {lesson.title}" + (f" / {lesson.title_en}" if bilingual and lesson.title_en else ''), '']
    cs = sorted([q for q in lesson.questions if (q.lang or 'cs') == 'cs'], key=lambda q:(q.area, q.order, q.id))
    en = sorted([q for q in lesson.questions if (q.lang or 'cs') == 'en'], key=lambda q:(q.area, q.order, q.id))
    en_map = {(q.area, q.order): q for q in en}
    lines.append('OTÁZKY / QUESTIONS' if bilingual else 'OTÁZKY')
    lines.append('-' * 58)
    for i,q in enumerate(cs,1):
        lines.append(f"{i}. {q.question}")
        lines.append(f"   Správná odpověď: {_answer_text(q)}")
        if bilingual:
            qe = en_map.get((q.area, q.order))
            if qe:
                lines.append(f"   EN: {qe.question}")
                lines.append(f"   Correct answer: {_answer_text(qe)}")
        lines.append('')
    acts = sorted([a for a in lesson.practical_activities if (a.lang or 'cs') != 'en'], key=lambda a:(a.order,a.id))
    if acts:
        lines += ['PRAKTICKÉ AKTIVITY / PRACTICAL ACTIVITIES' if bilingual else 'PRAKTICKÉ AKTIVITY', '-' * 58]
        for i,a in enumerate(acts,1):
            lines.append(f"{i}. {a.title or 'Praktická aktivita'}")
            if a.prompt: lines.append(f"   Zadání: {a.prompt}")
            lines.append(f"   Řešení: {_activity_solution(a, False)}")
            if bilingual:
                lines.append(f"   EN: {a.title_en or a.title or 'Practical activity'}")
                if a.prompt_en: lines.append(f"   Task: {a.prompt_en}")
                lines.append(f"   Solution: {_activity_solution(a, True)}")
            lines.append('')
    lines += ['=' * 58, 'Soubor byl automaticky vytvořen Digitální učebnicí.']
    return '\n'.join(lines)


@app.route('/teacher/lesson/<int:lesson_id>/solutions/<mode>')
def download_lesson_solutions(lesson_id, mode):
    r = require_teacher()
    if r: return r
    lesson = db.session.get(Lesson, lesson_id)
    if not lesson:
        return 'Lekce nebyla nalezena.', 404
    bilingual = mode == 'bilingual'
    if mode not in ('cs', 'bilingual'):
        return 'Neplatný režim.', 400
    text_data = _lesson_solution_text(lesson, bilingual=bilingual)
    slug = re.sub(r'[^a-zA-Z0-9_-]+', '_', strip_accents(lesson.title or 'lekce')).strip('_') or 'lekce'
    suffix = 'CZ_EN' if bilingual else 'CZ'
    # UTF-8 BOM zajistí správnou češtinu i při otevření v běžných editorech ve Windows.
    data = io.BytesIO(('\ufeff' + text_data).encode('utf-8'))
    data.seek(0)
    return send_file(data, mimetype='text/plain; charset=utf-8', as_attachment=True,
                     download_name=f'RESENI_{slug}_{suffix}.txt')

@app.route('/teacher/database')
def teacher_database():
    r=require_teacher()
    if r: return r
    html_results = Result.query.order_by(Result.created_at.desc()).all()
    interactive_results = InteractiveResult.query.order_by(InteractiveResult.completed_at.desc()).all()
    informatics_results = InformaticsSubmission.query.filter(
        InformaticsSubmission.status != 'kontrola'
    ).order_by(InformaticsSubmission.created_at.desc()).all()
    math_results = MathAttempt.query.filter(
        MathAttempt.status != 'rozpracováno'
    ).order_by(MathAttempt.updated_at.desc()).all()
    return render_template(
        'database.html',
        course=course_from_lesson(None),
        student_rows=student_overview_rows(),
        html_results=html_results,
        interactive_results=interactive_results,
        informatics_results=informatics_results,
        math_results=math_results
    )


@app.route('/teacher/result/<int:result_id>/delete', methods=['POST'])
def delete_result(result_id):
    r=require_teacher()
    if r: return r
    result = db.session.get(Result, result_id)
    if result:
        user_id = result.user_id
        lesson_id = result.lesson_id
        reset_student_lesson_progress(user_id, lesson_id)
        flash('Postup studenta v této lekci byl kompletně resetován. Při příštím otevření začne od začátku.')
    return redirect(url_for('teacher_database'))


@app.route('/teacher/interactive-result/<int:result_id>/delete', methods=['POST'])
def delete_interactive_result(result_id):
    r = require_teacher()
    if r:
        return r
    result = db.session.get(InteractiveResult, result_id)
    if result:
        db.session.delete(result)
        db.session.commit()
        flash('Výsledek interaktivní lekce byl smazán. Pokrok studenta zůstal zachovaný.')
    return redirect(url_for('teacher_database'))


@app.route('/teacher/informatics-result/<int:result_id>/delete', methods=['POST'])
def delete_informatics_result(result_id):
    r = require_teacher()
    if r:
        return r
    result = db.session.get(InformaticsSubmission, result_id)
    if result:
        db.session.delete(result)
        db.session.commit()
        flash('Výsledek informatického úkolu byl smazán.')
    return redirect(url_for('teacher_database'))


@app.route('/teacher/results/delete-all', methods=['POST'])
def delete_all_results():
    r=require_teacher()
    if r: return r
    Result.query.delete()
    InteractiveResult.query.delete()
    InformaticsSubmission.query.delete()
    MathAttempt.query.delete()
    db.session.commit()
    flash('Všechny výsledky byly smazány. Studenti a jejich průběžný pokrok ve všech předmětech zůstali zachováni.')
    return redirect(url_for('teacher_database'))

@app.route('/teacher/student/<int:user_id>/delete', methods=['POST'])
def delete_student(user_id):
    r=require_teacher()
    if r: return r
    stu = db.session.get(User, user_id)
    if stu and stu.role == 'student':
        Result.query.filter_by(user_id=stu.id).delete()
        StudentProgress.query.filter_by(user_id=stu.id).delete()
        StudentSectionProgress.query.filter_by(user_id=stu.id).delete()
        StudyQuestionProgress.query.filter_by(user_id=stu.id).delete()
        StudentActivityProgress.query.filter_by(user_id=stu.id).delete()
        FinalItemProgress.query.filter_by(user_id=stu.id).delete()
        StudentLessonReset.query.filter_by(user_id=stu.id).delete()
        InteractiveResult.query.filter_by(user_id=stu.id).delete()
        InteractiveProgress.query.filter_by(user_id=stu.id).delete()
        InformaticsSubmission.query.filter_by(user_id=stu.id).delete()
        MathAttempt.query.filter_by(user_id=stu.id).delete()
        db.session.delete(stu)
        db.session.commit()
        flash('Student byl smazán včetně jeho uloženého postupu a výsledků.')
    return redirect(url_for('teacher_students'))

@app.route('/teacher/lesson/<int:lesson_id>/archive', methods=['POST'])
def archive_lesson(lesson_id):
    r = require_teacher()
    if r: return r
    les = db.session.get(Lesson, lesson_id)
    if not les:
        flash('Lekce nebyla nalezena.')
        return redirect(url_for('teacher_home'))
    les.is_published = False
    db.session.commit()
    flash(f'Lekce „{les.title}“ byla archivována. Studentům se už nebude zobrazovat, ale výsledky zůstaly uložené.')
    return redirect(url_for('teacher_home'))

@app.route('/teacher/lesson/<int:lesson_id>/restore', methods=['POST'])
def restore_lesson(lesson_id):
    r = require_teacher()
    if r: return r
    les = db.session.get(Lesson, lesson_id)
    if not les:
        flash('Lekce nebyla nalezena.')
        return redirect(url_for('teacher_home'))
    les.is_published = True
    db.session.commit()
    flash(f'Lekce „{les.title}“ byla obnovena a studenti ji znovu uvidí.')
    return redirect(url_for('teacher_home'))

@app.route('/teacher/lesson/<int:lesson_id>/delete', methods=['POST'])
def delete_lesson(lesson_id):
    r = require_teacher()
    if r: return r
    les = db.session.get(Lesson, lesson_id)
    if not les:
        flash('Lekce nebyla nalezena.')
        return redirect(url_for('teacher_home'))
    title = les.title
    # Trvalé mazání: odstraníme výsledky, otázky, obrázky ve výkladu a sekce navázané na lekci.
    Result.query.filter_by(lesson_id=les.id).delete()
    StudentProgress.query.filter_by(lesson_id=les.id).delete()
    StudentSectionProgress.query.filter_by(lesson_id=les.id).delete()
    StudyQuestionProgress.query.filter_by(lesson_id=les.id).delete()
    FinalItemProgress.query.filter_by(lesson_id=les.id).delete()
    StudentLessonReset.query.filter_by(lesson_id=les.id).delete()
    activity_ids = [a.id for a in les.practical_activities]
    if activity_ids:
        StudentActivityProgress.query.filter(StudentActivityProgress.activity_id.in_(activity_ids)).delete(synchronize_session=False)
    PracticalActivity.query.filter_by(lesson_id=les.id).delete(synchronize_session=False)
    Question.query.filter_by(lesson_id=les.id).delete()
    section_ids = [sec.id for sec in les.sections]
    if section_ids:
        InlineImage.query.filter(InlineImage.section_id.in_(section_ids)).delete(synchronize_session=False)
        Section.query.filter(Section.id.in_(section_ids)).delete(synchronize_session=False)
    # Lekci smažeme bulk DELETE. Po předchozím bulk smazání sekcí by
    # db.session.delete(les) mohl přes načtený vztah les.sections zkoušet
    # nastavovat section.lesson_id = NULL, přestože je sloupec NOT NULL.
    Lesson.query.filter_by(id=les.id).delete(synchronize_session=False)
    db.session.flush()
    cleanup_empty_curriculum()
    db.session.commit()
    flash(f'Lekce „{title}“ byla trvale smazána. Prázdné staré odkazy byly také odstraněny.')
    return redirect(url_for('teacher_home'))



def import_docx_to_html(file_storage):
    """Starší nepoužívaná funkce pro DOCX import. Nově používáme CKEditor."""
    import mammoth
    if not file_storage or not file_storage.filename:
        return '', None
    if not file_storage.filename.lower().endswith('.docx'):
        raise ValueError('Podporovaný je pouze formát .docx')

    original_name = secure_filename(file_storage.filename)

    def convert_image(image):
        ext = 'png'
        if image.content_type and '/' in image.content_type:
            ext = image.content_type.split('/')[-1].replace('jpeg', 'jpg')
        img_name = datetime.now().strftime('%Y%m%d%H%M%S%f_') + 'docx_image.' + ext
        target = UPLOADS / img_name
        with image.open() as image_bytes:
            target.write_bytes(image_bytes.read())
        return {'src': url_for('uploads', filename=img_name)}

    try:
        file_storage.stream.seek(0)
        style_map = """
        p[style-name='Title'] => h1:fresh
        p[style-name='Subtitle'] => h2:fresh
        p[style-name='Heading 1'] => h2:fresh
        p[style-name='Heading 2'] => h3:fresh
        p[style-name='Heading 3'] => h4:fresh
        table => table.docx-table
        """
        result = mammoth.convert_to_html(
            file_storage.stream,
            convert_image=mammoth.images.img_element(convert_image),
            style_map=style_map
        )
    except Exception as exc:
        raise ValueError('DOCX se nepodařilo načíst. Zkontroluj, že jde opravdu o soubor .docx uložený z Wordu nebo LibreOffice.') from exc

    html_value = (result.value or '').strip()
    if not html_value:
        raise ValueError('DOCX se načetl, ale neobsahuje žádný převoditelný výklad.')

    messages = ''.join(f'<li>{html.escape(str(m))}</li>' for m in (result.messages or []))
    info = f'<div class="docx-import-note"><b>Importováno z DOCX:</b> {html.escape(original_name)}</div>'
    if messages:
        info += f'<details class="docx-import-warnings"><summary>Upozornění z převodu</summary><ul>{messages}</ul></details>'
    return info + '<div class="imported-docx-content">' + html_value + '</div>', original_name


def practical_activities_editor_json(lesson, lang='cs'):
    arr = []
    for a in sorted([x for x in lesson.practical_activities if (x.lang or 'cs') != 'en'], key=lambda x: x.order) if lesson else []:
        try: cfg = json.loads(a.config_json or '{}')
        except Exception: cfg = {}
        try: cfg_en = json.loads(a.config_en_json or '{}')
        except Exception: cfg_en = {}
        arr.append({
            'id': a.id, 'type': a.activity_type, 'title': a.title, 'title_en': a.title_en or '',
            'prompt': a.prompt, 'prompt_en': a.prompt_en or '', 'config': cfg, 'config_en': cfg_en,
            'image': a.image_file or '', 'video': a.video_file or '', 'include_final': bool(a.include_final)
        })
    return json.dumps(arr, ensure_ascii=False)


def _save_activity_file(ref, prefix):
    ref = str(ref or '').strip()
    if ref.startswith('__file__:'):
        field = ref.split(':', 1)[1]
        f = request.files.get(field)
        if f and f.filename:
            return save_upload(f)
        return ''
    return ref


def save_practical_activities_from_payload(lesson, section, payload, lang='cs'):
    try:
        data = json.loads(payload or '[]')
    except Exception:
        data = []

    # Média praktických aktivit mohou být v payloadu odkazována vícekrát
    # (např. stejný obrázek v české i anglické konfiguraci). FileStorage ale
    # nelze bezpečně ukládat opakovaně, protože po prvním save() je ukazatel
    # na konci souboru. Proto každý upload uložíme právě jednou a dále už
    # používáme stejné uložené jméno. Tím se zároveň opravuje občasné 404 u
    # aktivit „Najdi na obrázku“ a u obrázkových kartiček.
    activity_upload_cache = {}

    def _resolve_activity_media(ref):
        ref = str(ref or '').strip()
        if not ref.startswith('__file__:'):
            return ref
        field = ref.split(':', 1)[1]
        if field in activity_upload_cache:
            return activity_upload_cache[field]
        f = request.files.get(field)
        saved = ''
        if f and f.filename:
            try:
                f.stream.seek(0)
            except Exception:
                pass
            saved = save_upload(f)
        activity_upload_cache[field] = saved
        return saved

    PracticalActivity.query.filter_by(lesson_id=lesson.id).delete(synchronize_session=False)
    order = 1
    for item in data if isinstance(data, list) else []:
        typ = str(item.get('type') or '').strip()
        if typ not in {'cards','video_find','video_observe','find_image','real_world'}:
            continue
        title = str(item.get('title') or 'Praktická aktivita').strip()
        prompt = str(item.get('prompt') or '').strip()
        title_en = str(item.get('title_en') or '').strip()
        prompt_en = str(item.get('prompt_en') or '').strip()
        cfg = item.get('config') if isinstance(item.get('config'), dict) else {}
        cfg_en = item.get('config_en') if isinstance(item.get('config_en'), dict) else {}
        # Praktické aktivity mohou mít média i uvnitř configu (např. více obrázků
        # u obrázkových kartiček). Projdeme config rekurzivně a upload reference uložíme.
        def _save_cfg_media(value):
            if isinstance(value, dict):
                return {k: _save_cfg_media(v) for k, v in value.items()}
            if isinstance(value, list):
                return [_save_cfg_media(v) for v in value]
            if isinstance(value, str) and value.startswith('__file__:'):
                return _resolve_activity_media(value)
            return value
        cfg = _save_cfg_media(cfg)
        cfg_en = _save_cfg_media(cfg_en)
        image = _resolve_activity_media(item.get('image'))
        video = _resolve_activity_media(item.get('video'))
        db.session.add(PracticalActivity(
            lesson_id=lesson.id, section_id=section.id, order=order, lang=lang,
            activity_type=typ, title=title, title_en=title_en, prompt=prompt, prompt_en=prompt_en,
            config_json=json.dumps(cfg, ensure_ascii=False), config_en_json=json.dumps(cfg_en, ensure_ascii=False), image_file=image, video_file=video,
            include_final=bool(item.get('include_final', True))
        ))
        order += 1


@app.route('/teacher/lesson/new', methods=['GET','POST'])
def new_lesson():
    r=require_teacher();
    if r: return r
    if request.method == 'POST':
        subject_name = request.form.get('subject','').strip()
        icon = request.form.get('icon','').strip() or '🌱'
        grade_name = request.form.get('grade','').strip()
        block_title = request.form.get('block','').strip()
        title = request.form.get('title','').strip()
        if not subject_name or not grade_name or not block_title or not title:
            flash('Vyplň předmět, školu a ročník, téma i název lekce. Podle těchto údajů se lekce automaticky zařadí.')
            return render_template('lesson_form.html', course=course_from_lesson(None), lesson=None, section=None, subjects=Subject.query.all(), questions_json=request.form.get('questions_json','[]'), activities_json=request.form.get('activities_json','[]'), questions_json_en=request.form.get('questions_json_en','[]'), activities_json_en=request.form.get('activities_json_en','[]'), gallery_images=[])
        sub = Subject.query.filter_by(name=subject_name).first() or Subject(name=subject_name, icon=icon)
        sub.icon = icon
        db.session.add(sub); db.session.flush()
        gr = Grade.query.filter_by(subject_id=sub.id, name=grade_name).first() or Grade(subject_id=sub.id, name=grade_name)
        db.session.add(gr); db.session.flush()
        bl = Block.query.filter_by(grade_id=gr.id, title=block_title).first() or Block(grade_id=gr.id, title=block_title, order=Block.query.filter_by(grade_id=gr.id).count()+1)
        db.session.add(bl); db.session.flush()
        les = Lesson(block_id=bl.id, title=title, title_en=request.form.get('title_en','').strip(), tip=request.form.get('tip',''), tip_en=request.form.get('tip_en',''), order=Lesson.query.filter_by(block_id=bl.id).count()+1)
        db.session.add(les); db.session.flush()
        html_import = import_html_to_lesson_html(request.files.get('html_file'), request.files.getlist('html_assets'))
        sec_text = html_import if html_import is not None else process_inline_images(request.form.get('text',''))
        # U hotového HTML je nadpis součástí samotného souboru, nepřidáváme proto
        # starý výchozí nadpis aplikace nad něj.
        sec_heading = '' if html_import is not None else request.form.get('heading','Výklad')
        html_import_en = import_html_to_lesson_html(request.files.get('html_file_en'), request.files.getlist('html_assets_en'))
        sec_text_en = html_import_en if html_import_en is not None else process_inline_images(request.form.get('text_en',''))
        sec_heading_en = '' if html_import_en is not None else request.form.get('heading_en','')
        sec = Section(lesson_id=les.id, heading=sec_heading, text=sec_text, interest=request.form.get('interest',''), activity=request.form.get('activity',''),
                      heading_en=sec_heading_en, text_en=sec_text_en, interest_en=request.form.get('interest_en',''), activity_en=request.form.get('activity_en',''), order=1)
        db.session.add(sec); db.session.flush()
        image_map = save_question_images()
        handle_images(les, sec, image_map, 'cs')
        handle_images(les, sec, image_map, 'en')
        add_questions_from_payload(les.id, sec.id, 'study', request.form.get('questions_json',''), request.form.get('study_questions',''), image_map, lang='cs')
        add_questions_from_payload(les.id, sec.id, 'study', request.form.get('questions_json_en',''), '', image_map, lang='en')
        save_practical_activities_from_payload(les, sec, request.form.get('activities_json','[]'), 'cs')
        db.session.commit()
        return redirect(url_for('lesson', lesson_id=les.id))
    return render_template('lesson_form.html', course=course_from_lesson(None), lesson=None, section=None, subjects=Subject.query.all(), questions_json='[]', activities_json='[]', questions_json_en='[]', activities_json_en='[]', gallery_images=[])

@app.route('/teacher/lesson/<int:lesson_id>/edit', methods=['GET','POST'])
def edit_lesson(lesson_id):
    r=require_teacher();
    if r: return r
    les = db.session.get(Lesson, lesson_id)
    if not les: return 'Lekce nenalezena', 404
    sec = sorted(les.sections, key=lambda s:s.order)[0] if les.sections else Section(lesson_id=les.id, heading='Výklad')
    if request.method == 'POST':
        les.block.grade.subject.name = request.form.get('subject', les.block.grade.subject.name)
        les.block.grade.subject.icon = request.form.get('icon', les.block.grade.subject.icon)
        les.block.grade.name = request.form.get('grade', les.block.grade.name)
        les.block.title = request.form.get('block', les.block.title)
        les.title = request.form.get('title', les.title)
        les.title_en = request.form.get('title_en','').strip()
        les.tip = request.form.get('tip','')
        les.tip_en = request.form.get('tip_en','')
        requested_heading = request.form.get('heading','Výklad')
        existing_text = sec.text
        html_import = import_html_to_lesson_html(request.files.get('html_file'), request.files.getlist('html_assets'))
        if html_import is not None:
            sec.heading = ''
            sec.text = html_import
        else:
            # Při běžné editaci zachováme vlastní nadpis. U již importovaného HTML
            # ale starý demonstrační nadpis nezobrazujeme.
            sec.heading = '' if 'imported-html-content' in (existing_text or '') else requested_heading
            sec.text = process_inline_images(request.form.get('text', existing_text) or existing_text)
        sec.interest = request.form.get('interest','')
        sec.activity = request.form.get('activity','')
        existing_text_en = sec.text_en or ''
        html_import_en = import_html_to_lesson_html(request.files.get('html_file_en'), request.files.getlist('html_assets_en'))
        if html_import_en is not None:
            sec.heading_en = ''
            sec.text_en = html_import_en
        else:
            sec.heading_en = request.form.get('heading_en','')
            sec.text_en = process_inline_images(request.form.get('text_en', existing_text_en) or existing_text_en)
        sec.interest_en = request.form.get('interest_en','')
        sec.activity_en = request.form.get('activity_en','')
        Question.query.filter_by(lesson_id=les.id).delete()
        image_map = save_question_images()
        handle_images(les, sec, image_map, 'cs')
        handle_images(les, sec, image_map, 'en')
        add_questions_from_payload(les.id, sec.id, 'study', request.form.get('questions_json',''), request.form.get('study_questions',''), image_map, lang='cs')
        add_questions_from_payload(les.id, sec.id, 'study', request.form.get('questions_json_en',''), '', image_map, lang='en')
        save_practical_activities_from_payload(les, sec, request.form.get('activities_json','[]'), 'cs')
        db.session.commit()
        return redirect(url_for('lesson', lesson_id=les.id))
    return render_template('lesson_form.html', course=course_from_lesson(les), lesson=les, section=sec, subjects=Subject.query.all(), questions_json=questions_editor_json(les, 'study', 'cs'), activities_json=practical_activities_editor_json(les, 'cs'), questions_json_en=questions_editor_json(les, 'study', 'en'), activities_json_en=practical_activities_editor_json(les, 'en'), gallery_images=lesson_gallery(les))

def questions_editor_json(lesson, area, lang='cs'):
    arr=[]
    for q in sorted([q for q in lesson.questions if q.area==area and (q.lang or 'cs') == lang], key=lambda x:x.order):
        if q.qtype == 'choice':
            arr.append({'type':'choice','question':q.question,'options':json.loads(q.options_json or '[]'),'correct':json.loads(q.correct_json or '0')})
        elif q.qtype == 'image_choice':
            arr.append({'type':'image_choice','question':q.question,'images':json.loads(q.options_json or '[]'),'correct':json.loads(q.correct_json or '0')})
        else:
            try:
                opts = json.loads(q.options_json or '{}')
            except Exception:
                opts = {}
            arr.append({'type':'text','question':q.question,'roots':json.loads(q.roots_json or '[]'),'image': opts.get('image','') if isinstance(opts, dict) else ''})
    return json.dumps(arr, ensure_ascii=False)

def add_questions_from_payload(lesson_id, section_id, area, payload, fallback_raw='', image_map=None, lang='cs'):
    try:
        data = json.loads(payload or '[]')
    except Exception:
        data = []
    image_map = image_map or {}
    order = 1
    if data:
        for item in data:
            typ = item.get('type','choice')
            question = (item.get('question') or item.get('prompt') or '').strip()
            if not question: continue
            if typ == 'text':
                roots = [r.strip() for r in item.get('roots',[]) if str(r).strip()]
                img = str(item.get('image','') or '').strip()
                img = image_map.get(img, img)
                db.session.add(Question(lesson_id=lesson_id, section_id=section_id, area=area, lang=lang, qtype='text', question=question, options_json=json.dumps({'image': img}, ensure_ascii=False), roots_json=json.dumps(roots, ensure_ascii=False), hint='Odpověď najdeš ve výkladu.', order=order))
            elif typ == 'image_choice':
                imgs = [image_map.get(str(o).strip(), str(o).strip()) for o in item.get('images',[]) if str(o).strip()]
                if len(imgs) < 2:
                    continue
                while len(imgs) < 4:
                    imgs.append(imgs[-1])
                correct = int(item.get('correct',0) or 0)
                db.session.add(Question(lesson_id=lesson_id, section_id=section_id, area=area, lang=lang, qtype='image_choice', question=question, options_json=json.dumps(imgs[:4], ensure_ascii=False), correct_json=json.dumps(correct), hint='Odpověď najdeš ve výkladu.', order=order))
            else:
                opts = [str(o).strip() for o in item.get('options',[]) if str(o).strip()]
                while len(opts) < 2: opts.append('')
                correct = int(item.get('correct',0) or 0)
                db.session.add(Question(lesson_id=lesson_id, section_id=section_id, area=area, lang=lang, qtype='choice', question=question, options_json=json.dumps(opts, ensure_ascii=False), correct_json=json.dumps(correct), hint='Odpověď najdeš ve výkladu.', order=order))
            order += 1
        return
    add_questions_from_text(lesson_id, section_id, area, fallback_raw or '', lang=lang)

def add_questions_from_text(lesson_id, section_id, area, raw, lang='cs'):
    # formát: otázka | odpověď A | odpověď B | odpověď C | číslo správné odpovědi 1-3
    order = 1
    for line in raw.splitlines():
        line=line.strip()
        if not line or line.startswith('#'): continue
        parts=[p.strip() for p in line.split('|')]
        if len(parts)>=5:
            correct = max(0, int(parts[4])-1) if parts[4].isdigit() else 0
            db.session.add(Question(lesson_id=lesson_id, section_id=section_id, area=area, lang=lang, qtype='choice', question=parts[0], options_json=json.dumps(parts[1:4], ensure_ascii=False), correct_json=json.dumps(correct), hint='Odpověď najdeš ve výkladu.', order=order))
            order += 1




def _save_raw_image_bytes(raw, ext='png'):
    ext = (ext or 'png').lower().strip('.').replace('jpeg','jpg')
    if ext not in {'png','jpg','jpeg','gif','webp','svg'}:
        ext = 'png'
    name = datetime.now().strftime('%Y%m%d%H%M%S_') + uuid.uuid4().hex[:10] + '.' + ext
    (UPLOADS / name).write_bytes(raw)
    return name


def _decode_text_bytes(data):
    for enc in ('utf-8-sig', 'utf-8', 'cp1250', 'windows-1250', 'latin-1'):
        try:
            return data.decode(enc)
        except Exception:
            pass
    return data.decode('utf-8', errors='ignore')


def import_html_to_lesson_html(html_file, asset_files):
    """Naimportuje hotový HTML výklad.

    Učitel si připraví výklad mimo aplikaci (např. převod DOCX -> HTML přes pandoc/Word).
    V aplikaci vybere HTML soubor a případně obrázky ze stejné složky.
    Funkce zkopíruje obrázky do trvalého úložiště uploads a přepíše cesty v HTML.
    """
    if not html_file or not html_file.filename:
        return None
    filename = secure_filename(html_file.filename or '')
    if not filename.lower().endswith(('.html', '.htm')):
        raise ValueError('Vyber soubor ve formátu .html nebo .htm.')

    html_text = _decode_text_bytes(html_file.read())

    # Pokud je to celá stránka, vezmeme hlavně obsah body, aby se do lekce netáhla hlavička dokumentu.
    m = re.search(r'<body[^>]*>(.*?)</body>', html_text, flags=re.I | re.S)
    if m:
        html_text = m.group(1)

    # Odstraníme prvky, které do vloženého výkladu nepatří.
    html_text = re.sub(r'<script\b[^>]*>.*?</script>', '', html_text, flags=re.I | re.S)
    html_text = re.sub(r'<style\b[^>]*>.*?</style>', '', html_text, flags=re.I | re.S)
    html_text = re.sub(r'<link\b[^>]*>', '', html_text, flags=re.I | re.S)
    html_text = re.sub(r'<meta\b[^>]*>', '', html_text, flags=re.I | re.S)

    # Uložíme všechny obrázky, které učitel přiložil k HTML, a namapujeme je podle názvu souboru.
    image_map = {}
    for f in asset_files or []:
        if not f or not f.filename:
            continue
        if not (f.mimetype or '').startswith('image/') and not f.filename.lower().endswith(('.png','.jpg','.jpeg','.gif','.webp','.svg')):
            continue
        original_name = Path(secure_filename(Path(f.filename).name)).name
        saved = save_upload(f)
        if saved:
            image_map[original_name] = saved
            image_map[original_name.lower()] = saved

    # Přepíšeme src obrázků. Umíme i data:image z HTML.
    def replace_src(match):
        prefix, src, suffix = match.group(1), match.group(2), match.group(3)
        src_clean = src.strip()
        if src_clean.startswith('data:image/'):
            dm = re.match(r'data:(image/[^;]+);base64,(.*)', src_clean, flags=re.I | re.S)
            if dm:
                mime = dm.group(1).lower()
                data = dm.group(2)
                ext = 'png'
                if 'jpeg' in mime or 'jpg' in mime: ext = 'jpg'
                elif 'gif' in mime: ext = 'gif'
                elif 'webp' in mime: ext = 'webp'
                elif 'svg' in mime: ext = 'svg'
                try:
                    saved = _save_raw_image_bytes(base64.b64decode(data), ext)
                    return f'{prefix}{url_for("uploads", filename=saved)}{suffix}'
                except Exception:
                    return match.group(0)
        if src_clean.startswith(('http://','https://','/uploads/','/static/','data:')):
            return match.group(0)
        parsed = urllib.parse.urlparse(src_clean)
        base = Path(urllib.parse.unquote(parsed.path)).name
        safe_base = secure_filename(base)
        saved = image_map.get(base) or image_map.get(base.lower()) or image_map.get(safe_base) or image_map.get(safe_base.lower())
        if saved:
            return f'{prefix}{url_for("uploads", filename=saved)}{suffix}'
        # Když obrázek nebyl přiložen, necháme cestu být a zobrazíme varování v náhledu přes alt/title pro snazší hledání.
        return match.group(0)

    html_text = re.sub(r'(src\s*=\s*["\'])([^"\']+)(["\'])', replace_src, html_text, flags=re.I)
    html_text = html_text.strip()

    # Některé exporty HTML přidají na úplný začátek samostatné číslo stránky/oddílu
    # (např. <p>1</p> nebo <h1>1</h1>). V lekci se pak zobrazovalo osamocené "1".
    # Odstraníme ho jen tehdy, když jde opravdu o první samostatný blok s číslem.
    html_text = re.sub(
        r'^\s*<(?:p|div|h[1-6])(?:\s+[^>]*)?>\s*\d{1,3}\s*</(?:p|div|h[1-6])>\s*',
        '',
        html_text,
        count=1,
        flags=re.I | re.S,
    )

    if not html_text:
        raise ValueError('HTML soubor neobsahuje žádný výklad.')
    info = '<div class="docx-import-note"><b>Importováno z HTML:</b> ' + html.escape(filename) + '</div>'
    return info + '<div class="imported-html-content">' + html_text + '</div>'

def process_inline_images(html_text):
    """Uloží obrázky vložené do CKEditoru jako data:image/... a přepíše je na /uploads/...
    Tím zachráníme obrázky vložené přes Ctrl+V, které se jinak po uložení ztratí.
    """
    if not html_text:
        return ''

    def repl(match):
        mime = match.group(1).lower()
        data = match.group(2)
        ext = 'png'
        if 'jpeg' in mime or 'jpg' in mime:
            ext = 'jpg'
        elif 'gif' in mime:
            ext = 'gif'
        elif 'webp' in mime:
            ext = 'webp'
        elif 'svg' in mime:
            ext = 'svg'
        try:
            raw = base64.b64decode(data)
        except Exception:
            return match.group(0)
        name = datetime.now().strftime('%Y%m%d%H%M%S_') + uuid.uuid4().hex[:10] + '.' + ext
        (UPLOADS / name).write_bytes(raw)
        return 'src="' + url_for('uploads', filename=name) + '"'

    # src="data:image/png;base64,..."
    html_text = re.sub(r'src=["\']data:(image/[^;]+);base64,([^"\']+)["\']', repl, html_text)
    return html_text

def save_upload(file):
    if not file or not file.filename: return ''
    name = datetime.now().strftime('%Y%m%d%H%M%S_') + secure_filename(file.filename)
    file.save(UPLOADS / name)
    return name


def save_math_example_image(file):
    """Uloží nepovinný obrázek/náčrt k matematickému příkladu."""
    if not file or not file.filename:
        return ''
    ext = Path(file.filename).suffix.lower()
    allowed = {'.png', '.jpg', '.jpeg', '.webp', '.gif', '.svg'}
    if ext not in allowed or not (file.mimetype or '').startswith('image/'):
        return ''
    return save_upload(file)

def save_question_images():
    """Uloží obrázky vložené přímo u otázek.
    V JSONu editor používá odkaz ve tvaru __file__:nazev_pole.
    Tady ho převedeme na reálně uložený soubor v trvalém úložišti uploads.
    """
    image_map = {}
    for field_name, f in request.files.items():
        if not field_name.startswith('qimg_'):
            continue
        if not f or not f.filename:
            continue
        original = secure_filename(f.filename)
        saved = save_upload(f)
        if saved:
            image_map[f'__file__:{field_name}'] = saved
            image_map[original] = saved
            image_map[f.filename] = saved
    return image_map

# Starší název necháváme kvůli kompatibilitě, kdyby někde zůstal odkaz.
def save_gallery_images():
    return save_question_images()

def handle_images(les, sec, image_map=None, lang='cs'):
    image_map = image_map or {}
    field = 'hero_image_en' if lang == 'en' else 'hero_image'
    h = save_upload(request.files.get(field))
    if h:
        if lang == 'en': les.hero_image_en = h
        else: les.hero_image = h
    # Obrázky výkladu se vkládají přímo přes CKEditor a ukládají se v endpointu /teacher/upload-image.

@app.route('/uploads/<filename>')
def uploads(filename):
    return send_from_directory(UPLOADS, filename)

@app.route('/teacher/upload-image', methods=['POST'])
def upload_editor_image():
    r = require_teacher()
    if r:
        return jsonify({'error': 'Nepřihlášený uživatel.'}), 401
    f = request.files.get('upload')
    if not f or not f.filename:
        return jsonify({'error': 'Nebyl vybrán žádný obrázek.'}), 400
    if not (f.mimetype or '').startswith('image/'):
        return jsonify({'error': 'Soubor musí být obrázek.'}), 400
    saved = save_upload(f)
    if not saved:
        return jsonify({'error': 'Obrázek se nepodařilo uložit.'}), 400
    return jsonify({'url': url_for('uploads', filename=saved)})

@app.route('/lessons/<slug>/images/<filename>')
def old_img(slug, filename):
    return send_from_directory(BASE/'lessons'/slug/'images', filename)

@app.route('/img/<filename>')
def img(filename):
    return send_from_directory(UPLOADS, filename)

def ensure_schema_updates():
    inspector = inspect(db.engine)
    required = {
        'lesson': {
            'title_en': "VARCHAR(200) DEFAULT ''", 'tip_en': "TEXT DEFAULT ''", 'hero_image_en': "VARCHAR(255) DEFAULT ''"
        },
        'section': {
            'heading_en': "VARCHAR(200) DEFAULT ''", 'text_en': "TEXT DEFAULT ''", 'interest_en': "TEXT DEFAULT ''",
            'image_en': "VARCHAR(255) DEFAULT ''", 'activity_en': "TEXT DEFAULT ''"
        },
        'question': {'lang': "VARCHAR(2) DEFAULT 'cs'"},
        'practical_activity': {'lang': "VARCHAR(2) DEFAULT 'cs'", 'title_en': "VARCHAR(220) DEFAULT ''", 'prompt_en': "TEXT DEFAULT ''", 'config_en_json': "TEXT DEFAULT '{}'"},
        'result': {
            'focus_lost': 'INTEGER DEFAULT 0',
            'status': "VARCHAR(60) DEFAULT 'dokončeno'"
        },
        'interactive_result': {
            'focus_lost': 'INTEGER DEFAULT 0',
            'status': "VARCHAR(60) DEFAULT 'dokončeno'"
        },
        'informatics_submission': {
            'grade': 'INTEGER DEFAULT 5',
            'status': "VARCHAR(60) DEFAULT 'kontrola'",
            'focus_lost': 'INTEGER DEFAULT 0'
        },
        'math_attempt': {
            'answers_json': "TEXT DEFAULT '[]'",
            'variant_json': "TEXT DEFAULT '{}'"
        },
        'math_step': {
            'result_formula': "TEXT DEFAULT ''",
            'result_decimals': 'INTEGER DEFAULT 2'
        },
        'math_example': {
            'prose_problem': "TEXT DEFAULT ''",
            'image_stored': "VARCHAR(255) DEFAULT ''",
            'variant_enabled': 'BOOLEAN DEFAULT 0',
            'variant_values': "TEXT DEFAULT ''",
            'variant_condition': "TEXT DEFAULT ''",
            'variant_min': 'FLOAT DEFAULT 1',
            'variant_max': 'FLOAT DEFAULT 30',
            'variant_step': 'FLOAT DEFAULT 1',
            'variant_result_kind': "VARCHAR(20) DEFAULT 'any'",
            'variant_result_sign': "VARCHAR(20) DEFAULT 'any'",
            'variant_result_min': 'FLOAT',
            'variant_result_max': 'FLOAT',
            'variant_result_decimals': 'INTEGER DEFAULT -1'
        }
    }
    for table_name, columns in required.items():
        existing = {c['name'] for c in inspector.get_columns(table_name)} if inspector.has_table(table_name) else set()
        for column_name, sql_type in columns.items():
            if column_name not in existing:
                db.session.execute(text(
                    f'ALTER TABLE {table_name} ADD COLUMN {column_name} {sql_type}'
                ))
    db.session.commit()


def seed():
    db.create_all()
    ensure_informatics_columns()
    ensure_schema_updates()
    restore_interactive_lessons_from_files()
    tu = os.getenv('TEACHER_USERNAME', 'dnadler').lower()
    tp = os.getenv('TEACHER_PASSWORD', 'change-me')
    tn = os.getenv('TEACHER_NAME', 'Učitel')
    teacher = User.query.filter_by(username=tu).first()
    if not teacher:
        db.session.add(User(username=tu, name=tn, role='teacher', password_hash=generate_password_hash(tp)))
    else:
        # Učitelský účet je řízený přes .env, aby heslo nebylo natvrdo v kódu.
        teacher.name = tn
        teacher.role = 'teacher'
        teacher.password_hash = generate_password_hash(tp)

    # Starý anonymní demo účet student/student nechceme.
    demo = User.query.filter_by(username='student').first()
    if demo:
        db.session.delete(demo)
        db.session.flush()

    # Ukázkový konkrétní student pro vyzkoušení přihlášení. Další studenty vytvoří učitel v editoru.
    if not User.query.filter_by(username='jan.novak').first():
        db.session.add(User(username='jan.novak', name='Jan Novák', role='student', password_hash=generate_password_hash('zive123')))
    if os.getenv('SEED_DEMO_LESSON', '0') == '1' and Subject.query.count()==0:
        # Volitelná demo lekce. Ve skutečném provozu se sama neobnovuje po smazání.
        # zkopíruj ukázkové obrázky ze staré lekce do uploads
        old = BASE/'lessons'/'bio6_01_co_je_zive'/'images'
        for n in ['1.jpg','2.jpg','3.jpg']:
            if (old/n).exists() and not (UPLOADS/n).exists():
                (UPLOADS/n).write_bytes((old/n).read_bytes())
        sub=Subject(name='Biologie', icon='🌱'); db.session.add(sub); db.session.flush()
        gr=Grade(subject_id=sub.id, name='6. ročník'); db.session.add(gr); db.session.flush()
        bl=Block(grade_id=gr.id, title='Blok 1 – Život kolem nás', order=1); db.session.add(bl); db.session.flush()
        les=Lesson(block_id=bl.id, title='Co je živé a neživé', tip='Čti výklad jako detektiv. Každá odpověď je někde v textu.', hero_image='1.jpg', order=1); db.session.add(les); db.session.flush()
        sec=Section(lesson_id=les.id, heading='Jak poznáme živé organismy?', text='Živé organismy rostou, dýchají, přijímají živiny, reagují na okolí, rozmnožují se a skládají se z buněk. Neživé věci tyto znaky života samy nevykazují. Pes, strom nebo houba jsou živé organismy. Kámen, lavice nebo sklenice jsou neživé věci.', interest='Některé věci mohou vypadat jako živé, například plamen svíčky se pohybuje, ale není organismus.', image='2.jpg', activity='Rozhlédni se kolem sebe a napiš si 3 živé organismy a 3 neživé věci.', order=1); db.session.add(sec); db.session.flush()
        db.session.add(InlineImage(section_id=sec.id, file='3.jpg', caption='Ukázka přírody: živé organismy a neživé prostředí.', order=1))
        raw='''Který příklad je živý organismus? | pes | kámen | lavice | 1
Co patří mezi znaky života? | přijímání živin | tvrdost kamene | barva lavice | 1
Která věc je neživá? | strom | houba | sklenice | 3'''
        add_questions_from_text(les.id, sec.id, 'study', raw)
    # Odstraní i prázdné odkazy, které mohly zůstat v databázi po starších verzích.
    cleanup_empty_curriculum()
    db.session.commit()

with app.app_context():
    seed()

if __name__ == '__main__':
    app.run(debug=True)
