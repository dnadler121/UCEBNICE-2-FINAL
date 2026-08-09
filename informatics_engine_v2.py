import ast
import json
import math
import os
import re
import subprocess
import sys
import tempfile
import zipfile
from collections import Counter
from pathlib import Path
from xml.etree import ElementTree as ET

NS_W = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'


def _xml(z, name):
    try:
        return ET.fromstring(z.read(name))
    except Exception:
        return None


def _attr(el, ns, name, default=''):
    if el is None:
        return default
    return el.attrib.get(f'{{{ns}}}{name}', el.attrib.get(name, default))


def _norm_value(v):
    if v is None:
        return None
    if isinstance(v, float):
        if math.isfinite(v):
            return round(v, 10)
        return str(v)
    if hasattr(v, 'isoformat'):
        try:
            return v.isoformat()
        except Exception:
            pass
    return str(v).strip() if not isinstance(v, (int, bool)) else v


def _formula_result_map(path):
    """Returns cached results for formula cells when the workbook contains them."""
    try:
        import openpyxl
        wb_formula = openpyxl.load_workbook(path, data_only=False, read_only=False)
        wb_values = openpyxl.load_workbook(path, data_only=True, read_only=False)
        out = {}
        for ws in wb_formula.worksheets:
            if ws.title == '__UCEBNICE_ID__':
                continue
            wsv = wb_values[ws.title]
            for row in ws.iter_rows():
                for c in row:
                    if isinstance(c.value, str) and c.value.startswith('='):
                        out[f'{ws.title}!{c.coordinate}'] = _norm_value(wsv[c.coordinate].value)
        return out
    except Exception:
        return {}




def _excel_ooxml_profile(path):
    """Read structural Excel features directly from the OOXML package."""
    prof = {'sheet_xml_count':0,'formula_cells':{},'merge_refs':{},'auto_filters':{},
            'conditional_counts':{},'table_part_counts':{},'drawing_counts':{},
            'table_xml_count':0,'chart_xml_count':0,'pivot_xml_count':0}
    try:
        with zipfile.ZipFile(path) as z:
            names=z.namelist()
            prof['table_xml_count']=sum(1 for n in names if re.match(r'xl/tables/table\d+\.xml$',n))
            prof['chart_xml_count']=sum(1 for n in names if re.match(r'xl/charts/chart\d+\.xml$',n))
            prof['pivot_xml_count']=sum(1 for n in names if re.match(r'xl/pivotTables/pivotTable\d+\.xml$',n))
            for n in sorted(x for x in names if re.match(r'xl/worksheets/sheet\d+\.xml$',x)):
                root=_xml(z,n)
                if root is None: continue
                prof['sheet_xml_count'] += 1
                key=Path(n).stem
                ns='{http://schemas.openxmlformats.org/spreadsheetml/2006/main}'
                formulas={}
                for c in root.iter(ns+'c'):
                    f=c.find(ns+'f')
                    if f is not None:
                        formulas[c.attrib.get('r','')]=(f.text or '').strip()
                prof['formula_cells'][key]=formulas
                merges=root.find(ns+'mergeCells')
                prof['merge_refs'][key]=sorted(mc.attrib.get('ref','') for mc in (list(merges) if merges is not None else []))
                af=root.find(ns+'autoFilter')
                prof['auto_filters'][key]=af.attrib.get('ref','') if af is not None else ''
                prof['conditional_counts'][key]=len(root.findall(ns+'conditionalFormatting'))
                tp=root.find(ns+'tableParts')
                prof['table_part_counts'][key]=len(list(tp)) if tp is not None else 0
                prof['drawing_counts'][key]=len(root.findall(ns+'drawing'))
    except Exception:
        pass
    return prof


def _ppt_ooxml_profile(path):
    """Read per-slide PowerPoint structure directly from the OOXML package."""
    slides=[]
    try:
        with zipfile.ZipFile(path) as z:
            slide_names=sorted((n for n in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml$',n)),
                               key=lambda n:int(re.search(r'(\d+)',Path(n).stem).group(1)))
            for n in slide_names:
                root=_xml(z,n)
                if root is None: continue
                counts=Counter(); positions=[]; texts=[]
                for sp in root.iter('{%s}sp' % NS_P):
                    counts['shapes'] += 1
                    tx=sp.find('.//{%s}txBody' % NS_P)
                    if tx is not None:
                        t=''.join((e.text or '') for e in tx.iter('{%s}t' % NS_A)).strip()
                        if t: texts.append(t)
                for _ in root.iter('{%s}pic' % NS_P): counts['pictures'] += 1
                for gf in root.iter('{%s}graphicFrame' % NS_P):
                    counts['graphic_frames'] += 1
                    xml=ET.tostring(gf,encoding='unicode')
                    if '/chart' in xml: counts['charts'] += 1
                    if '/table' in xml or '<a:tbl' in xml: counts['tables'] += 1
                for xfrm in root.iter('{%s}xfrm' % NS_A):
                    off=xfrm.find('{%s}off' % NS_A); ext=xfrm.find('{%s}ext' % NS_A)
                    if off is not None and ext is not None:
                        positions.append((off.attrib.get('x'),off.attrib.get('y'),ext.attrib.get('cx'),ext.attrib.get('cy')))
                trans=root.find('{%s}transition' % NS_P)
                transition = ET.tostring(trans,encoding='unicode') if trans is not None else ''
                anim=0
                tags={'anim','animClr','animEffect','animMotion','animRot','animScale','set','cmd'}
                for e in root.iter():
                    if e.tag.startswith('{%s}' % NS_P) and e.tag.split('}',1)[1] in tags: anim += 1
                slides.append({'counts':dict(counts),'positions':positions,'texts':texts,'transition':bool(transition),'animations':anim})
    except Exception:
        pass
    return {'slides':slides}

def analyze_excel(path, original_name):
    import openpyxl
    from openpyxl.styles.borders import Border

    wb = openpyxl.load_workbook(path, data_only=False, read_only=False, keep_vba=Path(original_name).suffix.lower()=='.xlsm')
    visible_sheets = [s for s in wb.sheetnames if s != '__UCEBNICE_ID__']
    total_nonempty = total_formulas = chart_count = table_count = 0
    functions = set()
    sheet_specs = []
    formatting = Counter()
    font_names = Counter(); font_sizes = Counter(); number_formats = Counter()
    merged_count = conditional_count = filter_count = 0
    chart_types = Counter()
    formula_cells = []

    for ws in wb.worksheets:
        if ws.title == '__UCEBNICE_ID__':
            continue
        nonempty = 0; formulas = []
        max_row = min(ws.max_row or 1, 1000); max_col = min(ws.max_column or 1, 100)
        headers = [str(c.value) for c in ws[1][:max_col] if c.value not in (None, '')]
        for row in ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col):
            for cell in row:
                if cell.value not in (None, ''):
                    nonempty += 1
                    if cell.font:
                        if cell.font.name: font_names[str(cell.font.name)] += 1
                        if cell.font.sz: font_sizes[str(round(float(cell.font.sz),2))] += 1
                        if cell.font.bold: formatting['bold'] += 1
                        if cell.font.italic: formatting['italic'] += 1
                        if cell.font.underline: formatting['underline'] += 1
                    if cell.fill and getattr(cell.fill, 'fill_type', None): formatting['fill'] += 1
                    if cell.alignment:
                        if cell.alignment.horizontal: formatting[f'halign:{cell.alignment.horizontal}'] += 1
                        if cell.alignment.vertical: formatting[f'valign:{cell.alignment.vertical}'] += 1
                        if cell.alignment.wrap_text: formatting['wrap_text'] += 1
                    if cell.border and any(getattr(getattr(cell.border, side, None), 'style', None) for side in ('left','right','top','bottom')):
                        formatting['bordered'] += 1
                    if cell.number_format and cell.number_format != 'General':
                        number_formats[str(cell.number_format)] += 1
                if isinstance(cell.value, str) and cell.value.startswith('='):
                    total_formulas += 1
                    formula_cells.append(f'{ws.title}!{cell.coordinate}')
                    formulas.append({'cell': cell.coordinate, 'formula': cell.value})
                    for fn in re.findall(r'([A-Z][A-Z0-9_.]*)\s*\(', cell.value.upper()): functions.add(fn)
        total_nonempty += nonempty
        charts = list(getattr(ws, '_charts', []) or [])
        chart_count += len(charts)
        for ch in charts:
            chart_types[type(ch).__name__] += 1
        tables = getattr(ws, 'tables', {})
        table_count += len(tables)
        merged_count += len(ws.merged_cells.ranges)
        try: conditional_count += len(ws.conditional_formatting)
        except Exception: pass
        if ws.auto_filter and ws.auto_filter.ref: filter_count += 1
        sheet_specs.append({'name': ws.title, 'rows': ws.max_row or 0, 'cols': ws.max_column or 0,
                            'headers': headers, 'nonempty': nonempty, 'formulas': formulas[:80]})

    pivot_count = 0
    try:
        with zipfile.ZipFile(path) as z:
            pivot_count = len([n for n in z.namelist() if n.startswith('xl/pivotTables/pivotTable') and n.endswith('.xml')])
    except Exception:
        pass

    result_map = _formula_result_map(path)
    info = {
        'extension': Path(original_name).suffix.lower(), 'name': original_name, 'type': 'Excel',
        'sheets': visible_sheets, 'sheet_specs': sheet_specs, 'nonempty_count': total_nonempty,
        'formula_count': total_formulas, 'formula_functions': sorted(functions), 'formula_cells': formula_cells,
        'formula_results': result_map, 'chart_count': chart_count, 'chart_types': dict(chart_types),
        'table_count': table_count, 'pivot_count': pivot_count, 'merged_count': merged_count,
        'conditional_format_count': conditional_count, 'filter_count': filter_count,
        'font_names': dict(font_names), 'font_sizes': dict(font_sizes), 'number_formats': dict(number_formats),
        'formatting': dict(formatting), 'ooxml': _excel_ooxml_profile(path),
    }
    checks = [
        ('excel_sheets','Počet a názvy listů'), ('excel_headers','Záhlaví tabulky'),
        ('excel_size','Rozsah tabulky / počet řádků a sloupců'), ('excel_filled','Vyplněné části tabulky'),
        ('excel_results','Správné výsledky výpočtů (vzorec může být jiný)'), ('excel_functions','Použití požadovaných funkcí'),
        ('excel_format','Formátování buněk – písmo, velikost, zarovnání, ohraničení, výplň'),
        ('excel_number_format','Formát čísel – měna, procenta, datum, desetinná místa'),
        ('excel_merged','Sloučené buňky'), ('excel_conditional','Podmíněné formátování'),
        ('excel_filter','Filtr / automatický filtr'), ('excel_table','Excelová tabulka'),
        ('excel_chart','Grafy'), ('excel_chart_type','Typy grafů'), ('excel_pivot','Kontingenční tabulka'),
    ]
    # Only offer structure features if the teacher file uses them, but always offer core checks.
    core = {'excel_sheets','excel_headers','excel_size','excel_filled','excel_results','excel_functions','excel_format','excel_number_format'}
    feature_pred = {
        'excel_merged': merged_count, 'excel_conditional': conditional_count, 'excel_filter': filter_count,
        'excel_table': table_count, 'excel_chart': chart_count, 'excel_chart_type': chart_count, 'excel_pivot': pivot_count,
    }
    return info, [{'code':c,'label':l} for c,l in checks if c in core or feature_pred.get(c,0)]


def _docx_fields(doc_xml_text):
    return [x.strip() for x in re.findall(r'<w:instrText[^>]*>(.*?)</w:instrText>', doc_xml_text, flags=re.I|re.S)]


def analyze_word(path, original_name):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = Document(path)
    paragraphs = [p for p in doc.paragraphs if p.text.strip()]
    headings = [p for p in paragraphs if str(p.style.name or '').lower().startswith(('heading','nadpis'))]
    words = [w for p in paragraphs for w in p.text.split()]
    heading_styles = Counter(str(p.style.name) for p in headings)
    paragraph_align = Counter(str(p.alignment) for p in paragraphs if p.alignment is not None)
    centered_paragraphs = sum(1 for p in paragraphs if p.alignment == WD_ALIGN_PARAGRAPH.CENTER)
    captions = [p for p in paragraphs if str(p.style.name or '').lower() in ('caption','titulek')]
    centered_captions = sum(1 for p in captions if p.alignment == WD_ALIGN_PARAGRAPH.CENTER)

    image_specs=[]
    for shp in doc.inline_shapes:
        image_specs.append({'width_cm': round(shp.width/360000,2), 'height_cm': round(shp.height/360000,2)})

    image_paragraphs=[p for p in doc.paragraphs if p._p.xpath('.//w:drawing')]
    centered_image_count=sum(1 for p in image_paragraphs if p.alignment == WD_ALIGN_PARAGRAPH.CENTER)
    section_specs=[]
    for idx, sec in enumerate(doc.sections,1):
        valign=''
        try:
            vel=sec._sectPr.find('{%s}vAlign' % NS_W)
            valign=_attr(vel,NS_W,'val','')
        except Exception:
            pass
        section_specs.append({
            'index':idx,
            'header_linked': bool(sec.header.is_linked_to_previous) if idx>1 else False,
            'footer_linked': bool(sec.footer.is_linked_to_previous) if idx>1 else False,
            'vertical_alignment': valign,
        })

    info = {
        'extension':'.docx','name':original_name,'type':'Word','paragraph_count':len(paragraphs),
        'word_count':len(words),'heading_count':len(headings),'headings':[p.text.strip() for p in headings[:40]],
        'heading_styles':dict(heading_styles),'table_count':len(doc.tables),'image_count':len(doc.inline_shapes),
        'image_specs':image_specs,'section_count':len(doc.sections),'section_specs':section_specs,
        'paragraph_alignments':dict(paragraph_align),'centered_paragraph_count':centered_paragraphs,
        'caption_count':len(captions),'centered_caption_count':centered_captions,'centered_image_count':centered_image_count,
        'first_paragraph_text': paragraphs[0].text.strip() if paragraphs else '',
        'first_paragraph_style': str(paragraphs[0].style.name) if paragraphs else '',
    }

    with zipfile.ZipFile(path) as z:
        names=set(z.namelist())
        document_xml = z.read('word/document.xml').decode('utf-8','replace') if 'word/document.xml' in names else ''
        root = _xml(z,'word/document.xml')
        fields=[]
        if root is not None:
            # Čti Word pole po jednotlivých odstavcích. U složených polí Word často
            # rozdělí instrukci do více w:instrText uzlů, proto je spojíme jen v
            # rámci stejného odstavce (nikoli přes celý dokument).
            for p_el in root.findall('.//{%s}p' % NS_W):
                parts=[]
                for instr in p_el.findall('.//{%s}instrText' % NS_W):
                    if instr.text:
                        parts.append(instr.text)
                if parts:
                    val=''.join(parts).strip()
                    if val: fields.append(val)
                for fld in p_el.findall('.//{%s}fldSimple' % NS_W):
                    val=_attr(fld,NS_W,'instr','')
                    if val: fields.append(val.strip())
        field_text=' '.join(fields).upper()
        # Automatický obsah uznáme jen tehdy, když dokument skutečně obsahuje
        # pole TOC. Samotný content-control/docPartGallery nestačí: Word jej může
        # v DOCX ponechat i po smazání viditelného obsahu, což dříve vedlo k
        # falešnému výsledku „Splněno“.
        toc_fields = [f for f in fields if re.search(r'(^|\s)TOC(?=\s|$|\\)', f.upper())]
        # Word používá stejné pole TOC i pro Seznam obrázků/tabulek.
        # Např. `TOC \h \z \c "Obrázek"` je seznam obrázků, NE obsah kapitol.
        # Automatický obsah proto uznáme jen u TOC pole bez přepínače \c / \a.
        # Běžný obsah z Wordu mívá např. `TOC \o "1-3" \h \z \u`.
        def _is_list_toc(field):
            up = (' ' + str(field).upper() + ' ')
            return bool(re.search(r'\\[CA](?=\s|$)', up))
        info['has_toc'] = any(not _is_list_toc(f) for f in toc_fields)
        bibliography_fields = [f for f in fields if re.search(r'(^|\s)BIBLIOGRAPHY(?=\s|$|\\)', f.upper())]
        info['has_bibliography'] = bool(bibliography_fields)
        info['citation_field_count'] = len(re.findall(r'\bCITATION\b', field_text))
        info['has_list_of_figures'] = any(_is_list_toc(f) for f in toc_fields)
        info['page_break_count'] = document_xml.count('w:type="page"') + document_xml.count("w:type='page'")
        # Číslování nadpisů může být přímo v odstavci NEBO zděděné ze stylu.
        style_num_ids={}
        styles_root=_xml(z,'word/styles.xml') if 'word/styles.xml' in names else None
        if styles_root is not None:
            for st in styles_root.findall('.//{%s}style' % NS_W):
                sid=_attr(st,NS_W,'styleId','')
                ppr=st.find('{%s}pPr' % NS_W)
                numpr=ppr.find('{%s}numPr' % NS_W) if ppr is not None else None
                numid=numpr.find('{%s}numId' % NS_W) if numpr is not None else None
                if sid and numid is not None:
                    style_num_ids[sid]=_attr(numid,NS_W,'val','')
        numbered=0
        if root is not None:
            for p_el in root.findall('.//{%s}p' % NS_W):
                ppr=p_el.find('{%s}pPr' % NS_W)
                if ppr is None: continue
                pstyle=ppr.find('{%s}pStyle' % NS_W)
                sid=_attr(pstyle,NS_W,'val','')
                low=sid.lower()
                if not ('heading' in low or 'nadpis' in low): continue
                direct=ppr.find('{%s}numPr' % NS_W) is not None
                inherited=bool(style_num_ids.get(sid))
                if direct or inherited: numbered += 1
        info['numbered_heading_count']=numbered
        # Page numbering starts: pgNumType in section properties.
        starts=[]
        if root is not None:
            for i,sect in enumerate(root.findall('.//{%s}sectPr' % NS_W),1):
                pg=sect.find('{%s}pgNumType' % NS_W)
                if pg is not None and _attr(pg,NS_W,'start','') != '':
                    starts.append({'section':i,'start':_attr(pg,NS_W,'start','')})
        info['page_number_starts']=starts
        # Saved page count from docProps/app.xml (Word/LibreOffice updates this when saved).
        pages=None
        appxml=_xml(z,'docProps/app.xml') if 'docProps/app.xml' in names else None
        if appxml is not None:
            for e in appxml.iter():
                if e.tag.endswith('Pages'):
                    try: pages=int(e.text or '0')
                    except Exception: pass
        info['saved_page_count']=pages
        # Bibliography sources are typically stored in customXml.
        source_count=0
        for n in names:
            if n.startswith('customXml/') and n.endswith('.xml'):
                try:
                    txt=z.read(n).decode('utf-8','replace')
                    source_count += len(re.findall(r'<(?:[A-Za-z0-9_]+:)?Source\b',txt))
                except Exception: pass
        info['bibliography_source_count']=source_count

    checks=[
        ('word_sections','Počet oddílů'), ('word_unlinked','Vypnuté „Propojit s předchozím“ v oddílech'),
        ('word_page_numbering','Číslování stránek od určeného oddílu/čísla'),
        ('word_alignment','Vodorovné a svislé zarovnání'), ('word_images','Počet vložených obrázků'),
        ('word_image_size','Velikost obrázků'), ('word_image_center','Obrázky zarovnané na střed'),
        ('word_title_first','Nadpis / titul na první stránce'), ('word_toc','Automatický obsah'),
        ('word_heading_styles','Použité styly nadpisů'), ('word_heading_numbering','Číslování nadpisů'),
        ('word_captions','Titulky pod obrázky a jejich zarovnání'), ('word_citations','Vložené citace'),
        ('word_sources','Vložené zdroje citací'), ('word_pages','Počet stran dokumentu'), ('word_list_figures','Seznam obrázků'),
        ('word_bibliography','Seznam citací / bibliografie'),
    ]
    pred={
        'word_sections':info['section_count']>1, 'word_unlinked':info['section_count']>1,
        'word_page_numbering':bool(info['page_number_starts']), 'word_alignment':True,
        'word_images':info['image_count']>0, 'word_image_size':info['image_count']>0, 'word_image_center':info['image_count']>0,
        'word_title_first':True, 'word_toc':info['has_toc'], 'word_heading_styles':bool(info['heading_styles']),
        'word_heading_numbering':info['numbered_heading_count']>0, 'word_captions':info['caption_count']>0,
        'word_citations':info['citation_field_count']>0, 'word_sources':info['bibliography_source_count']>0,
        'word_pages':bool(info['saved_page_count']), 'word_list_figures':info['has_list_of_figures'],
        'word_bibliography':info['has_bibliography'],
    }
    # U Wordu zobrazujeme učiteli všechny podporované kontroly. Automatická detekce
    # pouze určuje, co je ve vzoru nalezené; učitel ale musí mít možnost zvolit i
    # obsah, citace, zdroje nebo bibliografii, i když Word daný prvek uloží jiným způsobem.
    return info,[{'code':c,'label':l,'detected':bool(pred.get(c,False))} for c,l in checks]


def _ppt_shape_stats(prs):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    stats=Counter(); titles=[]; font_names=Counter(); font_sizes=Counter(); alignments=Counter()
    for slide in prs.slides:
        if slide.shapes.title and getattr(slide.shapes.title,'text','').strip(): titles.append(slide.shapes.title.text.strip())
        for sh in slide.shapes:
            st=getattr(sh,'shape_type',None)
            if st == MSO_SHAPE_TYPE.PICTURE: stats['images']+=1
            elif st == MSO_SHAPE_TYPE.TABLE: stats['tables']+=1
            elif st == MSO_SHAPE_TYPE.CHART: stats['charts']+=1
            elif st == MSO_SHAPE_TYPE.TEXT_BOX: stats['textboxes']+=1
            elif st == MSO_SHAPE_TYPE.AUTO_SHAPE: stats['autoshapes']+=1
            if hasattr(sh,'text_frame') and sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    if p.alignment is not None: alignments[str(p.alignment)]+=1
                    if p.level>0 or getattr(p,'_p',None) is not None:
                        # bullet presence is stored in pPr; text level alone is also useful.
                        xml=p._p.xml
                        if '<a:bu' in xml: stats['bullets']+=1
                    for run in p.runs:
                        if run.font.name: font_names[str(run.font.name)]+=1
                        if run.font.size: font_sizes[str(round(run.font.size.pt,2))]+=1
                        if run.font.bold: stats['bold_runs']+=1
                        if run.font.italic: stats['italic_runs']+=1
    return stats,titles,font_names,font_sizes,alignments


def analyze_powerpoint(path, original_name):
    from pptx import Presentation
    prs=Presentation(path)
    stats,titles,font_names,font_sizes,alignments=_ppt_shape_stats(prs)
    transition_slides=0; animation_count=0
    with zipfile.ZipFile(path) as z:
        for n in z.namelist():
            if re.match(r'ppt/slides/slide\d+\.xml$',n):
                root=_xml(z,n)
                if root is None: continue
                if root.find('{%s}transition' % NS_P) is not None: transition_slides += 1
                # Count user-visible animation action elements, not timing container nodes.
                tags={'anim','animClr','animEffect','animMotion','animRot','animScale','set','cmd'}
                for e in root.iter():
                    if e.tag.startswith('{%s}' % NS_P) and e.tag.split('}',1)[1] in tags:
                        animation_count += 1
    info={'extension':'.pptx','name':original_name,'type':'PowerPoint','slide_count':len(prs.slides),
          'titles':titles[:50],'title_count':len(titles),'image_count':stats['images'],'table_count':stats['tables'],
          'chart_count':stats['charts'],'shape_count':stats['autoshapes'],'textbox_count':stats['textboxes'],
          'bullet_count':stats['bullets'],'font_names':dict(font_names),'font_sizes':dict(font_sizes),
          'alignments':dict(alignments),'animation_count':animation_count,'transition_slide_count':transition_slides,
          'has_transition':transition_slides>0,'ooxml':_ppt_ooxml_profile(path)}
    checks=[('ppt_slides','Počet snímků'),('ppt_titles','Nadpisy snímků'),('ppt_images','Obrázky'),
            ('ppt_tables','Tabulky'),('ppt_charts','Grafy'),('ppt_shapes','Tvary a textová pole'),('ppt_bullets','Odrážky'),
            ('ppt_format','Formátování textu – písmo, velikost a zarovnání'),('ppt_animations','Počet animací'),
            ('ppt_transition','Použitý přechod mezi snímky'),('ppt_transition_count','Počet snímků s přechodem')]
    pred={'ppt_slides':True,'ppt_titles':info['title_count']>0,'ppt_images':info['image_count']>0,
          'ppt_tables':info['table_count']>0,'ppt_charts':info['chart_count']>0,
          'ppt_shapes':info['shape_count']+info['textbox_count']>0,'ppt_bullets':info['bullet_count']>0,
          'ppt_format':bool(info['font_names'] or info['font_sizes'] or info['alignments']),
          'ppt_animations':info['animation_count']>0,'ppt_transition':info['has_transition'],'ppt_transition_count':info['has_transition']}
    return info,[{'code':c,'label':l} for c,l in checks if pred.get(c,False)]



BLOCKED_IMPORTS={'os','sys','subprocess','socket','pathlib','shutil','requests','urllib','http','ftplib','pickle','marshal','ctypes','multiprocessing','threading','signal','resource','importlib','builtins','inspect','types','tempfile','glob'}
ALLOWED_IMPORTS={'math','statistics','decimal','fractions'}
BLOCKED_CALLS={'open','exec','eval','compile','__import__','breakpoint','globals','locals','vars','getattr','setattr','delattr','input'}
# input is allowed separately; it is removed below for safety gate.
BLOCKED_CALLS.remove('input')


def _python_safety(source):
    try:
        tree=ast.parse(source)
    except Exception as exc:
        return False, f'Syntaktická chyba: {exc}', None
    for n in ast.walk(tree):
        if isinstance(n, ast.Import):
            for a in n.names:
                if a.name.split('.')[0] in BLOCKED_IMPORTS or a.name.split('.')[0] not in ALLOWED_IMPORTS: return False, f'Nepovolený import: {a.name}', tree
        elif isinstance(n, ast.ImportFrom):
            if (n.module or '').split('.')[0] in BLOCKED_IMPORTS or (n.module or '').split('.')[0] not in ALLOWED_IMPORTS: return False, f'Nepovolený import: {n.module}', tree
        elif isinstance(n, ast.Call):
            if isinstance(n.func, ast.Name) and n.func.id in BLOCKED_CALLS: return False, f'Zakázané volání: {n.func.id}', tree
        elif isinstance(n, ast.Attribute):
            if str(n.attr).startswith('__'): return False, 'Dvojité podtržítkové atributy nejsou povolené.', tree
    return True,'',tree


def _count_inputs(tree):
    return sum(1 for n in ast.walk(tree) if isinstance(n,ast.Call) and isinstance(n.func,ast.Name) and n.func.id=='input')


def _python_cases(input_count):
    banks=[['2','3','4','5','6'],['10','4','2','8','3'],['-2','5','1','7','9'],['0','1','2','3','4']]
    if input_count<=0: return [[]]
    return [b[:input_count] for b in banks]


def _canon_number(text):
    t=str(text).replace(',', '.').strip()
    try:
        x=float(t)
        if math.isfinite(x): return format(x, '.12g')
    except Exception:
        pass
    return t


def _python_semantic(stdout):
    lines=[x.strip() for x in str(stdout or '').replace('\r','').split('\n') if x.strip()]
    if not lines: return ''
    # Prefer an explicitly labelled result line. This lets students use different
    # prompts/variable names while still comparing the actual result.
    result_line = None
    for ln in reversed(lines):
        low=ln.casefold()
        if ':' in ln and any(k in low for k in ('výsledek','vysledek','součet','soucet','celkem','result')):
            result_line=ln; break
    tail=result_line or lines[-1]
    if ':' in tail:
        value=tail.rsplit(':',1)[-1].strip()
        nums=re.findall(r'[-+]?\d+(?:[.,]\d+)?',value)
        if nums: return _canon_number(nums[-1])
        if value: return re.sub(r'\s+',' ',value).strip().casefold()
    nums=re.findall(r'[-+]?\d+(?:[.,]\d+)?',tail)
    if nums: return _canon_number(nums[-1])
    return re.sub(r'\s+',' ',tail).strip().casefold()


def _run_python(path, inputs, timeout=2.5):
    source=Path(path).read_text(encoding='utf-8',errors='replace')
    safe,msg,_=_python_safety(source)
    if not safe: return {'ok':False,'stdout':'','stderr':msg,'returncode':-1,'semantic':''}
    def limits():
        try:
            import resource
            resource.setrlimit(resource.RLIMIT_CPU,(2,2))
            resource.setrlimit(resource.RLIMIT_AS,(512*1024*1024,512*1024*1024))
            resource.setrlimit(resource.RLIMIT_FSIZE,(1024*1024,1024*1024))
            resource.setrlimit(resource.RLIMIT_NOFILE,(16,16))
        except Exception:
            pass
    try:
        cp=subprocess.run([sys.executable,'-I',str(Path(path).resolve())],input='\n'.join(inputs)+('\n' if inputs else ''),
                          text=True,capture_output=True,timeout=timeout,cwd=tempfile.mkdtemp(prefix='ucebnice_py_'),
                          env={'PYTHONIOENCODING':'utf-8','PATH':os.environ.get('PATH','')},preexec_fn=limits if os.name=='posix' else None)
        return {'ok':cp.returncode==0,'stdout':cp.stdout[-8000:],'stderr':cp.stderr[-3000:],'returncode':cp.returncode,
                'semantic':_python_semantic(cp.stdout)}
    except subprocess.TimeoutExpired:
        return {'ok':False,'stdout':'','stderr':'Program překročil časový limit.','returncode':-2,'semantic':''}
    except Exception as exc:
        return {'ok':False,'stdout':'','stderr':str(exc),'returncode':-3,'semantic':''}


def analyze_python(path, original_name):
    source=Path(path).read_text(encoding='utf-8',errors='replace')
    safe,msg,tree=_python_safety(source)
    info={'extension':'.py','name':original_name,'type':'Python','syntax_error':'','safe':safe,'safety_error':msg,'python_tests':[]}
    if tree is None:
        info['syntax_error']=msg
        return info,[{'code':'py_functionality','label':'Stejná funkčnost programu a správné výsledky'}]
    info['input_count']=_count_inputs(tree)
    info['line_count']=len(source.splitlines())
    tests=[]
    for inputs in _python_cases(info['input_count']):
        r=_run_python(path,inputs)
        if r['ok']:
            tests.append({'inputs':inputs,'expected_output':r['stdout'],'expected_semantic':r['semantic']})
    info['python_tests']=tests
    if tests:
        info['visible_example']={'inputs':tests[0]['inputs'],'output':tests[0]['expected_output'],'semantic':tests[0]['expected_semantic']}
    return info,[{'code':'py_functionality','label':'Stejná funkčnost programu a správné výsledky'}]


def analyze_file(path, original_name):
    ext=Path(original_name).suffix.lower()
    try:
        if ext in ('.xlsx','.xlsm'): return analyze_excel(path,original_name)
        if ext=='.docx': return analyze_word(path,original_name)
        if ext=='.pptx': return analyze_powerpoint(path,original_name)
        if ext=='.py': return analyze_python(path,original_name)
    except Exception as exc:
        typ={'.xlsx':'Excel','.xlsm':'Excel','.docx':'Word','.pptx':'PowerPoint','.py':'Python'}.get(ext,'Soubor')
        return {'extension':ext,'name':original_name,'type':typ,'analysis_error':str(exc)},[]
    return {'extension':ext,'name':original_name,'type':'Soubor'},[{'code':'file_type','label':'Správný typ souboru'}]


def generated_assignment(info):
    ext=info.get('extension','')
    lines=[]
    if ext in ('.xlsx','.xlsm'):
        lines.append('Vypracuj úkol v přiděleném pracovním souboru Excel. Není nutné použít stejné vzorce jako učitel; důležité jsou správné výsledky a vybrané dovednosti.')
        if info.get('sheets'): lines.append('• Požadované listy: '+', '.join(info['sheets'])+'.')
        for sp in (info.get('sheet_specs') or [])[:5]:
            if sp.get('headers'): lines.append(f'• List „{sp["name"]}“ – záhlaví: '+' | '.join(sp['headers'])+'.')
        if info.get('chart_count'): lines.append(f'• Soubor obsahuje {info["chart_count"]} graf/grafy.')
        if info.get('pivot_count'): lines.append(f'• Soubor obsahuje {info["pivot_count"]} kontingenční tabulku/tabulky.')
        lines.append('• Přesné požadavky, které se hodnotí, jsou uvedené níže.')
    elif ext=='.docx':
        lines.append('Vypracuj úkol v přiděleném pracovním souboru Word. Hodnotí se vybrané dovednosti a struktura dokumentu, ne přesná vizuální kopie učitelova vzoru.')
        if info.get('section_count',1)>1: lines.append(f'• Vzor obsahuje {info["section_count"]} oddíly.')
        if info.get('saved_page_count'): lines.append(f'• Uložený dokument má {info["saved_page_count"]} stran.')
        if info.get('image_count'): lines.append(f'• Vzor obsahuje {info["image_count"]} obrázek/obrázky.')
        lines.append('• Přesné požadavky, které se hodnotí, jsou uvedené níže.')
    elif ext=='.pptx':
        lines.append('Vypracuj úkol v přiděleném pracovním souboru PowerPoint. Hodnotí se použití funkcí, nikoli přesné rozložení nebo vzhled vzoru.')
        lines.append(f'• Vzor obsahuje {info.get("slide_count",0)} snímků.')
        if info.get('animation_count'): lines.append(f'• Vzor obsahuje {info["animation_count"]} animací.')
        if info.get('has_transition'): lines.append('• Ve vzoru je použit přechod mezi snímky.')
        lines.append('• Přesné požadavky, které se hodnotí, jsou uvedené níže.')
    elif ext=='.py':
        lines.append('Naprogramuj řešení v přiděleném souboru .py. Zdrojový kód učitele se nezobrazuje a tvoje řešení může být napsané jinak.')
        lines.append('• Rozhodující je stejná funkčnost: pro stejné vstupy musí program vracet správné výsledky.')
        ex=info.get('visible_example') or {}
        if ex:
            if ex.get('inputs'): lines.append('• Ukázkový vstup: '+', '.join(ex['inputs']))
            if ex.get('semantic')!='': lines.append('• Očekávaný výsledek: '+str(ex['semantic']))
    return '\n'.join(lines)


HINTS={
'excel_sheets':'Zkontroluj počet a názvy listů.','excel_headers':'Zkontroluj záhlaví tabulek.',
'excel_size':'Zkontroluj rozsah tabulek.','excel_filled':'Zkontroluj, zda jsou požadované části vyplněné.',
'excel_results':'Vzorec může být jiný, ale výsledná hodnota musí být správná.','excel_functions':'Zkontroluj použité funkce.',
'excel_format':'Zkontroluj písmo, velikost, zarovnání, výplň a ohraničení.','excel_number_format':'Zkontroluj formát měny, procent, data a desetinných míst.',
'excel_merged':'Zkontroluj sloučené buňky.','excel_conditional':'Zkontroluj podmíněné formátování.',
'excel_filter':'Zkontroluj filtr tabulky.','excel_table':'Zkontroluj, že data jsou vytvořená jako Excelová tabulka.',
'excel_chart':'Zkontroluj počet grafů.','excel_chart_type':'Zkontroluj typ grafu.','excel_pivot':'Zkontroluj kontingenční tabulku.',
'word_sections':'Zkontroluj počet oddílů dokumentu.','word_unlinked':'V záhlaví/zápatí nového oddílu vypni „Propojit s předchozím“.',
'word_page_numbering':'Zkontroluj nastavení číslování stránek v oddílu.','word_alignment':'Zkontroluj vodorovné i svislé zarovnání.',
'word_images':'Zkontroluj počet obrázků.','word_image_size':'Zkontroluj rozměry obrázků.','word_image_center':'Zarovnej odstavce s obrázky na střed.',
'word_title_first':'Zkontroluj nadpis/titul na začátku dokumentu.','word_toc':'Vlož automatický obsah pomocí funkce Wordu.',
'word_heading_styles':'Použij skutečné styly Nadpis 1, Nadpis 2 atd.','word_heading_numbering':'Nadpisy očísluj pomocí víceúrovňového seznamu.',
'word_captions':'Použij titulky obrázků a zarovnej je na střed.','word_citations':'Citace vlož pomocí nástroje Citace ve Wordu.',
'word_sources':'Zkontroluj, že jsou zdroje vložené přes Správce pramenů / Citace ve Wordu.',
'word_pages':'Zkontroluj požadovaný rozsah dokumentu.','word_list_figures':'Vlož automatický seznam obrázků.','word_bibliography':'Vlož seznam citací/bibliografii.',
'ppt_slides':'Zkontroluj počet snímků.','ppt_titles':'Zkontroluj nadpisy snímků.','ppt_images':'Zkontroluj obrázky.',
'ppt_tables':'Zkontroluj tabulky.','ppt_charts':'Zkontroluj grafy.','ppt_shapes':'Zkontroluj tvary a textová pole.','ppt_bullets':'Zkontroluj odrážky.',
'ppt_format':'Zkontroluj písmo, velikost a zarovnání textu.','ppt_animations':'Zkontroluj počet animací objektů.',
'ppt_transition':'Použij alespoň jeden přechod mezi snímky.','ppt_transition_count':'Zkontroluj počet snímků s přechodem.',
'py_functionality':'Spusť program s testovacími vstupy. Nevadí jiný kód, ale výsledek musí odpovídat očekávanému chování.'}

def check_hint(code): return HINTS.get(code,'Vrať se k zadání a zkontroluj tuto část práce.')


def _subset_counts(want,have):
    return all(int(have.get(k,0) or 0)>=int(v or 0) for k,v in want.items())


def _approx_image_specs(want,have,tol=0.12):
    if len(have)<len(want): return False
    unused=list(have)
    for w in want:
        found=None
        for i,h in enumerate(unused):
            ww=float(w.get('width_cm') or 0); wh=float(w.get('height_cm') or 0)
            hw=float(h.get('width_cm') or 0); hh=float(h.get('height_cm') or 0)
            if (not ww or abs(hw-ww)<=max(.15,ww*tol)) and (not wh or abs(hh-wh)<=max(.15,wh*tol)):
                found=i; break
        if found is None:return False
        unused.pop(found)
    return True


def evaluate(student_path, student_name, teacher, raw_checks):
    student,_=analyze_file(student_path,student_name)
    checks=[]
    for x in raw_checks:
        if isinstance(x,str): checks.append({'code':x,'question':'','hint':''})
        elif isinstance(x,dict) and x.get('code'): checks.append(x)
    results=[]
    for ch in checks:
        code=ch.get('code',''); ok=True; label=code
        if code=='excel_sheets': ok=(teacher.get('sheets') or [])==(student.get('sheets') or []); label='Počet a názvy listů'
        elif code=='excel_headers':
            wm={x['name']:x.get('headers',[]) for x in teacher.get('sheet_specs',[])}; hm={x['name']:x.get('headers',[]) for x in student.get('sheet_specs',[])}
            ok=all(hm.get(k)==v for k,v in wm.items()); label='Záhlaví tabulky'
        elif code=='excel_size':
            wm={x['name']:(x.get('rows',0),x.get('cols',0)) for x in teacher.get('sheet_specs',[])}; hm={x['name']:(x.get('rows',0),x.get('cols',0)) for x in student.get('sheet_specs',[])}
            ok=all(k in hm and hm[k][0]>=v[0] and hm[k][1]>=v[1] for k,v in wm.items()); label='Rozsah tabulky'
        elif code=='excel_filled': ok=int(student.get('nonempty_count',0))>=int(teacher.get('nonempty_count',0)); label='Vyplněné části tabulky'
        elif code=='excel_results':
            want={k:v for k,v in (teacher.get('formula_results') or {}).items() if v is not None}; have=student.get('formula_results') or {}
            if want: ok=all(k in have and _norm_value(have[k])==_norm_value(v) for k,v in want.items())
            else: ok=int(student.get('formula_count',0))>=int(teacher.get('formula_count',0))
            label='Správné výsledky výpočtů'
        elif code=='excel_functions': ok=set(teacher.get('formula_functions') or []).issubset(set(student.get('formula_functions') or [])); label='Požadované funkce'
        elif code=='excel_format':
            ok=_subset_counts(teacher.get('font_names') or {},student.get('font_names') or {}) and _subset_counts(teacher.get('font_sizes') or {},student.get('font_sizes') or {}) and _subset_counts(teacher.get('formatting') or {},student.get('formatting') or {}); label='Formátování buněk'
        elif code=='excel_number_format': ok=_subset_counts(teacher.get('number_formats') or {},student.get('number_formats') or {}); label='Formát čísel'
        elif code=='excel_merged': ok=(teacher.get('ooxml',{}).get('merge_refs',{})==student.get('ooxml',{}).get('merge_refs',{})); label='Sloučené buňky'
        elif code=='excel_conditional': ok=(teacher.get('ooxml',{}).get('conditional_counts',{})==student.get('ooxml',{}).get('conditional_counts',{})); label='Podmíněné formátování'
        elif code=='excel_filter': ok=(teacher.get('ooxml',{}).get('auto_filters',{})==student.get('ooxml',{}).get('auto_filters',{})); label='Filtr'
        elif code=='excel_table': ok=(teacher.get('ooxml',{}).get('table_xml_count',0)==student.get('ooxml',{}).get('table_xml_count',0) and teacher.get('ooxml',{}).get('table_part_counts',{})==student.get('ooxml',{}).get('table_part_counts',{})); label='Excelová tabulka'
        elif code=='excel_chart': ok=int(student.get('ooxml',{}).get('chart_xml_count',0))==int(teacher.get('ooxml',{}).get('chart_xml_count',0)); label='Grafy'
        elif code=='excel_chart_type': ok=_subset_counts(teacher.get('chart_types') or {},student.get('chart_types') or {}); label='Typy grafů'
        elif code=='excel_pivot': ok=int(student.get('ooxml',{}).get('pivot_xml_count',0))==int(teacher.get('ooxml',{}).get('pivot_xml_count',0)); label='Kontingenční tabulka'
        elif code=='word_sections': ok=int(student.get('section_count',0))>=int(teacher.get('section_count',0)); label='Počet oddílů'
        elif code=='word_unlinked':
            want=teacher.get('section_specs') or []; have=student.get('section_specs') or []
            ok=len(have)>=len(want) and all((i==0 or (not have[i].get('header_linked') and not have[i].get('footer_linked'))) for i in range(len(want)) if i==0 or (not want[i].get('header_linked') and not want[i].get('footer_linked'))); label='Propojení s předchozím vypnuto'
        elif code=='word_page_numbering': ok=(student.get('page_number_starts') or [])==(teacher.get('page_number_starts') or []); label='Číslování stránek'
        elif code=='word_alignment':
            ws=teacher.get('section_specs') or []; hs=student.get('section_specs') or []
            ok=_subset_counts(teacher.get('paragraph_alignments') or {},student.get('paragraph_alignments') or {}) and len(hs)>=len(ws) and all(hs[i].get('vertical_alignment')==sp.get('vertical_alignment') for i,sp in enumerate(ws)); label='Vodorovné a svislé zarovnání'
        elif code=='word_images': ok=int(student.get('image_count',0))>=int(teacher.get('image_count',0)); label='Obrázky'
        elif code=='word_image_size': ok=_approx_image_specs(teacher.get('image_specs') or [],student.get('image_specs') or []); label='Velikost obrázků'
        elif code=='word_image_center': ok=int(student.get('centered_image_count',0))>=int(teacher.get('centered_image_count',0)); label='Obrázky na střed'
        elif code=='word_title_first': ok=bool(student.get('first_paragraph_text')) and (str(student.get('first_paragraph_style','')).lower() in ('title','nadpis') or student.get('first_paragraph_text')==teacher.get('first_paragraph_text')); label='Nadpis na první stránce'
        elif code=='word_toc': ok=bool(student.get('has_toc')); label='Automatický obsah'
        elif code=='word_heading_styles': ok=(teacher.get('heading_styles') or {})==(student.get('heading_styles') or {}); label='Styly a počet nadpisů'
        elif code=='word_heading_numbering': ok=int(student.get('numbered_heading_count',0))==int(teacher.get('numbered_heading_count',0)) and int(teacher.get('numbered_heading_count',0))>0; label='Číslování nadpisů'
        elif code=='word_captions': ok=int(student.get('caption_count',0))>=int(teacher.get('caption_count',0)) and int(student.get('centered_caption_count',0))>=int(teacher.get('centered_caption_count',0)); label='Titulky obrázků'
        elif code=='word_citations':
            want=max(1,int(teacher.get('citation_field_count',0) or 0)); ok=int(student.get('citation_field_count',0) or 0)>=want; label='Vložené citace'
        elif code=='word_sources':
            want=max(1,int(teacher.get('bibliography_source_count',0) or 0)); ok=int(student.get('bibliography_source_count',0) or 0)>=want; label='Vložené zdroje citací'
        elif code=='word_pages': ok=bool(student.get('saved_page_count')) and int(student.get('saved_page_count'))>=int(teacher.get('saved_page_count') or 0); label='Počet stran'
        elif code=='word_list_figures': ok=bool(student.get('has_list_of_figures')); label='Seznam obrázků'
        elif code=='word_bibliography': ok=bool(student.get('has_bibliography')); label='Bibliografie'
        elif code=='ppt_slides': ok=int(student.get('slide_count',0))==int(teacher.get('slide_count',0)); label='Počet snímků'
        elif code=='ppt_titles': ok=(student.get('titles') or [])==(teacher.get('titles') or []); label='Nadpisy snímků'
        elif code=='ppt_images': ok=int(student.get('image_count',0))==int(teacher.get('image_count',0)); label='Obrázky'
        elif code=='ppt_tables': ok=int(student.get('table_count',0))==int(teacher.get('table_count',0)); label='Tabulky'
        elif code=='ppt_charts': ok=int(student.get('chart_count',0))==int(teacher.get('chart_count',0)); label='Grafy'
        elif code=='ppt_shapes': ok=[x.get('counts',{}) for x in student.get('ooxml',{}).get('slides',[])]==[x.get('counts',{}) for x in teacher.get('ooxml',{}).get('slides',[])]; label='Tvary a textová pole'
        elif code=='ppt_bullets': ok=int(student.get('bullet_count',0))==int(teacher.get('bullet_count',0)); label='Odrážky'
        elif code=='ppt_format': ok=_subset_counts(teacher.get('font_names') or {},student.get('font_names') or {}) and _subset_counts(teacher.get('font_sizes') or {},student.get('font_sizes') or {}); label='Formátování textu'
        elif code=='ppt_animations': ok=[x.get('animations',0) for x in student.get('ooxml',{}).get('slides',[])]==[x.get('animations',0) for x in teacher.get('ooxml',{}).get('slides',[])]; label='Počet animací'
        elif code=='ppt_transition': ok=[x.get('transition',False) for x in student.get('ooxml',{}).get('slides',[])]==[x.get('transition',False) for x in teacher.get('ooxml',{}).get('slides',[])]; label='Přechod mezi snímky'
        elif code=='ppt_transition_count': ok=int(student.get('transition_slide_count',0))==int(teacher.get('transition_slide_count',0)); label='Počet snímků s přechodem'
        elif code=='py_functionality':
            ok=True
            for tc in teacher.get('python_tests') or []:
                r=_run_python(student_path,tc.get('inputs') or [])
                if not r.get('ok') or r.get('semantic')!=tc.get('expected_semantic'):
                    ok=False; break
            if not (teacher.get('python_tests') or []): ok=False
            label='Stejná funkčnost programu'
        results.append({'code':code,'label':label,'question':ch.get('question',''),'ok':bool(ok),'hint':ch.get('hint') or check_hint(code)})
    return results


def preview(path, original_name, teacher=False):
    ext=Path(original_name).suffix.lower()
    if ext=='.py':
        info,_=analyze_python(path,original_name)
        ex=info.get('visible_example') or {}
        if teacher:
            text='Zdrojový kód učitele je skrytý.'
            if ex:
                text += '\n\nUkázkový test:\n'
                if ex.get('inputs'): text += 'Vstup: '+', '.join(ex['inputs'])+'\n'
                text += 'Výsledek: '+str(ex.get('semantic',''))
            return {'kind':'text','text':text}
        return {'kind':'code','text':Path(path).read_text(encoding='utf-8',errors='replace')[:18000]}
    if ext in ('.xlsx','.xlsm'):
        try:
            import openpyxl
            wb=openpyxl.load_workbook(path,data_only=False); ws=next(w for w in wb.worksheets if w.title!='__UCEBNICE_ID__')
            rows=[[c.value for c in row] for row in ws.iter_rows(min_row=1,max_row=min(ws.max_row,25),max_col=min(ws.max_column,12))]
            return {'kind':'table','title':ws.title,'rows':rows}
        except Exception as exc:return {'kind':'text','text':f'Náhled se nepodařilo vytvořit: {exc}'}
    if ext=='.docx':
        try:
            from docx import Document
            return {'kind':'text','text':'\n'.join(p.text for p in Document(path).paragraphs if p.text.strip())[:18000]}
        except Exception as exc:return {'kind':'text','text':f'Náhled se nepodařilo vytvořit: {exc}'}
    if ext=='.pptx':
        try:
            from pptx import Presentation
            prs=Presentation(path); lines=[]
            for i,slide in enumerate(prs.slides[:25],1):
                texts=[sh.text.strip() for sh in slide.shapes if hasattr(sh,'text') and sh.text.strip()]
                lines.append(f'Snímek {i}: '+' | '.join(texts))
            return {'kind':'text','text':'\n'.join(lines)}
        except Exception as exc:return {'kind':'text','text':f'Náhled se nepodařilo vytvořit: {exc}'}
    return {'kind':'text','text':'Soubor byl nahrán.'}
